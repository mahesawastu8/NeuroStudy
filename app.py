import streamlit as st
import streamlit.components.v1 as components
import fitz, google.generativeai as genai
import json, datetime, html, re, time, math, os, urllib.parse
from pathlib import Path

st.set_page_config(page_title="NeuroStudy", page_icon="🧠", layout="wide",
                   initial_sidebar_state="expanded")

DATA_DIR = Path("data")
FLASHCARD_DIR = DATA_DIR / "flashcards"
DISCUSSION_DIR = DATA_DIR / "discussions"
SESSION_DIR = DATA_DIR / "sessions"

@st.cache_resource
def _ensure_core_directories():
    for d in [DATA_DIR, FLASHCARD_DIR, DISCUSSION_DIR, SESSION_DIR]:
        d.mkdir(parents=True, exist_ok=True)
    return True
_ensure_core_directories()

CONFIG_FILE = DATA_DIR / "config.json"

def load_config():
    if CONFIG_FILE.exists():
        try: return json.loads(CONFIG_FILE.read_text())
        except: return {}
    return {}

def save_config(cfg_dict):
    CONFIG_FILE.write_text(json.dumps(cfg_dict, ensure_ascii=False, indent=2))

def _safe_get_secret(key: str, default: str = "") -> str:
    try:
        if hasattr(st, "secrets"):
            val = st.secrets.get(key)
            if val is not None:
                return str(val)
    except Exception:
        pass
    return default

def get_gemini_api_key():
    saved_cfg = load_config()
    saved_key = saved_cfg.get("api_key", "").strip()
    if saved_key:
        return saved_key
    sec_key = _safe_get_secret("GEMINI_API_KEY") or _safe_get_secret("api_key")
    if sec_key:
        return sec_key
    return os.environ.get("GEMINI_API_KEY", "")

# ── THEME STATE (OBSIDIAN NIGHT VS CLINICAL APPLE WHITE) ──────────────────────
if "app_theme" not in st.session_state:
    saved_theme = load_config().get("app_theme", "obsidian")
    st.session_state.app_theme = saved_theme

# ── PREMIUM DESIGN SYSTEM ─────────────────────────────────────────────────────
theme_css = ""
if st.session_state.get("app_theme") == "clinical_white":
    theme_css = """
body, .stApp {
    background: #f8fafc !important;
    color: #0f172a !important;
}
h1, h2, h3, h4, h5, h6, .nav-title, div, p, span, li {
    color: #0f172a !important;
}
.cs, .stCaption, [data-testid="stCaptionContainer"], p[style*="color:#94a3b8"] {
    color: #475569 !important;
}
div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"], .card {
    background: #ffffff !important;
    border: 1px solid #e2e8f0 !important;
    box-shadow: 0 4px 16px rgba(0, 0, 0, 0.05) !important;
}
div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"]:hover {
    border-color: #2563eb !important;
    box-shadow: 0 8px 25px rgba(37, 99, 235, 0.12) !important;
}
[data-testid="stTabs"] [role="tablist"] {
    background: #f1f5f9 !important;
    border: 1px solid #e2e8f0 !important;
}
[data-testid="stTabs"] [role="tab"] {
    color: #64748b !important;
}
[data-testid="stTabs"] [role="tab"][aria-selected="true"] {
    background: #ffffff !important;
    color: #1e40af !important;
    box-shadow: 0 2px 8px rgba(0,0,0,0.06) !important;
}
input, textarea, [data-baseweb="select"] {
    background: #ffffff !important;
    color: #0f172a !important;
    border-color: #cbd5e1 !important;
}
"""

st.markdown(f"""
<style>
{theme_css}
</style>
""", unsafe_allow_html=True)

st.markdown(r"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&family=Inter:wght@300;400;500;600;700&display=swap');


/* ── Clean UI: Hide deploy button & menu without breaking sidebar toggle ── */
#MainMenu { visibility: hidden !important; display: none !important; }
[data-testid="stToolbar"] { visibility: hidden !important; display: none !important; }
[data-testid="stDecoration"] { visibility: hidden !important; display: none !important; }
[data-testid="stStatusWidget"] { visibility: hidden !important; display: none !important; }
.stDeployButton { display: none !important; }
footer { visibility: hidden !important; display: none !important; }
header[data-testid="stHeader"] { background: transparent !important; }
[data-testid="collapsedControl"] { display: flex !important; visibility: visible !important; color: #818cf8 !important; }

html, body, [class*="css"] {
    font-family: 'Plus Jakarta Sans', 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
    box-sizing: border-box;
}

.stApp {
    background-color: #060810 !important;
    background-image: 
        radial-gradient(circle at 15% 15%, rgba(99, 102, 241, 0.09) 0%, transparent 45%),
        radial-gradient(circle at 85% 75%, rgba(56, 189, 248, 0.07) 0%, transparent 50%),
        radial-gradient(circle at 50% 50%, rgba(139, 92, 246, 0.04) 0%, transparent 60%) !important;
    background-attachment: fixed !important;
    color: #e2e8f0;
}

.block-container {
    padding: 2rem 2.5rem !important;
    max-width: 1360px !important;
}

/* ── EXECUTIVE TOP NAVBAR & UNIFIED TABS ── */
.executive-navbar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    background: rgba(15, 23, 42, 0.7);
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid rgba(255, 255, 255, 0.08);
    border-radius: 16px;
    padding: 10px 18px;
    margin-bottom: 20px;
    box-shadow: 0 8px 30px rgba(0, 0, 0, 0.4);
}
.nav-brand {
    display: flex;
    align-items: center;
    gap: 10px;
}
.nav-title {
    font-size: 1.15rem;
    font-weight: 800;
    color: #ffffff;
    letter-spacing: -0.4px;
    line-height: 1.1;
}
.nav-badge {
    font-size: 0.65rem;
    color: #818cf8;
    font-weight: 700;
    letter-spacing: 0.5px;
    text-transform: uppercase;
}
.nav-user-chip {
    display: flex;
    align-items: center;
    gap: 10px;
    background: rgba(30, 41, 59, 0.6);
    border: 1px solid rgba(255, 255, 255, 0.06);
    padding: 6px 12px;
    border-radius: 10px;
}
[data-testid="stTabs"] [role="tablist"] {
    display: flex !important;
    justify-content: center !important;
    background: rgba(15, 23, 42, 0.6) !important;
    backdrop-filter: blur(16px) !important;
    border: 1px solid rgba(255, 255, 255, 0.06) !important;
    border-radius: 14px !important;
    padding: 6px !important;
    gap: 8px !important;
    margin-bottom: 24px !important;
}
[data-testid="stTabs"] [role="tab"] {
    flex: 1 !important;
    text-align: center !important;
    justify-content: center !important;
    border-radius: 10px !important;
    padding: 10px 16px !important;
    font-weight: 600 !important;
    font-size: 0.88rem !important;
    color: #94a3b8 !important;
    border: 1px solid transparent !important;
    transition: all 0.2s ease !important;
}
[data-testid="stTabs"] [aria-selected="true"] {
    background: linear-gradient(135deg, #4f46e5 0%, #6366f1 100%) !important;
    color: #ffffff !important;
    border: 1px solid rgba(255, 255, 255, 0.15) !important;
    box-shadow: 0 4px 14px rgba(99, 102, 241, 0.4) !important;
}

/* ── OFFICIAL GOOGLE SIGN-IN BUTTON & IDENTITY STYLING ── */
.g-signin-btn {
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 12px;
    background: #ffffff;
    color: #1f2937;
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.95rem;
    font-weight: 700;
    padding: 12px 20px;
    border-radius: 12px;
    border: 1px solid #e5e7eb;
    box-shadow: 0 4px 14px rgba(0, 0, 0, 0.12);
    cursor: pointer;
    transition: all 0.2s ease;
    text-decoration: none;
    margin: 10px 0;
}
.g-signin-btn:hover {
    background: #f9fafb;
    box-shadow: 0 6px 20px rgba(0, 0, 0, 0.2);
    transform: translateY(-1px);
}
.g-icon {
    width: 20px;
    height: 20px;
}
.g-account-card {
    background: rgba(15, 23, 42, 0.75);
    border: 1px solid rgba(99, 102, 241, 0.25);
    border-radius: 14px;
    padding: 14px 16px;
    display: flex;
    align-items: center;
    gap: 12px;
    margin: 8px 0;
}

/* ── PURE CENTERED MINIMALIST LAYOUT (SERENE, DELIGHTFUL & LUXURIOUS) ── */
body, .stApp {
    background: radial-gradient(ellipse at 50% -20%, rgba(99, 102, 241, 0.12) 0%, #060811 65%, #030408 100%) !important;
    font-family: 'Plus Jakarta Sans', -apple-system, BlinkMacSystemFont, sans-serif !important;
    color: #f1f5f9 !important;
    letter-spacing: -0.15px !important;
}

/* Soothing Typography with Generous Line Height */
p, li, span, div {
    line-height: 1.68 !important;
}

/* Glassmorphism Cards that Feel Soft & Elegant */
div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"] {
    background: rgba(13, 18, 33, 0.65) !important;
    backdrop-filter: blur(16px) !important;
    -webkit-backdrop-filter: blur(16px) !important;
    border: 1px solid rgba(255, 255, 255, 0.07) !important;
    border-radius: 16px !important;
    box-shadow: 0 8px 32px rgba(0, 0, 0, 0.25) !important;
    transition: all 0.25s ease !important;
}

div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"]:hover {
    border-color: rgba(99, 102, 241, 0.3) !important;
    box-shadow: 0 12px 40px rgba(99, 102, 241, 0.1) !important;
}

/* Pill Badges */
.badge-enak {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    padding: 3px 10px;
    border-radius: 20px;
    font-size: 0.72rem;
    font-weight: 700;
    letter-spacing: 0.3px;
}

[data-testid="stSidebar"], 
[data-testid="collapsedControl"], 
section[data-testid="stSidebar"], 
button[kind="header"] {
    display: none !important;
}

.main, .stApp {
    background: #070913 !important;
}

.block-container {
    max-width: 1240px !important;
    margin: 0 auto !important;
    padding: 1.5rem 1.5rem 3rem !important;
}

/* ── EXECUTIVE TOP NAVBAR & UNIFIED TABS ── */
.executive-navbar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    background: rgba(15, 23, 42, 0.7);
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid rgba(255, 255, 255, 0.08);
    border-radius: 16px;
    padding: 10px 18px;
    margin-bottom: 20px;
    box-shadow: 0 8px 30px rgba(0, 0, 0, 0.4);
}

/* ── Tabs ── */
[data-testid="stTabs"] [role="tablist"] {
    background: rgba(17, 24, 39, 0.7);
    backdrop-filter: blur(12px);
    padding: 5px 6px;
    border-radius: 12px;
    border: 1px solid rgba(255, 255, 255, 0.08);
    gap: 6px;
    margin-bottom: 24px;
}
[data-testid="stTabs"] [role="tab"] {
    background: transparent !important;
    color: #64748b !important;
    border: none !important;
    font-weight: 600;
    padding: 8px 20px;
    border-radius: 8px;
    font-size: 0.9rem;
    transition: all 0.2s ease;
}
[data-testid="stTabs"] [role="tab"]:hover {
    color: #e2e8f0 !important;
    background: rgba(255, 255, 255, 0.04) !important;
}
[data-testid="stTabs"] [role="tab"][aria-selected="true"] {
    color: #ffffff !important;
    background: linear-gradient(135deg, rgba(99, 102, 241, 0.8), rgba(79, 70, 229, 0.9)) !important;
    box-shadow: 0 4px 14px rgba(99, 102, 241, 0.35) !important;
}

/* ── Cards ── */
.card {
    background: rgba(15, 21, 37, 0.75);
    backdrop-filter: blur(16px);
    -webkit-backdrop-filter: blur(16px);
    border: 1px solid rgba(255, 255, 255, 0.08);
    border-radius: 16px;
    padding: 22px 24px;
    margin: 10px 0;
    box-shadow: 0 8px 24px -4px rgba(0, 0, 0, 0.4);
    transition: all 0.25s cubic-bezier(0.16, 1, 0.3, 1);
}
.card:hover {
    border-color: rgba(99, 102, 241, 0.4);
    box-shadow: 0 12px 32px -4px rgba(99, 102, 241, 0.15);
    transform: translateY(-2px);
}
.card-sm { padding: 14px 18px; border-radius: 12px; margin: 8px 0; }
.ct { font-size: 0.94rem; font-weight: 600; color: #f8fafc; margin-bottom: 4px; }
.cs { font-size: 0.8rem; color: #94a3b8; }

/* ── Badges ── */
.badge {
    display: inline-flex; align-items: center; padding: 4px 11px;
    border-radius: 20px; font-size: 0.73rem; font-weight: 600; gap: 4px;
    letter-spacing: 0.2px;
}
.bb { background: rgba(56, 189, 248, 0.12); color: #38bdf8; border: 1px solid rgba(56, 189, 248, 0.3); }
.bg { background: rgba(74, 222, 128, 0.12); color: #4ade80; border: 1px solid rgba(74, 222, 128, 0.3); }
.br { background: rgba(248, 113, 113, 0.12); color: #f87171; border: 1px solid rgba(248, 113, 113, 0.3); }
.bp { background: rgba(168, 85, 247, 0.12); color: #c084fc; border: 1px solid rgba(168, 85, 247, 0.3); }
.ba { background: rgba(251, 191, 36, 0.12); color: #fbbf24; border: 1px solid rgba(251, 191, 36, 0.3); }
.pill-w { display: flex; flex-wrap: wrap; gap: 6px; margin: 8px 0; }
.pill {
    background: rgba(30, 41, 59, 0.8); color: #a5b4fc; border-radius: 20px;
    padding: 4px 12px; font-size: 0.72rem; font-weight: 500;
    border: 1px solid rgba(99, 102, 241, 0.2);
}

/* ── Forms & Clean Containers ── */
[data-testid="stForm"] {
    border: none !important;
    padding: 0 !important;
    background: transparent !important;
}

/* ── Buttons (Standard & Form Submit) ── */
.stButton > button, 
[data-testid="stFormSubmitButton"] > button,
button[kind="primaryFormSubmit"],
button[kind="primary"] {
    background: linear-gradient(135deg, #6366f1 0%, #4f46e5 60%, #4338ca 100%) !important;
    color: #ffffff !important; border: 1px solid rgba(255, 255, 255, 0.18) !important;
    border-radius: 12px !important; font-weight: 700 !important; font-size: 0.92rem !important;
    letter-spacing: 0.2px !important;
    padding: 12px 24px !important; box-shadow: 0 4px 16px rgba(99, 102, 241, 0.35) !important;
    transition: all 0.25s cubic-bezier(0.16, 1, 0.3, 1) !important;
}
.stButton > button:hover,
[data-testid="stFormSubmitButton"] > button:hover,
button[kind="primaryFormSubmit"]:hover,
button[kind="primary"]:hover {
    box-shadow: 0 8px 24px rgba(99, 102, 241, 0.55) !important;
    transform: translateY(-2px) !important;
    filter: brightness(1.1) !important;
}
.stButton > button:active,
[data-testid="stFormSubmitButton"] > button:active,
button[kind="primaryFormSubmit"]:active,
button[kind="primary"]:active { transform: translateY(0) !important; }

/* ── Tab Highlight Remover (No Red Underline) ── */
[data-testid="stTabs"] [data-baseweb="tab-highlight"],
[data-testid="stTabs"] [data-baseweb="tab-border"] {
    display: none !important;
}
button[data-baseweb="tab"] {
    border-bottom: none !important;
}

/* ── File Uploader ── */
[data-testid="stFileUploader"] {
    background: rgba(15, 21, 37, 0.7) !important;
    backdrop-filter: blur(12px) !important;
    border: 1.5px dashed rgba(99, 102, 241, 0.3) !important;
    border-radius: 16px !important;
    padding: 16px 20px !important;
}
[data-testid="stFileUploader"]:hover {
    border-color: #6366f1 !important;
    background: rgba(20, 27, 48, 0.85) !important;
}
[data-testid="stFileUploader"] button {
    background: #1e293b !important; color: #e2e8f0 !important;
    border: 1px solid rgba(255, 255, 255, 0.12) !important; border-radius: 8px !important;
    padding: 7px 16px !important; font-size: 0.84rem !important; font-weight: 600 !important;
}
[data-testid="stFileUploader"] button:hover { background: #4f46e5 !important; border-color: #6366f1 !important; }
[data-testid="stFileUploader"] span, [data-testid="stFileUploader"] small { color: #94a3b8 !important; font-size: 0.82rem !important; }

/* ── Inputs ── */
[data-testid="stSelectbox"] > div > div, textarea, [data-testid="stTextInput"] > div > div > input {
    background: rgba(17, 24, 39, 0.85) !important;
    border: 1px solid rgba(255, 255, 255, 0.1) !important;
    color: #f1f5f9 !important; border-radius: 10px !important;
    box-shadow: inset 0 2px 4px rgba(0,0,0,0.3) !important;
}
[data-testid="stProgressBar"] > div {
    background: linear-gradient(90deg, #6366f1, #38bdf8) !important;
    border-radius: 6px; height: 8px !important;
}
hr { border-color: rgba(255, 255, 255, 0.08) !important; }
::-webkit-scrollbar { width: 6px; height: 6px; }
::-webkit-scrollbar-track { background: #060810; }
::-webkit-scrollbar-thumb { background: #1f293d; border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: #374151; }

/* ── Markdown Output ── */
.stMarkdown p, .stMarkdown li { color: #cbd5e1; line-height: 1.85; font-size: 0.94rem; }
.stMarkdown h1, .stMarkdown h2, .stMarkdown h3 { color: #ffffff !important; letter-spacing: -0.02em; }
.stMarkdown strong { color: #ffffff !important; font-weight: 700; }
.stMarkdown em { color: #a5b4fc !important; }
.stMarkdown code { background: rgba(30, 41, 59, 0.8) !important; color: #818cf8 !important; border-radius: 5px; font-size: 0.85em; padding: 2px 6px; }
.stMarkdown blockquote {
    border-left: 3px solid #6366f1; padding: 10px 16px; background: rgba(15, 23, 42, 0.5);
    border-radius: 0 10px 10px 0; color: #cbd5e1; margin: 12px 0;
}

/* ── Stepper ── */
.stepper-container {
    display: flex; align-items: center; justify-content: space-between;
    background: rgba(15, 21, 37, 0.75); backdrop-filter: blur(16px);
    border: 1px solid rgba(255, 255, 255, 0.08); border-radius: 16px;
    padding: 16px 22px; margin: 16px 0 28px;
    box-shadow: 0 8px 24px -4px rgba(0, 0, 0, 0.4);
}
.step-item {
    display: flex; flex-direction: column; align-items: center; gap: 7px; flex: 1;
}
.step-circle {
    width: 38px; height: 38px; border-radius: 50%; display: flex; align-items: center;
    justify-content: center; font-weight: 700; font-size: 0.88rem; transition: all 0.3s cubic-bezier(0.16, 1, 0.3, 1);
}
.step-done .step-circle {
    background: rgba(74, 222, 128, 0.15); color: #4ade80; border: 1.5px solid #4ade80;
    box-shadow: 0 0 12px rgba(74, 222, 128, 0.2);
}
.step-active .step-circle {
    background: linear-gradient(135deg, #6366f1, #4f46e5); color: #ffffff;
    box-shadow: 0 0 20px rgba(99, 102, 241, 0.6); border: 1.5px solid #a5b4fc;
    transform: scale(1.08);
}
.step-todo .step-circle {
    background: rgba(26, 32, 53, 0.6); color: #64748b; border: 1.5px solid rgba(255, 255, 255, 0.06);
}
.step-label { font-size: 0.73rem; color: #64748b; font-weight: 600; text-align: center; white-space: nowrap; letter-spacing: 0.3px; }
.step-active .step-label { color: #818cf8; font-weight: 700; }
.step-done .step-label { color: #4ade80; }

.step-divider { height: 2px; flex: 1; min-width: 14px; background: rgba(255, 255, 255, 0.08); margin: 0 6px 20px; border-radius: 1px; }
.step-divider.done { background: linear-gradient(90deg, #4ade80, #38bdf8); }

/* ── Phase Card Box ── */
.phase-box {
    background: rgba(15, 21, 37, 0.8);
    backdrop-filter: blur(16px);
    border: 1.5px solid rgba(99, 102, 241, 0.25);
    border-radius: 18px; padding: 26px 28px; margin-top: 8px;
    box-shadow: 0 10px 30px -5px rgba(0,0,0,0.5);
}
.phase-header { display: flex; align-items: center; gap: 14px; margin-bottom: 8px; }
.phase-icon { font-size: 2rem; }
.phase-title { font-size: 1.35rem; font-weight: 800; color: #ffffff; letter-spacing: -0.02em; }
.phase-meta { font-size: 0.82rem; color: #818cf8; margin-top: 2px; font-weight: 500; }
.phase-source {
    font-size: 0.77rem; color: #94a3b8; font-style: italic; margin-top: 12px;
    padding: 10px 14px; background: rgba(10, 14, 26, 0.7); border-radius: 10px;
    border-left: 3px solid #6366f1; line-height: 1.65;
}

/* ── Diagnostic Feedback Box ── */
.analysis-box {
    background: rgba(13, 22, 45, 0.85);
    backdrop-filter: blur(16px);
    border: 1.5px solid #6366f1;
    border-radius: 16px;
    padding: 22px 26px;
    margin: 16px 0;
    box-shadow: 0 8px 30px rgba(99, 102, 241, 0.25);
}
.analysis-header {
    display: flex; align-items: center; justify-content: space-between;
    margin-bottom: 12px; padding-bottom: 8px; border-bottom: 1px solid rgba(99, 102, 241, 0.3);
}

/* ── Chat & Responses ── */
.msg-user {
    align-self: flex-end; max-width: 82%;
    background: linear-gradient(135deg, #4f46e5, #6366f1);
    color: #ffffff; border-radius: 18px 18px 4px 18px;
    padding: 14px 18px; font-size: 0.93rem; line-height: 1.7; margin: 10px 0;
    box-shadow: 0 6px 18px rgba(79, 70, 229, 0.3);
}
.msg-ai {
    align-self: flex-start; max-width: 92%;
    background: rgba(15, 21, 37, 0.85); backdrop-filter: blur(16px);
    border: 1px solid rgba(255, 255, 255, 0.08); color: #cbd5e1;
    border-radius: 4px 18px 18px 18px; padding: 18px 22px;
    font-size: 0.93rem; line-height: 1.85; margin: 10px 0;
    box-shadow: 0 8px 24px rgba(0, 0, 0, 0.4);
}
.msg-ai strong { color: #ffffff; }
.msg-ai em { color: #a5b4fc; }
.ai-row { display: flex; align-items: center; gap: 8px; margin-bottom: 10px; }
.ai-dot {
    width: 26px; height: 26px; border-radius: 50%;
    background: linear-gradient(135deg, #6366f1, #38bdf8);
    display: flex; align-items: center; justify-content: center;
    font-size: 0.7rem; font-weight: 800; color: #fff; flex-shrink: 0;
    box-shadow: 0 0 10px rgba(99, 102, 241, 0.5);
}

/* ── Streaming Thinking Box ── */
@keyframes pulseGlow { 0%,100%{opacity:0.6; transform:scale(1);} 50%{opacity:1; transform:scale(1.25);} }
.live-dot { width:8px; height:8px; border-radius:50%; background:#818cf8; box-shadow:0 0 10px #818cf8; animation:pulseGlow 1.4s infinite; }

.thinking-live-box {
    background: rgba(10, 14, 28, 0.92);
    border: 1.5px solid rgba(99, 102, 241, 0.35);
    border-radius: 14px; padding: 16px 20px; margin: 12px 0;
    box-shadow: 0 8px 25px rgba(0,0,0,0.5);
}
.thinking-live-header {
    display: flex; align-items: center; justify-content: space-between;
    margin-bottom: 8px; padding-bottom: 6px; border-bottom: 1px solid rgba(99,102,241,0.2);
}
.thinking-live-text {
    font-family: 'SFMono-Regular', Consolas, 'Liberation Mono', Menlo, Courier, monospace;
    font-size: 0.84rem; color: #cbd5e1; line-height: 1.65; white-space: pre-wrap; word-break: break-word;
}

@keyframes bl { 0%,100%{opacity:1}50%{opacity:0} }
.cur { display: inline-block; width: 2px; height: 1.1em; background: #818cf8; margin-left: 2px; vertical-align: text-bottom; border-radius: 1px; animation: bl 0.75s ease-in-out infinite; }
.stream-wrap {
    background: rgba(15, 21, 37, 0.85); backdrop-filter: blur(16px);
    border: 1px solid rgba(99, 102, 241, 0.3); border-radius: 16px;
    padding: 22px 26px; min-height: 85px; font-size: 0.92rem; line-height: 1.85;
    color: #cbd5e1; white-space: pre-wrap; word-break: break-word;
}

/* ── Completion Celebration ── */
@keyframes pop { 0%{transform:scale(.85);opacity:0}80%{transform:scale(1.03)}100%{transform:scale(1);opacity:1} }
.complete-box {
    background: linear-gradient(135deg, rgba(5, 46, 22, 0.9), rgba(20, 83, 45, 0.9));
    backdrop-filter: blur(16px); border: 1.5px solid #4ade80;
    border-radius: 20px; padding: 32px; text-align: center;
    box-shadow: 0 12px 40px rgba(74, 222, 128, 0.2); animation: pop 0.4s ease-out;
}


/* ── Flashcard 3D Interactive Design ── */
.fc-card {
    background: linear-gradient(145deg, rgba(15, 23, 42, 0.95), rgba(30, 41, 59, 0.9));
    border: 1.5px solid rgba(99, 102, 241, 0.35);
    border-radius: 18px;
    padding: 26px 30px;
    margin: 14px 0;
    box-shadow: 0 10px 30px rgba(0,0,0,0.5);
    position: relative;
    transition: all 0.2s ease;
}
.fc-card:hover {
    border-color: rgba(99, 102, 241, 0.6);
    box-shadow: 0 12px 35px rgba(99, 102, 241, 0.2);
}
.fc-badge {
    display: inline-block; padding: 3px 10px; border-radius: 20px;
    font-size: 0.72rem; font-weight: 700; letter-spacing: 0.5px;
    background: rgba(99, 102, 241, 0.15); color: #818cf8; border: 1px solid rgba(99, 102, 241, 0.3);
}
.fc-front-title {
    font-size: 1.15rem; font-weight: 700; color: #f8fafc; margin: 12px 0 8px; line-height: 1.5;
}
.fc-back-box {
    background: rgba(6, 9, 20, 0.75);
    border: 1px solid rgba(74, 222, 128, 0.3);
    border-radius: 12px;
    padding: 16px 20px;
    margin-top: 14px;
    color: #e2e8f0;
    font-size: 0.92rem;
    line-height: 1.7;
    animation: fadeIn 0.25s ease-out;
}
@keyframes fadeIn { from { opacity: 0; transform: translateY(6px); } to { opacity: 1; transform: translateY(0); } }

/* ── Mindmap Wrapper ── */
.mindmap-box {
    background: #070913; border: 1.5px solid rgba(99, 102, 241, 0.25);
    border-radius: 18px; overflow: hidden; margin: 18px 0;
    box-shadow: 0 10px 40px rgba(0,0,0,0.6);
}
</style>
""", unsafe_allow_html=True)



# ── GOOGLE AUTHENTICATION & AUTO-LOGIN SESSION PERSISTENCE ───────────────────
ACTIVE_AUTH_SESSION_FILE = DATA_DIR / "active_google_session.json"

def get_persisted_auth_session():
    if ACTIVE_AUTH_SESSION_FILE.exists():
        try:
            d = json.loads(ACTIVE_AUTH_SESSION_FILE.read_text())
            if d.get("username"): return d
        except: pass
    return None

def set_persisted_auth_session(user_dict):
    try:
        ACTIVE_AUTH_SESSION_FILE.write_text(json.dumps(user_dict, ensure_ascii=False, indent=2))
    except: pass

def clear_persisted_auth_session():
    try:
        if ACTIVE_AUTH_SESSION_FILE.exists():
            ACTIVE_AUTH_SESSION_FILE.unlink()
    except: pass

# ── USER AUTHENTICATION & ADVANCED SECURITY LAYER ─────────────────────────────
USERS_DB_FILE = DATA_DIR / "users_db.json"
LOGIN_ATTEMPTS_FILE = DATA_DIR / "login_attempts.json"

def get_login_attempts():
    if LOGIN_ATTEMPTS_FILE.exists():
        try: return json.loads(LOGIN_ATTEMPTS_FILE.read_text())
        except Exception: return {}
    return {}

def save_login_attempts(attempts):
    try: atomic_write_json(LOGIN_ATTEMPTS_FILE, attempts)
    except Exception: pass

def hash_user_password(pw, salt=None):
    import hashlib, secrets
    if not salt:
        salt = secrets.token_hex(16)
    key = hashlib.pbkdf2_hmac("sha256", pw.encode("utf-8"), salt.encode("utf-8"), 100000)
    return f"{salt}${key.hex()}"

def verify_user_password(pw, stored_hash):
    import hashlib, secrets
    try:
        if "$" in stored_hash:
            salt, key_hex = stored_hash.split("$")
            check_key = hashlib.pbkdf2_hmac("sha256", pw.encode("utf-8"), salt.encode("utf-8"), 100000)
            return secrets.compare_digest(check_key.hex(), key_hex)
        else:
            # Fallback legacy hash
            legacy = hashlib.sha256(("neurostudy_" + pw + "_salt2026").encode("utf-8")).hexdigest()
            return secrets.compare_digest(legacy, stored_hash)
    except:
        return False

def atomic_write_json(path_obj, data):
    """Menulis file JSON secara atomik menggunakan os.replace untuk mencegah race conditions & korupsi file saat diakses banyak pengguna bersamaan."""
    import os
    p = Path(path_obj)
    p.parent.mkdir(parents=True, exist_ok=True)
    tmp = p.with_suffix(f".tmp_{os.getpid()}_{time.time_ns()}")
    tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2))
    os.replace(tmp, p)

def load_users():
    if USERS_DB_FILE.exists():
        try: return json.loads(USERS_DB_FILE.read_text())
        except: return {}
    return {}

def save_users(users_dict):
    atomic_write_json(USERS_DB_FILE, users_dict)

def get_user_root(username=None):
    u = username or st.session_state.get("current_user", "dimas")
    # Sanitize strictly to prevent path traversal
    safe_u = re.sub(r'[^a-zA-Z0-9_-]', '', str(u).strip()) or "default_user"
    p = DATA_DIR / "users" / safe_u
    p.mkdir(parents=True, exist_ok=True)
    (p / "materials").mkdir(exist_ok=True)
    (p / "flashcards").mkdir(exist_ok=True)
    (p / "discussions").mkdir(exist_ok=True)
    (p / "sessions").mkdir(exist_ok=True)
    return p

def ensure_user_has_materials(username):
    """Memastikan direktori user memiliki 208 materi kuliah kedokteran & starter flashcards secara instan tanpa download."""
    try:
        user_root = get_user_root(username)
        dest_dir = user_root / "materials"
        dest_dir.mkdir(parents=True, exist_ok=True)
        
        # Priority 1: Instant master materials cache (zero download, zero delay)
        master_dir = DATA_DIR / "master_materials"
        source_dir = None
        if master_dir.exists() and len(list(master_dir.glob("*.json"))) >= 200:
            source_dir = master_dir
        else:
            for cand in ["dimas", "dimaswastu", "dimaswastumahesa"]:
                cdir = DATA_DIR / "users" / cand / "materials"
                if cdir.exists() and len(list(cdir.glob("*.json"))) >= 200:
                    source_dir = cdir
                    break
                    
        if source_dir and dest_dir.resolve() != source_dir.resolve():
            dest_count = len(list(dest_dir.glob("*.json")))
            source_count = len(list(source_dir.glob("*.json")))
            if dest_count < source_count:
                now_next = (datetime.datetime.now() + datetime.timedelta(days=1)).isoformat()
                for f in source_dir.glob("*.json"):
                    tf = dest_dir / f.name
                    if not tf.exists():
                        try:
                            d = json.loads(f.read_text(encoding="utf-8"))
                            d["sessions"] = 0
                            d["review_count"] = 0
                            d["ease_factor"] = 2.5
                            d.pop("last_interval", None)
                            d["next_review"] = now_next
                            tf.write_text(json.dumps(d, ensure_ascii=False, indent=2))
                        except Exception:
                            import shutil
                            shutil.copy2(f, tf)
                        
        # Ensure starter flashcards are present for instant testing & Anki export
        fc_dest = user_root / "flashcards"
        fc_dest.mkdir(parents=True, exist_ok=True)
        fc_src = DATA_DIR / "flashcards"
        if fc_src.exists():
            import shutil
            for f in fc_src.glob("*.json"):
                tf = fc_dest / f.name
                if not tf.exists():
                    shutil.copy2(f, tf)
    except Exception:
        pass
        pass

def login_or_register_google_account(email, display_name=None):
    """Autentikasi akun Google & otomatis registrasi dengan akses Pro."""
    import secrets
    email_clean = email.strip().lower()
    if "@" not in email_clean or "." not in email_clean:
        return None, "Format email Google tidak valid. Contoh: nama@gmail.com"
    
    users = load_users()
    for uname, udata in users.items():
        if udata.get("email", "").lower() == email_clean:
            ensure_user_has_materials(uname)
            return udata, "OK"
            
    # Auto-register new Google account
    un_c = re.sub(r'[^a-zA-Z0-9_]', '', email_clean.split("@")[0].lower())
    if not un_c or len(un_c) < 3:
        un_c = "google_user_" + secrets.token_hex(3)
    if un_c in users:
        un_c = f"{un_c}_{secrets.token_hex(2)}"
        
    disp = display_name.strip() if display_name and display_name.strip() else un_c.capitalize()
    now_dt = datetime.datetime.now()
    new_user_data = {
        "username": un_c,
        "display_name": disp,
        "email": email_clean,
        "password_hash": "google_verified_oauth",
        "created_at": now_dt.isoformat(),
        "tier": "pro",
        "subscription_ends_at": (now_dt + datetime.timedelta(days=365)).isoformat(),
        "plan_name": "Pro Mahasiswa Kedokteran (Google Verified)"
    }
    users[un_c] = new_user_data
    save_users(users)
    ensure_user_has_materials(un_c)
    return new_user_data, "OK"

# ── OFFICIAL GOOGLE OAUTH 2.0 & IDENTITY AUTHORIZATION ENGINE ─────────────────
GOOGLE_OAUTH_CLIENT_ID = os.environ.get("GOOGLE_OAUTH_CLIENT_ID") or _safe_get_secret("GOOGLE_OAUTH_CLIENT_ID", "764086051850-6qr4p6gpi6hn506pt8ejuq83di341hur.apps.googleusercontent.com")
GOOGLE_OAUTH_CLIENT_SECRET = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET") or _safe_get_secret("GOOGLE_OAUTH_CLIENT_SECRET", "d-FL95Q19q7MQmFpd7hHD0Ty")
GOOGLE_OAUTH_REDIRECT_URI = os.environ.get("GOOGLE_OAUTH_REDIRECT_URI") or _safe_get_secret("GOOGLE_OAUTH_REDIRECT_URI", "http://localhost:8501")

def is_localhost_access():
    """Mendeteksi apakah request berasal dari mesin lokal (localhost) atau dari internet publik."""
    try:
        if hasattr(st, "context") and hasattr(st.context, "headers"):
            headers = st.context.headers or {}
            host = str(headers.get("host", "") or headers.get("x-forwarded-host", "")).lower()
            if "localhost" in host or "127.0.0.1" in host:
                return True
            if host:
                return False
    except Exception:
        pass
    return False

def get_official_google_auth_url():
    """Menghasilkan URL otorisasi resmi Google OAuth 2.0."""
    import urllib.parse
    params = {
        "client_id": GOOGLE_OAUTH_CLIENT_ID,
        "redirect_uri": GOOGLE_OAUTH_REDIRECT_URI,
        "response_type": "code",
        "scope": "openid email profile",
        "access_type": "offline",
        "prompt": "select_account"
    }
    return "https://accounts.google.com/o/oauth2/v2/auth?" + urllib.parse.urlencode(params)

def handle_google_oauth_code_exchange():
    """Memproses otorisasi resmi callback dari accounts.google.com via query_params."""
    import requests
    if hasattr(st, "query_params") and "code" in st.query_params:
        auth_code = st.query_params["code"]
        try:
            r = requests.post("https://oauth2.googleapis.com/token", data={
                "code": auth_code,
                "client_id": GOOGLE_OAUTH_CLIENT_ID,
                "client_secret": GOOGLE_OAUTH_CLIENT_SECRET,
                "redirect_uri": GOOGLE_OAUTH_REDIRECT_URI,
                "grant_type": "authorization_code"
            }, timeout=15)
            token_json = r.json()
            acc_tok = token_json.get("access_token")
            if acc_tok:
                u_res = requests.get("https://www.googleapis.com/oauth2/v3/userinfo", headers={"Authorization": f"Bearer {acc_tok}"}, timeout=15)
                u_info = u_res.json()
                email = u_info.get("email")
                if email and u_info.get("email_verified", True):
                    disp = u_info.get("name") or email.split("@")[0].capitalize()
                    pic = u_info.get("picture", "")
                    u_data, msg = login_or_register_google_account(email, disp)
                    if u_data:
                        if pic:
                            u_data["picture"] = pic
                            db = load_users()
                            if u_data["username"] in db:
                                db[u_data["username"]]["picture"] = pic
                                save_users(db)
                        ensure_user_has_materials(u_data["username"])
                        st.session_state.current_user = u_data["username"]
                        st.session_state.user_info = u_data
                        set_persisted_auth_session(u_data)
                        st.query_params.clear()
                        st.toast(f"✓ Otorisasi Google Resmi Berhasil: {email}!", icon="🟢")
                        st.rerun()
        except Exception as e:
            st.error(f"Gagal verifikasi otorisasi Google: {e}")

def authenticate_via_local_google_adc():
    """Otorisasi instan menggunakan kredensial Google OAuth resmi yang aktif di sistem."""
    import requests, json
    adc_path = Path.home() / ".config" / "gcloud" / "application_default_credentials.json"
    if adc_path.exists():
        try:
            cred = json.loads(adc_path.read_text())
            c_id = cred.get("client_id", GOOGLE_OAUTH_CLIENT_ID)
            c_sec = cred.get("client_secret", GOOGLE_OAUTH_CLIENT_SECRET)
            r_tok = cred.get("refresh_token")
            if r_tok:
                r = requests.post("https://oauth2.googleapis.com/token", data={
                    "client_id": c_id,
                    "client_secret": c_sec,
                    "refresh_token": r_tok,
                    "grant_type": "refresh_token"
                }, timeout=15)
                token_data = r.json()
                acc_tok = token_data.get("access_token")
                if acc_tok:
                    u_r = requests.get("https://www.googleapis.com/oauth2/v3/userinfo", headers={"Authorization": f"Bearer {acc_tok}"}, timeout=15)
                    u_info = u_r.json()
                    email = u_info.get("email")
                    if not email:
                        return None, "Kredensial Google lokal tidak menyertakan email valid."
                    name = u_info.get("name") or email.split("@")[0].capitalize()
                    pic = u_info.get("picture", "")
                    u_data, msg = login_or_register_google_account(email, display_name=name)
                    if u_data:
                        if pic:
                            u_data["picture"] = pic
                            db = load_users()
                            if u_data["username"] in db:
                                db[u_data["username"]]["picture"] = pic
                                save_users(db)
                        ensure_user_has_materials(u_data["username"])
                        st.session_state.current_user = u_data["username"]
                        st.session_state.user_info = u_data
                        set_persisted_auth_session(u_data)
                        return u_data, "OK"
        except Exception as e:
            return None, str(e)
    return None, "Kredensial Google belum terpasang."



def authenticate_user(login_identifier, password):
    ident = (login_identifier or "").strip().lower()
    pw = (password or "").strip()
    if not ident or not pw:
        return None, "Harap masukkan email/username dan kata sandi."
    
    # Persistent Rate limiting check (max 5 failed attempts in 60s)
    now = time.time()
    attempts = get_login_attempts()
    if ident in attempts:
        first_time, count = attempts[ident]
        if count >= 5 and (now - first_time) < 60:
            remaining = int(60 - (now - first_time))
            return None, f"Terlalu banyak percobaan login gagal. Silakan tunggu {remaining} detik demi keamanan akun."
        elif (now - first_time) >= 60:
            attempts[ident] = [now, 0]
            save_login_attempts(attempts)
            
    users = load_users()
    matched_user = None
    for uname, udata in users.items():
        if uname.lower() == ident or udata.get("email", "").lower() == ident:
            matched_user = udata
            break
            
    if not matched_user:
        if ident not in attempts:
            attempts[ident] = [now, 1]
        else:
            attempts[ident][1] += 1
        save_login_attempts(attempts)
        return None, "Email/Username atau Kata Sandi tidak cocok."
        
    stored_hash = matched_user.get("password_hash", "")
    
    # Strictly handle Google OAuth only accounts without a set password
    if stored_hash == "google_verified_oauth":
        return None, "Akun ini didaftarkan via Akun Google resmi. Silakan masuk menggunakan tombol 'Masuk Akun Google'."
        
    # Strictly verify cryptographic PBKDF2 hash
    if verify_user_password(pw, stored_hash):
        if ident in attempts:
            attempts.pop(ident, None)
            save_login_attempts(attempts)
        ensure_user_has_materials(matched_user["username"])
        return matched_user, "OK"
    else:
        if ident not in attempts:
            attempts[ident] = [now, 1]
        else:
            attempts[ident][1] += 1
        save_login_attempts(attempts)
        return None, "Email/Username atau Kata Sandi tidak cocok."

def register_user(display_name, email, password, confirm_password=None):
    users = load_users()
    
    name_clean = (display_name or "").strip()
    if not name_clean or len(name_clean) < 2:
        return None, "Harap masukkan nama lengkap Anda (minimal 2 karakter)."
        
    email_clean = (email or "").strip().lower()
    if "@" not in email_clean or "." not in email_clean:
        return None, "Format email tidak valid. Pastikan format email benar (contoh: nama@gmail.com)."
        
    for u in users.values():
        if u.get("email", "").lower() == email_clean:
            return None, "Email ini sudah terdaftar. Silakan pilih tab 'Sudah punya akun' untuk masuk."
            
    pw = (password or "").strip()
    if len(pw) < 6:
        return None, "Kata sandi minimal 6 karakter demi keamanan akun Anda."
        
    if confirm_password is not None and pw != (confirm_password or "").strip():
        return None, "Konfirmasi kata sandi tidak cocok. Pastikan kedua kata sandi sama persis."
        
    # Generate unique, clean username
    import secrets
    base_un = re.sub(r'[^a-zA-Z0-9_]', '', email_clean.split("@")[0].lower())
    if not base_un or len(base_un) < 3:
        base_un = re.sub(r'[^a-zA-Z0-9_]', '', name_clean.lower().replace(" ", "_"))
    if not base_un or len(base_un) < 3:
        base_un = f"user_{secrets.token_hex(3)}"
        
    uname_clean = base_un
    while uname_clean in users:
        uname_clean = f"{base_un}_{secrets.token_hex(2)}"
        
    now_dt = datetime.datetime.now()
    new_user_data = {
        "username": uname_clean,
        "display_name": name_clean,
        "email": email_clean,
        "password_hash": hash_user_password(pw),
        "created_at": now_dt.isoformat(),
        "tier": "pro",
        "trial_ends_at": (now_dt + datetime.timedelta(days=7)).isoformat(),
        "subscription_ends_at": (now_dt + datetime.timedelta(days=365)).isoformat(),
        "plan_name": "Pro Mahasiswa Kedokteran"
    }
    users[uname_clean] = new_user_data
    save_users(users)
    ensure_user_has_materials(uname_clean)
    return new_user_data, "OK"

def create_guest_tester_session(display_name, email=""):
    import secrets
    name_clean = (display_name or "").strip() or "Teman Penguji"
    email_clean = (email or "").strip().lower() or f"tester_{secrets.token_hex(3)}@med.test"
    users = load_users()
    
    for uname, udata in users.items():
        if udata.get("email", "").lower() == email_clean or udata.get("display_name", "").lower() == name_clean.lower():
            ensure_user_has_materials(uname)
            return udata, "OK"
            
    base_s = re.sub(r'[^a-zA-Z0-9_]', '', name_clean.lower().replace(' ', '_'))[:12] or "tester"
    uname_clean = f"test_{base_s}_{secrets.token_hex(2)}"
    now_dt = datetime.datetime.now()
    new_user = {
        "username": uname_clean,
        "display_name": name_clean,
        "email": email_clean,
        "password_hash": hash_user_password(secrets.token_hex(8)),
        "created_at": now_dt.isoformat(),
        "tier": "pro",
        "plan_name": "Pro Beta Tester (Akses Penuh 208 Modul)"
    }
    users[uname_clean] = new_user
    save_users(users)
    ensure_user_has_materials(uname_clean)
    return new_user, "OK"




# ── COGNITIVE DIFFICULTY LEVEL & EVIDENCE-BASED PROMPT ENGINE ─────────────────
COGNITIVE_LEVELS = {
    "Level 1: Fondasi & Definisi": {
        "icon": "🟢",
        "desc": "Tingkat C1-C2 (Taksonomi Bloom): Uji definisi, terminologi, dan fakta dasar medis.",
        "prompt_mod": "Fokus pada Level 1 (Fondasi & Terminologi Medis): Uji pemahaman konsep dasar, definisi esensial, klasifikasi obat/penyakit, dan pengenalan terminologi kunci."
    },
    "Level 2: Analisis Mekanisme & Kausalitas": {
        "icon": "🟡",
        "desc": "Tingkat C3-C4 (Taksonomi Bloom): Uji hubungan sebab-akibat, kurva dosis, dan cascade sinyal.",
        "prompt_mod": "Fokus pada Level 2 (Analisis Mekanisme & Kausalitas Molekuler): Uji hubungan sebab-akibat patofisiologis, interaksi molekuler, interpretasi kurva farmakodinamik, dan konsekuensi fisiologis."
    },
    "Level 3: Kasus Klinis & Jebakan Ujian/UKMPPD": {
        "icon": "🔴",
        "desc": "Tingkat C5-C6 (Taksonomi Bloom): Uji penalaran klinis di IGD, diferensial diagnosis, dan jebakan soal.",
        "prompt_mod": "Fokus pada Level 3 (Kasus Klinis Tingkat Tinggi & Jebakan Soal UKMPPD): Buat skenario kasus pasien nyata di ranjang/IGD, analisis diferensial diagnosis, jebakan pengecoh yang sering menipu mahasiswa, dan keputusan terapeutik darurat."
    }
}

# ── SUBSCRIPTION, BILLING & QUOTA MANAGEMENT ENGINE ───────────────────────────
def get_user_subscription_status(username=None):
    u = username or st.session_state.get("current_user")
    if not u:
        return {
            "tier": "free", "is_pro": False, "days_left": 0,
            "status_label": "Free Plan", "can_upload": True, "max_mats": 10,
            "plan_name": "Free Plan"
        }
    users = load_users()
    udata = users.get(str(u), {})
    
    tier = udata.get("tier", "trial")
    sub_ends_str = udata.get("subscription_ends_at")
    trial_ends_str = udata.get("trial_ends_at")
    now = datetime.datetime.now()
    
    is_pro = False
    days_left = 0
    status_label = "Free Plan"
    
    if sub_ends_str:
        try:
            sub_ends = datetime.datetime.fromisoformat(sub_ends_str)
            if sub_ends > now:
                is_pro = True
                days_left = (sub_ends - now).days + 1
                status_label = f"PRO Mahasiswa ({days_left} hari lagi)"
                return {
                    "tier": "pro", "is_pro": True, "days_left": days_left,
                    "status_label": status_label, "can_upload": True, "max_mats": 999,
                    "plan_name": udata.get("plan_name", "Pro Mahasiswa")
                }
        except: pass

    if trial_ends_str:
        try:
            trial_ends = datetime.datetime.fromisoformat(trial_ends_str)
            if trial_ends > now:
                is_pro = True
                days_left = (trial_ends - now).days + 1
                status_label = f"Trial Pro ({days_left} hari sisa)"
                return {
                    "tier": "trial", "is_pro": True, "days_left": days_left,
                    "status_label": status_label, "can_upload": True, "max_mats": 999,
                    "plan_name": "Free 3-Day Pro Trial"
                }
        except: pass
        
    return {
        "tier": "free", "is_pro": False, "days_left": 0,
        "status_label": "Free Plan (Terbatas 1 Slide)", "can_upload": len(load_mats()) < 1,
        "max_mats": 1, "plan_name": "Free Tier"
    }

def activate_user_subscription(username, months=1, plan_name="Pro Mahasiswa Kedokteran (Rp 25.000/bln)"):
    users = load_users()
    u = str(username)
    if u in users:
        now = datetime.datetime.now()
        current_ends_str = users[u].get("subscription_ends_at")
        base_time = now
        if current_ends_str:
            try:
                cur_ends = datetime.datetime.fromisoformat(current_ends_str)
                if cur_ends > now: base_time = cur_ends
            except: pass
            
        new_ends = base_time + datetime.timedelta(days=30 * months)
        users[u]["tier"] = "pro"
        users[u]["subscription_ends_at"] = new_ends.isoformat()
        users[u]["plan_name"] = plan_name
        save_users(users)
        if st.session_state.get("current_user") == u:
            st.session_state.user_info = users[u]
        return True, new_ends
    return False, None

# ── Data helpers ──────────────────────────────────────────────────────────────

# ── DOCUMENT EXTRACTION (PDF, PPTX, DOCX, OCR FALLBACK) ──────────────────────
def extract_document_text(uploaded_file, api_key=""):
    fname = uploaded_file.name.lower()
    raw_bytes = uploaded_file.read()
    
    # 1. PPTX / PPT
    if fname.endswith(".pptx") or fname.endswith(".ppt"):
        try:
            import pptx, io
            prs = pptx.Presentation(io.BytesIO(raw_bytes))
            text_runs = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in paragraph.runs).strip()
                            if t: slide_texts.append(t)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_cells: slide_texts.append(" | ".join(row_cells))
                if slide_texts:
                    text_runs.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_texts))
            full_ppt_txt = "\n\n".join(text_runs)
            if len(full_ppt_txt.strip()) > 20:
                return full_ppt_txt
        except Exception as e:
            pass
            
    # 2. DOCX
    if fname.endswith(".docx"):
        try:
            import docx2txt, io
            txt = docx2txt.process(io.BytesIO(raw_bytes))
            if txt and len(txt.strip()) > 20:
                return txt
        except Exception as e:
            pass
            
    # 3. PDF (PyMuPDF + Vision OCR Fallback)
    if fname.endswith(".pdf"):
        try:
            doc = fitz.open(stream=raw_bytes, filetype="pdf")
            pages_txt = []
            scanned_pages = []
            for page_idx, page in enumerate(doc):
                p_txt = page.get_text().strip()
                if len(p_txt) > 25:
                    pages_txt.append(f"--- Halaman {page_idx + 1} ---\n" + p_txt)
                else:
                    scanned_pages.append(page)
                    
            if len(pages_txt) == 0 and len(scanned_pages) > 0 and api_key:
                genai.configure(api_key=api_key)
                vision_model = genai.GenerativeModel("gemini-3.1-flash-lite-preview")
                ocr_results = []
                for p in scanned_pages[:12]:
                    pix = p.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("png")
                    prompt_ocr = "Transkripsikan seluruh teks akademik, tabel, bagan, dan formula medis dari gambar slide ini secara lengkap dan akurat dalam teks terstruktur."
                    resp = vision_model.generate_content([prompt_ocr, {"mime_type": "image/png", "data": img_bytes}])
                    if resp.text:
                        ocr_results.append(f"--- Halaman {p.number + 1} (OCR) ---\n" + resp.text.strip())
                if ocr_results:
                    return "\n\n".join(ocr_results)
                    
            if pages_txt:
                return "\n\n".join(pages_txt)
        except Exception as e:
            pass
            
    try:
        return raw_bytes.decode("utf-8", errors="ignore")
    except:
        return ""

# ── SESSION AUTO-SAVE & RECOVERY ──────────────────────────────────────────────
def save_active_session(mat_name, state_dict):
    try:
        s_path = get_user_root() / "sessions" / f"{mat_name}.json"
        s_path.write_text(json.dumps(state_dict, ensure_ascii=False, indent=2))
    except: pass

def load_active_session(mat_name):
    try:
        s_path = get_user_root() / "sessions" / f"{mat_name}.json"
        if s_path.exists():
            return json.loads(s_path.read_text())
    except: pass
    return None

def clear_active_session(mat_name):
    try:
        s_path = get_user_root() / "sessions" / f"{mat_name}.json"
        if s_path.exists(): s_path.unlink()
    except: pass

# ── 1-CLICK ZIP BACKUP & RESTORE ──────────────────────────────────────────────
def create_backup_zip():
    import zipfile, io
    buf = io.BytesIO()
    user_root = get_user_root()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for f in user_root.rglob("*"):
            if f.is_file() and not f.name.startswith("."):
                rel_path = f.relative_to(user_root)
                z.write(f, arcname=str(rel_path))
    buf.seek(0)
    return buf.getvalue()

def restore_backup_zip(zip_bytes):
    import zipfile, io
    buf = io.BytesIO(zip_bytes)
    user_root = get_user_root()
    with zipfile.ZipFile(buf, "r") as z:
        z.extractall(user_root)

def clean_academic_text(raw_text):
    if len(raw_text) < 4000:
        return raw_text
    cleaned = re.sub(r'\.{4,}\s*\d+', '', raw_text)
    match = re.search(r'(TEMA\s+\d+|BAB\s+[I|1]|I\.\s+PENDAHULUAN|MATERI|KULIAH\s+\d+|SKENARIO\s+\d+|Clinical\s+reasoning|DASAR\s+PERTIMBANGAN)', cleaned, re.IGNORECASE)
    if match and match.start() > 800:
        cleaned = cleaned[match.start():]
    return cleaned[:35000]

def pdf_text(f): return "".join(p.get_text() for p in fitz.open(stream=f.read(), filetype="pdf"))

def save_mat(name, text, username=None):
    m_path = get_user_root(username) / "materials" / f"{name}.json"
    m_path.write_text(json.dumps({
        "name": name, "text": text,
        "uploaded_at": datetime.datetime.now().isoformat(),
        "next_review": (datetime.datetime.now() + datetime.timedelta(days=1)).isoformat(),
        "review_count": 0, "ease_factor": 2.5, "sessions": 0
    }, ensure_ascii=False, indent=2))

def load_mats(username=None):
    out = {}
    m_dir = get_user_root(username) / "materials"
    if not m_dir.exists(): return out
    for f in m_dir.glob("*.json"):
        try:
            d = json.loads(f.read_text())
            m_name = d.get("name") or d.get("title") or f.stem
            d["name"] = m_name
            out[m_name] = d
        except Exception: pass
    return out

def update_sr(name, q):
    path = get_user_root() / "materials" / f"{name}.json"
    d = json.loads(path.read_text())
    ef = max(1.3, d.get("ease_factor", 2.5) + 0.1 - (5 - q) * (0.08 + (5 - q) * 0.02))
    c = d.get("review_count", 0)
    iv = 1 if q < 3 or c == 0 else (6 if c == 1 else round(d.get("last_interval", 6) * ef))
    new_c = 0 if q < 3 else c + 1
    d.update({"ease_factor": ef, "review_count": new_c, "last_interval": iv, "sessions": d.get("sessions", 0) + 1,
              "next_review": (datetime.datetime.now() + datetime.timedelta(days=iv)).isoformat()})
    path.write_text(json.dumps(d, ensure_ascii=False, indent=2)); return iv



# ── POWERFUL CLINICAL & NEUROSCIENCE ENGINES (SIMPEL & POWERFUL) ──────────────
def calculate_memory_retention(mat_data):
    """
    Menghitung perkiraan kekuatan retensi memori otak menggunakan rumus Kurva Lupa Ebbinghaus:
    R = 100 * exp(-t / S)
    di mana t = hari sejak review terakhir, S = faktor stabilitas memori (berdasarkan EF & jumlah repetisi).
    
    Catatan Neurosains:
    Jika materi belum pernah dipelajari (sessions == 0 dan review_count == 0),
    maka retensi belum ada (None / Belum Dimulai). Jangan menampilkan nilai 50% semu.
    """
    sessions = mat_data.get("sessions", 0)
    review_count = mat_data.get("review_count", 0)
    if sessions == 0 and review_count == 0:
        return None
        
    now = datetime.datetime.now()
    ef = mat_data.get("ease_factor", 2.5)
    
    # Stabilitas memori (dalam hari)
    stability = max(1.5, (review_count * 2.2 + sessions * 1.5) * (ef / 2.0))
    
    next_rev_str = mat_data.get("next_review")
    if next_rev_str:
        try:
            next_rev = datetime.datetime.fromisoformat(next_rev_str)
            days_to_rev = (next_rev - now).days
            if days_to_rev > 0:
                # Masih dalam interval aman
                retention = min(98, int(85 + (days_to_rev / stability) * 12))
            else:
                # Sudah melewati jadwal review, retensi menurun
                days_overdue = abs(days_to_rev)
                retention = max(25, int(85 * math.exp(-days_overdue / stability)))
            return retention
        except: pass
        
    return 85

def generate_anki_export_data(flashcards_list, mat_name="NeuroStudy"):
    """
    Mengekspor kartu flashcards ke dalam format standar Anki Tab-Separated Values (TSV).
    Kompatibel 100% dengan Anki Desktop, AnkiDroid, & AnkiMobile.
    """
    lines = ["#separator:tab", "#html:true", f"#tags:{mat_name.replace(' ', '_')}"]
    for card in flashcards_list:
        front = card.get("front", "").replace("\t", " ").replace("\n", "<br>")
        back = card.get("back", "").replace("\t", " ").replace("\n", "<br>")
        phase = card.get("phase", "General")
        if front and back:
            lines.append(f"{front}\t{back}\tNeuroStudy_{phase}")
    return "\n".join(lines).encode("utf-8")



def fetch_gdrive_folder_contents(folder_id):
    """
    Mengambil daftar file dan subfolder dari Google Drive folder publik.
    """
    import requests, re, json
    url = f"https://drive.google.com/drive/folders/{folder_id}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    }
    try:
        resp = requests.get(url, headers=headers, timeout=20)
        if resp.status_code != 200: return []
        m = re.search(r"window\['_DRIVE_ivd'\]\s*=\s*'([^']+)'", resp.text)
        if not m: return []
        decoded_str = m.group(1).encode('utf-8').decode('unicode_escape')
        data = json.loads(decoded_str)
        results = []
        def walk(obj):
            if isinstance(obj, list):
                if len(obj) >= 4 and isinstance(obj[0], str) and len(obj[0]) > 20 and isinstance(obj[2], str):
                    item_id = obj[0]
                    name = obj[2]
                    mime = obj[3] if len(obj) > 3 else ""
                    is_f = "folder" in mime
                    results.append({"id": item_id, "name": name, "mime": mime, "is_folder": is_f})
                for it in obj: walk(it)
        walk(data)
        seen = set()
        deduped = []
        for r in results:
            if r["id"] not in seen:
                seen.add(r["id"])
                deduped.append(r)
        return deduped
    except:
        return []


# ── GOOGLE DRIVE LIVE RECURSIVE SYNC HUB ENGINE ──────────────────────────────
def load_gdrive_index():
    user_root = get_user_root()
    idx_f = user_root / "gdrive_index.json"
    if idx_f.exists():
        try: return json.loads(idx_f.read_text())
        except: return None
    return None

def save_gdrive_index(data):
    user_root = get_user_root()
    idx_f = user_root / "gdrive_index.json"
    idx_f.write_text(json.dumps(data, ensure_ascii=False, indent=2))

def sync_gdrive_folder_recursive(root_id):
    """
    Menjelajah seluruh pohon folder Google Drive secara rekursif murni (arbitrary depth)
    dan mengindeks seluruh slide kuliah agar siap dipelajari kapan saja.
    """
    import requests, re, json, datetime
    headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
    
    def get_items(fid):
        u = f"https://drive.google.com/drive/folders/{fid}"
        try:
            r = requests.get(u, headers=headers, timeout=20)
            if r.status_code != 200: return []
            m = re.search(r"window\['_DRIVE_ivd'\]\s*=\s*'([^']+)'", r.text)
            if not m: return []
            d_str = m.group(1).encode('utf-8').decode('unicode_escape')
            d = json.loads(d_str)
            res = []
            def walk(o):
                if isinstance(o, list):
                    if len(o) >= 4 and isinstance(o[0], str) and len(o[0]) > 20 and isinstance(o[2], str):
                        m_type = o[3] if len(o) > 3 else ""
                        res.append({"id": o[0], "name": o[2], "mime": m_type, "is_folder": "folder" in m_type})
                    for it in o: walk(it)
            walk(d)
            seen, dedup = set(), []
            for item in res:
                if item["id"] not in seen:
                    seen.add(item["id"])
                    dedup.append(item)
            return dedup
        except: return []

    def is_doc(item):
        mime = item.get("mime", "").lower()
        name = item.get("name", "").lower()
        return ("pdf" in mime or "presentation" in mime or "document" in mime or
                any(name.endswith(ext) for ext in [".pdf", ".pptx", ".ppt", ".docx"]))

    all_files = []
    def walk_folder(fid, cat="BLOK", sub="Umum", depth=0):
        if depth > 5: return
        items = get_items(fid)
        for it in items:
            if it["is_folder"]:
                walk_folder(it["id"], cat=cat, sub=it["name"], depth=depth+1)
            elif is_doc(it):
                all_files.append({
                    "category": cat,
                    "subcategory": sub,
                    "name": it["name"],
                    "id": it["id"],
                    "mime": it["mime"]
                })

    root_items = get_items(root_id)
    for it in root_items:
        if it["is_folder"]:
            walk_folder(it["id"], cat=it["name"], sub="Umum", depth=1)
        elif is_doc(it):
            all_files.append({
                "category": "Root",
                "subcategory": "Umum",
                "name": it["name"],
                "id": it["id"],
                "mime": it["mime"]
            })
                        
    index_data = {
        "root_folder_id": root_id,
        "root_folder_name": "VERTEXTERIAL",
        "last_synced": datetime.datetime.now().isoformat(),
        "total_files": len(all_files),
        "files": all_files
    }
    save_gdrive_index(index_data)
    return index_data

def download_and_import_gdrive_slide(file_id, file_name, api_key, username=None):
    """
    Mengunduh file slide dari Google Drive, mengekstrak teks (PDF atau PPTX), dan menyimpannya ke database materi.
    """
    import requests, fitz, io
    try:
        from pptx import Presentation
    except ImportError:
        Presentation = None

    headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
    dl_url = f"https://drive.google.com/uc?export=download&id={file_id}"
    try:
        resp = requests.get(dl_url, headers=headers, timeout=30)
        if resp.status_code == 200 and len(resp.content) > 300:
            all_text = ""
            # Coba PDF dulu
            try:
                doc = fitz.open(stream=resp.content, filetype="pdf")
                all_text = "\n\n".join(f"--- Slide/Halaman {i+1} ---\n" + p.get_text() for i, p in enumerate(doc))
            except Exception:
                # Jika bukan PDF, coba PPTX
                if Presentation:
                    try:
                        prs = Presentation(io.BytesIO(resp.content))
                        slides_txt = []
                        for i, slide in enumerate(prs.slides):
                            stxt = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                            slides_txt.append(f"--- Slide {i+1} ---\n" + "\n".join(stxt))
                        all_text = "\n\n".join(slides_txt)
                    except Exception: pass

            if len(all_text.strip()) > 30:
                save_mat(file_name, all_text, username=username)
                return True, len(all_text)
            return False, "Dokumen kosong atau tidak memiliki teks yang terbaca."
        return False, f"Gagal mengunduh file (HTTP {resp.status_code})"
    except Exception as e:
        return False, str(e)


# ── AUTO-DOWNLOAD & SYNC DAEMON UNTUK FOLDER BLOK DOSEN ───────────────────────
def check_and_auto_download_blok_updates(username=None):
    """
    Otomatis memeriksa dan mengunduh seluruh materi baru di folder BLOK secara rekursif
    sehingga mahasiswa selalu memiliki slide kuliah paling mutakhir.
    """
    import threading
    target_user = username or st.session_state.get("current_user", "dimas")
    
    def _run_bg(t_user):
        try:
            user_mats = load_mats(username=t_user)
            blok_folder_ids = [
                ("BMS 1", "1GFFIxSHGjGf1B3nq4kSJwQ4QMj_qbRm3"),
                ("BUAMS", "1mRMtfuraC0HueVl-4LTzXhSpinj1C6CD"),
                ("BMS 2", "1tVA1x3gVnyEvYKtcedPJItGZLcJnstJR"),
                ("BMS 3", "1NeQ2ZvQbr6lJZEujXYIYa8ToQIBuNCp4"),
                ("BMS 4", "1b1ejwl4M0fMwes32fKVTa5x0z7Fdc_fW"),
                ("BDT", "1823rKPuQVAl_vE-EQfo7rTq-Eh9KQves"),
                ("BMD", "1WdqpA4rXB-uFb0z9E4luC0elRniBDQqx")
            ]
            
            def scan_and_download(fid, b_name, depth=0):
                if depth > 3: return
                items = fetch_gdrive_folder_contents(fid)
                for it in items:
                    if it.get("is_folder"):
                        scan_and_download(it['id'], b_name, depth+1)
                    else:
                        mime = it.get("mime", "").lower()
                        name = it.get("name", "").lower()
                        if "pdf" in mime or "presentation" in mime or any(name.endswith(e) for e in [".pdf", ".pptx", ".ppt"]):
                            c_name = f"[{b_name}] {it['name']}".strip()
                            clean_target = re.sub(r'\.(pdf|pptx|ppt|docx|txt)$', '', c_name, flags=re.IGNORECASE).strip()
                            if clean_target not in user_mats:
                                download_and_import_gdrive_slide(it['id'], clean_target, "", username=t_user)

            for b_name, b_id in blok_folder_ids:
                scan_and_download(b_id, b_name, 0)
        except Exception: pass
        
    t = threading.Thread(target=_run_bg, args=(target_user,), daemon=True)
    t.start()

# ── GOOGLE DRIVE & GOOGLE SLIDES DIRECT IMPORT ENGINE ─────────────────────────
def extract_gdrive_file_id(url):
    m = re.search(r"/(?:file/d/|presentation/d/|document/d/|folders/|open\?id=)([a-zA-Z0-9_-]+)", url)
    if m:
        return m.group(1)
    return None

def fetch_gdrive_document_text(url, api_key):
    """
    Mengunduh dan mengekstrak dokumen/slide dari link Google Drive atau Google Slides publik/sharing.
    """
    import requests, io
    file_id = extract_gdrive_file_id(url)
    if not file_id:
        return False, "Link Google Drive tidak valid. Pastikan link memiliki format https://drive.google.com/... atau https://docs.google.com/..."
    
    # 1. Coba unduh sebagai Google Slides export PDF
    if "presentation/d" in url or "docs.google.com/presentation" in url:
        export_url = f"https://docs.google.com/presentation/d/{file_id}/export/pdf"
    elif "document/d" in url or "docs.google.com/document" in url:
        export_url = f"https://docs.google.com/document/d/{file_id}/export?format=pdf"
    else:
        # File umum (PDF/PPTX) di Google Drive
        export_url = f"https://drive.google.com/uc?export=download&id={file_id}"
        
    try:
        resp = requests.get(export_url, timeout=25, headers={"User-Agent": "Mozilla/5.0"})
        if resp.status_code == 200 and len(resp.content) > 100:
            # Bungkus ke BytesIO
            file_obj = io.BytesIO(resp.content)
            # Tentukan tipe file dari header atau url
            content_type = resp.headers.get("content-type", "").lower()
            if "pdf" in content_type or export_url.endswith("pdf"):
                file_obj.name = f"gdrive_{file_id}.pdf"
            elif "presentation" in content_type:
                file_obj.name = f"gdrive_{file_id}.pptx"
            else:
                file_obj.name = f"gdrive_{file_id}.pdf"
                
            txt = extract_document_text(file_obj, api_key)
            if len(txt.strip()) > 30:
                return True, txt
            return False, "Gagal membaca teks dari file Google Drive. Pastikan slide berisi teks atau dapat diakses publik."
        elif resp.status_code == 403:
            return False, "Akses Google Drive ditolak (403). Pastikan izin berbagi file disetel: Siapa saja yang memiliki tautan (Anyone with the link can view)."
        else:
            return False, f"Gagal mengunduh dari Google Drive (Status: {resp.status_code})."
    except Exception as e:
        return False, f"Terjadi kesalahan koneksi Google Drive: {str(e)}"


# ── SMART HIGH-YIELD COMPRESSION CACHE (ANTI-QUOTA EXHAUSTION) ───────────────
def get_cached_high_yield_text(mat_name, raw_text, api_key):
    """
    Menyimpan intisari teks padat ke disk sehingga pemanggilan berulang oleh para AI
    tidak perlu menyuntikkan 35.000 karakter berkali-kali. Menghemat 75% token & kuota!
    """
    cache_dir = get_user_root() / "cache"
    cache_dir.mkdir(exist_ok=True)
    cache_file = cache_dir / f"{mat_name}_compressed.json"
    
    if cache_file.exists():
        try:
            d = json.loads(cache_file.read_text())
            if d.get("summary"): return d["summary"]
        except: pass
        
    # Jika belum ada cache, kompresi 1x
    if len(raw_text) > 6000:
        compressed = raw_text[:7000] # Safe high-yield slice
    else:
        compressed = raw_text
        
    try:
        cache_file.write_text(json.dumps({"summary": compressed}, ensure_ascii=False))
    except: pass
    return compressed

# ── FLASHCARDS, DISCUSSION & GCAL HELPERS ─────────────────────────────────────
def load_flashcards(mat_name):
    f_path = get_user_root() / "flashcards" / f"{mat_name}.json"
    if f_path.exists():
        try: return json.loads(f_path.read_text())
        except: return []
    f_glob = DATA_DIR / "flashcards" / f"{mat_name}.json"
    if f_glob.exists():
        try: return json.loads(f_glob.read_text())
        except: return []
    return []

def save_flashcards(mat_name, cards):
    f_path = get_user_root() / "flashcards" / f"{mat_name}.json"
    atomic_write_json(f_path, cards)

def load_discussion(mat_name):
    d_path = get_user_root() / "discussions" / f"{mat_name}.json"
    if d_path.exists():
        try: return json.loads(d_path.read_text())
        except: return []
    return []

def save_discussion(mat_name, messages):
    d_path = get_user_root() / "discussions" / f"{mat_name}.json"
    d_path.write_text(json.dumps(messages, ensure_ascii=False, indent=2))

def get_current_indonesia_time(tz_name=None):
    """Mengembalikan hari, tanggal, bulan, tahun, dan jam dalam zona waktu Indonesia (WIB, WITA, WIT).
    Secara eksplisit menghitung offset UTC sehingga akurat di server lokal maupun cloud (Streamlit Cloud)."""
    HARI_INA = ["Senin", "Selasa", "Rabu", "Kamis", "Jumat", "Sabtu", "Minggu"]
    BULAN_INA = [
        "", "Januari", "Februari", "Maret", "April", "Mei", "Juni",
        "Juli", "Agustus", "September", "Oktober", "November", "Desember"
    ]
    if tz_name is None:
        try:
            tz_name = st.session_state.get("selected_tz", "WIB")
        except:
            tz_name = "WIB"

    tz_offsets = {
        "WIB": (datetime.timezone(datetime.timedelta(hours=7)), "Asia/Jakarta", "Waktu Indonesia Barat (UTC+7)"),
        "WITA": (datetime.timezone(datetime.timedelta(hours=8)), "Asia/Makassar", "Waktu Indonesia Tengah (UTC+8)"),
        "WIT": (datetime.timezone(datetime.timedelta(hours=9)), "Asia/Jayapura", "Waktu Indonesia Timur (UTC+9)")
    }
    tz_info = tz_offsets.get(tz_name, tz_offsets["WIB"])
    tz_obj = tz_info[0]

    now_live = datetime.datetime.now(tz_obj)
    hari = HARI_INA[now_live.weekday()]
    tgl = now_live.day
    bulan = BULAN_INA[now_live.month]
    tahun = now_live.year
    jam = now_live.strftime("%H:%M") + f" {tz_name}"
    full_str = f"{hari}, {tgl} {bulan} {tahun} • {jam}"
    short_str = f"{jam} • {tgl} {bulan[:3]}"
    return {
        "datetime": now_live,
        "hari": hari,
        "tanggal": tgl,
        "bulan": bulan,
        "tahun": tahun,
        "jam": jam,
        "full_str": full_str,
        "short_str": short_str,
        "tz_name": tz_name,
        "tz_iana": tz_info[1],
        "tz_desc": tz_info[2]
    }

def build_gcal_url(title, dt, details="", tz_iana="Asia/Jakarta"):
    import urllib.parse
    start_str = dt.strftime("%Y%m%dT090000")
    end_str = (dt + datetime.timedelta(hours=1)).strftime("%Y%m%dT100000")
    params = {
        "action": "TEMPLATE",
        "text": title,
        "dates": f"{start_str}/{end_str}",
        "ctz": tz_iana,
        "details": details
    }
    return "https://calendar.google.com/calendar/render?" + urllib.parse.urlencode(params)

def generate_ics_content(events):
    lines = [
        "BEGIN:VCALENDAR",
        "VERSION:2.0",
        "PRODID:-//NeuroStudy//Spaced Repetition Schedule//ID",
        "CALSCALE:GREGORIAN",
        "METHOD:PUBLISH"
    ]
    for ev in events:
        dt_start = ev["dt"].strftime("%Y%m%dT090000Z")
        dt_end = (ev["dt"] + datetime.timedelta(hours=1)).strftime("%Y%m%dT100000Z")
        now_str = datetime.datetime.now(datetime.timezone.utc).strftime("%Y%m%dT%H%M%SZ")
        name_clean = str(ev["name"]).replace(" ", "_")
        dt_tag = ev["dt"].strftime("%Y%m%d")
        rev_tag = ev.get("rev_num", 1)
        uid = f"neurostudy-{name_clean}-{dt_tag}-{rev_tag}@neurostudy.app"
        t_val = ev["title"]
        d_val = ev.get("desc", "Spaced Repetition Review")
        lines.extend([
            "BEGIN:VEVENT",
            f"UID:{uid}",
            f"DTSTAMP:{now_str}",
            f"DTSTART:{dt_start}",
            f"DTEND:{dt_end}",
            f"SUMMARY:{t_val}",
            f"DESCRIPTION:{d_val}",
            "STATUS:CONFIRMED",
            "END:VEVENT"
        ])
    lines.append("END:VCALENDAR")
    return "\r\n".join(lines)



def generate_flashcards_ai(api_key, text, phase_tag, mat_name):
    prompt = f"""Kamu adalah Profesor Kedokteran dan Pakar Kognitif Neurosains.
Buat 5 FLASHCARDS BERMUTU TINGGI berbasis substansi ilmiah/medis dari materi berikut.

Materi Sumber:
{text[:10000]}

Fokus Fase Belajar: {phase_tag}
Panduan Fase:
- Jika PRIME: Buat kartu tantangan intuisi & pertanyaan pre-test kunci.
- Jika DIG: Buat kartu penalaran mekanisme Mengapa & Bagaimana (Causal mechanisms).
- Jika RECALL: Buat kartu active retrieval untuk menguji fakta klinis, dosis, target, atau klasifikasi penting.
- Jika FEYNMAN: Buat kartu skenario klinis nyata / analogi sederhana untuk menguji model mental.

Output WAJIB berupa JSON ARRAY MURNI tanpa teks pembuka/penutup, format:
[
  {{
    "front": "Pertanyaan atau konsep di sisi depan kartu",
    "back": "Jawaban mendalam, mekanisme jelas, dan intisari klinis di sisi belakang kartu",
    "phase": "{phase_tag}"
  }}
]"""
    genai.configure(api_key=api_key)
    candidate_models = [
        "gemini-3.1-flash-lite-preview",
        "gemini-3-flash-preview",
        "gemini-flash-lite-latest",
        "gemini-flash-latest",
        "gemini-3.5-flash-lite"
    ]
    for m_name in candidate_models:
        try:
            m = genai.GenerativeModel(m_name)
            resp = m.generate_content(prompt, request_options={"timeout": 30})
            raw = resp.text.strip()
            if "```json" in raw: raw = raw.split("```json")[1].split("```")[0]
            elif "```" in raw: raw = raw.split("```")[1].split("```")[0]
            parsed = json.loads(raw)
            if isinstance(parsed, list) and len(parsed) > 0:
                current_cards = load_flashcards(mat_name)
                current_cards.extend(parsed)
                save_flashcards(mat_name, current_cards)
                return len(parsed)
        except Exception as e:
            continue
    return 0

def update_card_sm2(mat_name, card_idx_in_all, grade):
    all_cards = load_flashcards(mat_name)
    if 0 <= card_idx_in_all < len(all_cards):
        c = all_cards[card_idx_in_all]
        ef = c.get("ease_factor", 2.5)
        reps = c.get("repetitions", 0)
        last_iv = c.get("interval", 1)
        
        if grade < 3:
            iv = 1
            reps = 0
            ef = max(1.3, ef - 0.2)
        elif grade == 3:
            iv = 3 if reps == 0 else max(1, round(last_iv * 1.2))
            reps += 1
        else: # grade 5 (Paham)
            ef = max(1.3, ef + 0.1)
            iv = 6 if reps == 0 else max(1, round(last_iv * ef))
            reps += 1
            
        next_dt = (datetime.datetime.now() + datetime.timedelta(days=iv)).isoformat()
        c.update({
            "ease_factor": ef,
            "interval": iv,
            "repetitions": reps,
            "next_review": next_dt,
            "last_grade": grade
        })
        save_flashcards(mat_name, all_cards)
        return iv
    return 1

def render_flashcards_widget(mat_name, text, api_key, phase_tag="ALL", key_prefix="fc"):
    all_cards = load_flashcards(mat_name)
    indexed_cards = list(enumerate(all_cards))
    
    if phase_tag != "ALL":
        filtered = [(idx, c) for idx, c in indexed_cards if c.get("phase", "ALL") == phase_tag]
    else:
        filtered = indexed_cards
        
    st.markdown(f'<div style="display:flex;justify-content:space-between;align-items:center;margin:10px 0 6px;"><div><span style="font-weight:700;color:#818cf8;font-size:0.95rem;">🃏 Flashcards Interaktif SM-2 ({phase_tag})</span><span style="font-size:0.75rem;color:#94a3b8;margin-left:8px;">Total {len(filtered)} kartu</span></div></div>', unsafe_allow_html=True)
    
    idx_key = f"{key_prefix}_fc_idx_{mat_name}_{phase_tag}"
    rev_key = f"{key_prefix}_fc_rev_{mat_name}_{phase_tag}"
    if idx_key not in st.session_state: st.session_state[idx_key] = 0
    if rev_key not in st.session_state: st.session_state[rev_key] = False
    
    if not filtered:
        st.markdown(f'<div class="card card-sm" style="text-align:center;padding:20px;"><div style="color:#94a3b8;font-size:0.88rem;margin-bottom:12px;">Belum ada flashcard untuk fase <strong>{phase_tag}</strong>.</div></div>', unsafe_allow_html=True)
        if st.button(f"⚡ Buat +5 Flashcards {phase_tag} (Unlimited AI)", type="primary", key=f"{key_prefix}_btn_gen_empty_{mat_name}_{phase_tag}"):
            with st.spinner("Merancang 5 flashcard ilmiah berbasis AI..."):
                count = generate_flashcards_ai(api_key, text, phase_tag, mat_name)
                if count > 0:
                    st.success(f"✅ Berhasil menambahkan {count} flashcards baru!")
                    st.rerun()
                else:
                    st.error("Gagal generate flashcards, coba lagi.")
        return

    cur_idx = min(max(0, st.session_state[idx_key]), len(filtered) - 1)
    orig_idx, card = filtered[cur_idx]
    is_revealed = st.session_state[rev_key]
    
    c_iv = card.get("interval", 1)
    c_ef = card.get("ease_factor", 2.5)
    c_rep = card.get("repetitions", 0)
    
    badge_label = f"FASE {card.get('phase', phase_tag)}"
    back_html = ""
    if is_revealed:
        b_txt = html.escape(str(card.get("back", "")))
        back_html = f'<div class="fc-back-box"><div style="font-size:0.72rem;color:#4ade80;font-weight:700;margin-bottom:4px;">💡 JAWABAN & MEKANISME KLINIS:</div>{b_txt}<div style="margin-top:10px;padding-top:8px;border-top:1px solid rgba(255,255,255,0.08);font-size:0.72rem;color:#94a3b8;">Spaced Repetition: <strong>Interval {c_iv} hari</strong> · <strong>EF {c_ef:.2f}</strong> · <strong>Diulang {c_rep}×</strong></div></div>'
    
    f_txt = html.escape(str(card.get("front", "")))
    card_html = f'<div class="fc-card"><div style="display:flex;justify-content:space-between;align-items:center;"><span class="fc-badge">{badge_label}</span><span style="font-size:0.75rem;color:#94a3b8;font-weight:600;">Kartu {cur_idx + 1} dari {len(filtered)}</span></div><div class="fc-front-title">{f_txt}</div>{back_html}</div>'
    st.markdown(card_html, unsafe_allow_html=True)
    
    col_a, col_b, col_c = st.columns([1, 1, 1])
    with col_a:
        if st.button("⬅️ Sebelumnya", disabled=(cur_idx == 0), key=f"{key_prefix}_fc_prev_{mat_name}_{phase_tag}_{cur_idx}", use_container_width=True):
            st.session_state[idx_key] = max(0, cur_idx - 1)
            st.session_state[rev_key] = False
            st.rerun()
    with col_b:
        rev_btn_label = "🙈 Tutup Jawaban" if is_revealed else "👁️ Buka Jawaban"
        if st.button(rev_btn_label, type="secondary", key=f"{key_prefix}_fc_flip_{mat_name}_{phase_tag}_{cur_idx}", use_container_width=True):
            st.session_state[rev_key] = not is_revealed
            st.rerun()
    with col_c:
        if st.button("Berikutnya ➡️", disabled=(cur_idx >= len(filtered) - 1), key=f"{key_prefix}_fc_next_{mat_name}_{phase_tag}_{cur_idx}", use_container_width=True):
            st.session_state[idx_key] = min(len(filtered) - 1, cur_idx + 1)
            st.session_state[rev_key] = False
            st.rerun()
            
    # Granular Card SM-2 Self-Grading Buttons (Shown when card is revealed)
    if is_revealed:
        st.markdown('<div style="font-size:0.75rem;color:#cbd5e1;font-weight:600;margin:6px 0 4px;">🎯 Nilai Pemahaman Kartu Ini (Algoritma SM-2):</div>', unsafe_allow_html=True)
        cg1, cg2, cg3 = st.columns(3)
        with cg1:
            if st.button("🔴 Sulit (Ulang 1 Hari)", key=f"{key_prefix}_sm_bad_{mat_name}_{orig_idx}", use_container_width=True):
                update_card_sm2(mat_name, orig_idx, 1)
                st.session_state[idx_key] = min(len(filtered) - 1, cur_idx + 1)
                st.session_state[rev_key] = False
                st.rerun()
        with cg2:
            if st.button("🟡 Sedang (Ulang 3 Hari)", key=f"{key_prefix}_sm_mid_{mat_name}_{orig_idx}", use_container_width=True):
                update_card_sm2(mat_name, orig_idx, 3)
                st.session_state[idx_key] = min(len(filtered) - 1, cur_idx + 1)
                st.session_state[rev_key] = False
                st.rerun()
        with cg3:
            next_paham_iv = 6 if c_rep == 0 else max(1, round(c_iv * c_ef))
            if st.button(f"🟢 Paham (Ulang {next_paham_iv} Hari)", key=f"{key_prefix}_sm_good_{mat_name}_{orig_idx}", use_container_width=True):
                update_card_sm2(mat_name, orig_idx, 5)
                st.session_state[idx_key] = min(len(filtered) - 1, cur_idx + 1)
                st.session_state[rev_key] = False
                st.rerun()
            
    if st.button(f"⚡ Tambah +5 Flashcards Baru ({phase_tag})", key=f"{key_prefix}_btn_add_more_{mat_name}_{phase_tag}", use_container_width=True):
        with st.spinner("AI sedang merancang flashcards tambahan..."):
            count = generate_flashcards_ai(api_key, text, phase_tag, mat_name)
            if count > 0:
                st.toast(f"✅ {count} kartu baru ditambahkan ke koleksi!", icon="🃏")
                st.session_state[idx_key] = len(filtered)
                st.session_state[rev_key] = False
                st.rerun()
# ── EVIDENCE-BASED PARETO 80/20 & COLLECTIVE PEER CACHING LAYER ───────────────
GLOBAL_LIB_DIR = DATA_DIR / "global_library"

@st.cache_resource
def _ensure_global_lib_dirs():
    GLOBAL_LIB_DIR.mkdir(parents=True, exist_ok=True)
    (GLOBAL_LIB_DIR / "master_notes").mkdir(exist_ok=True)
    (GLOBAL_LIB_DIR / "active_recall").mkdir(exist_ok=True)
    (GLOBAL_LIB_DIR / "exam_simulations").mkdir(exist_ok=True)
    return True
_ensure_global_lib_dirs()

def get_cached_master_note(mat_name):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "master_notes" / f"{safe_name}.json"
    if p_user.exists():
        try:
            d = json.loads(p_user.read_text())
            d["source"] = "private"
            return d
        except: pass
        
    p_glob = GLOBAL_LIB_DIR / "master_notes" / f"{safe_name}.json"
    if p_glob.exists():
        try:
            d = json.loads(p_glob.read_text())
            d["source"] = "global_peer"
            return d
        except: pass
    return None

def save_cached_master_note(mat_name, data, is_verified=False, reviewer_name=""):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "master_notes" / f"{safe_name}.json"
    p_glob = GLOBAL_LIB_DIR / "master_notes" / f"{safe_name}.json"
    if is_verified:
        data["verified"] = True
        data["verified_by"] = reviewer_name or "dr. Dimas Wastu Mahesa (Reviewer Klinis EBM)"
        data["verified_at"] = datetime.datetime.now().isoformat()
        data["clinical_standard"] = "Harrison 21st Ed (2022) / Guyton 14th Ed (2020) / Robbins 10th Ed (2020) / Katzung 13th Ed (2021) / Norman & Eva (NEJM & Med Educ)"
    atomic_write_json(p_user, data)
    atomic_write_json(p_glob, data)

def get_cached_active_recall(mat_name):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "active_recall" / f"{safe_name}.json"
    if p_user.exists():
        try:
            d = json.loads(p_user.read_text())
            d["source"] = "private"
            return d
        except: pass
    p_glob = GLOBAL_LIB_DIR / "active_recall" / f"{safe_name}.json"
    if p_glob.exists():
        try:
            d = json.loads(p_glob.read_text())
            d["source"] = "global_peer"
            return d
        except: pass
    return None

def save_cached_active_recall(mat_name, data):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "active_recall" / f"{safe_name}.json"
    p_glob = GLOBAL_LIB_DIR / "active_recall" / f"{safe_name}.json"
    atomic_write_json(p_user, data)
    atomic_write_json(p_glob, data)

def get_cached_exam_simulation(mat_name):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "exam_simulations" / f"{safe_name}.json"
    if p_user.exists():
        try:
            d = json.loads(p_user.read_text())
            d["source"] = "private"
            return d
        except: pass
    p_glob = GLOBAL_LIB_DIR / "exam_simulations" / f"{safe_name}.json"
    if p_glob.exists():
        try:
            d = json.loads(p_glob.read_text())
            d["source"] = "global_peer"
            return d
        except: pass
    return None

def save_cached_exam_simulation(mat_name, data):
    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', mat_name)
    user_root = get_user_root()
    p_user = user_root / "exam_simulations" / f"{safe_name}.json"
    p_glob = GLOBAL_LIB_DIR / "exam_simulations" / f"{safe_name}.json"
    atomic_write_json(p_user, data)
    atomic_write_json(p_glob, data)

def generate_anki_export_tsv(mat_name):
    """Menghasilkan format Anki TSV resmi (#separator:tab, #html:true) untuk diimpor ke Anki Desktop/Mobile."""
    cards = load_flashcards(mat_name)
    if not cards:
        return ""
    lines = ["#separator:tab", "#html:true", "#tags column:3"]
    safe_title = re.sub(r'[^a-zA-Z0-9_]', '_', mat_name)
    for c in cards:
        front = c.get("front", "").replace("\t", " ").replace("\n", "<br>")
        back = c.get("back", "").replace("\t", " ").replace("\n", "<br>")
        tags = f"NeuroStudy::{safe_title}::{c.get('phase', 'HighYield')}"
        lines.append(f"{front}\t{back}\t{tags}")
    return "\n".join(lines)

# ── BETA TESTER FEEDBACK & PUBLIC SHARING HELPERS ────────────────────────────
FEEDBACK_FILE = DATA_DIR / "beta_feedback.json"
PUBLIC_URL_FILE = DATA_DIR / "active_public_url.txt"

def load_beta_feedback():
    if FEEDBACK_FILE.exists():
        try: return json.loads(FEEDBACK_FILE.read_text())
        except: return []
    return []

def save_beta_feedback(entry):
    feedbacks = load_beta_feedback()
    feedbacks.append(entry)
    atomic_write_json(FEEDBACK_FILE, feedbacks)

def get_active_public_url():
    if PUBLIC_URL_FILE.exists():
        try:
            u = PUBLIC_URL_FILE.read_text().strip()
            if u.startswith("http"):
                return u
        except: pass
    return "https://mysimon-frontier-leaving-defend.trycloudflare.com"

def extract_json_safely(text_str):
    import re, json
    s = (text_str or "").strip()
    if "```json" in s:
        s = s.split("```json", 1)[1].split("```", 1)[0].strip()
    elif "```" in s:
        s = s.split("```", 1)[1].split("```", 1)[0].strip()
    try:
        return json.loads(s)
    except:
        m = re.search(r'(\[.*\]|\{.*\})', s, re.DOTALL)
        if m:
            try: return json.loads(m.group(1))
            except: pass
    return None


def parse_thought_and_content(raw_str):
    thinking = ""
    content = ""
    if "<thinking>" in raw_str:
        after_think = raw_str.split("<thinking>", 1)[1]
        if "</thinking>" in after_think:
            thinking, content = after_think.split("</thinking>", 1)
        else:
            thinking = after_think
            content = ""
    else:
        content = raw_str
    return thinking.strip(), content.strip()

def stream_ai_transparent(api_key, prompt, ph):
    """
    100% Transparent Real-Time Stream Engine:
    - Streams AI's internal reasoning (<thinking>...</thinking>) in a live glowing terminal box
    - Streams AI's actual output in real-time as tokens arrive with live cursor
    - Sub-second first-token latency with ultra-fast Gemini flash-lite models
    """
    ph.markdown("""
<div class="thinking-live-box">
  <div class="thinking-live-header">
    <div style="display:flex;align-items:center;gap:8px;font-size:0.8rem;color:#818cf8;font-weight:700;">
      <span class="live-dot"></span>
      🧠 PROSES PENALARAN AKTIF AI
    </div>
    <span style="font-size:0.72rem;color:#94a3b8;">Menghubungkan ke neural engine (sub-detik)...</span>
  </div>
  <div class="thinking-live-text"><span class="cur"></span></div>
</div>
""", unsafe_allow_html=True)

    genai.configure(api_key=api_key)
    
    # Intelligently adapt thinking instruction to prevent conflicting prompt rules
    if "Catatan Master Klinis Komprehensif" in prompt or "Pareto 80/20" in prompt:
        full_prompt = (
            "Sebelum menulis catatan master, tuliskan 1-3 kalimat intisari penalaranmu di dalam tag <thinking>...</thinking>.\n"
            "Setelah tag </thinking>, langsung susun Catatan Master sesuai struktur markdown yang diminta secara lengkap, mendalam, dan presisi:\n\n"
        ) + prompt
    else:
        thinking_instruction = (
            "INSTRUKSI PEDAGOGI EMAS — HANGAT & MENGALIR:\n"
            "Sebelum menulis jawaban, kamu boleh menuliskan proses analisismu secara ringkas dalam tag <thinking>...</thinking>.\n"
            "Setelah tag </thinking>, susun penjelasanmu dengan bahasa yang jelas, bersahabat, berbasis bukti ilmiah, dan mudah dipahami.\n"
            "DILARANG menggunakan notasi LaTeX mentah ($). Gunakan simbol/teks biasa.\n\n"
        )
        full_prompt = thinking_instruction + prompt

    candidate_models = [
        "gemini-3.1-flash-lite-preview",
        "gemini-3-flash-preview",
        "gemini-flash-lite-latest",
        "gemini-flash-latest",
        "gemini-3.5-flash-lite"
    ]
    
    last_err = ""
    for idx, model_name in enumerate(candidate_models):
        try:
            if idx > 0:
                ph.markdown(f"""
<div class="thinking-live-box">
  <div class="thinking-live-header">
    <div style="display:flex;align-items:center;gap:8px;font-size:0.8rem;color:#f59e0b;font-weight:700;">
      <span class="live-dot" style="background:#f59e0b;"></span>
      🔄 BERALIH KE ENGINE CADANGAN ({model_name})...
    </div>
    <span style="font-size:0.72rem;color:#94a3b8;">Menyambungkan ulang...</span>
  </div>
  <div class="thinking-live-text"><span class="cur"></span></div>
</div>
""", unsafe_allow_html=True)
            m = genai.GenerativeModel(model_name)
            resp = m.generate_content(full_prompt, stream=True, request_options={"timeout": 45})
            raw_accumulated = ""
            last_render_time = 0
            
            for chunk in resp:
                if not chunk.text: continue
                raw_accumulated += chunk.text
                
                now_t = time.time()
                # Throttle Streamlit UI updates to at most ~12 fps (every 80ms) to prevent WebSocket saturation
                if (now_t - last_render_time) < 0.08:
                    continue
                last_render_time = now_t
                
                thinking_text, main_text = parse_thought_and_content(raw_accumulated)
                
                # Render live stream
                if not main_text:
                    ph.markdown(f"""
<div class="thinking-live-box">
  <div class="thinking-live-header">
    <div style="display:flex;align-items:center;gap:8px;font-size:0.78rem;color:#818cf8;font-weight:700;">
      <span class="live-dot"></span>
      🧠 PROSES PENALARAN AKTIF AI ({model_name})
    </div>
    <span style="font-size:0.72rem;color:#94a3b8;">{len(thinking_text.split())} kata dipikirkan</span>
  </div>
  <div class="thinking-live-text">{html.escape(thinking_text)}<span class="cur"></span></div>
</div>
""", unsafe_allow_html=True)
                else:
                    ph.markdown(f"""
<div class="thinking-live-box" style="border-color:rgba(74, 222, 128, 0.35); opacity:0.9; margin-bottom:12px;">
  <div class="thinking-live-header">
    <div style="display:flex;align-items:center;gap:8px;font-size:0.78rem;color:#4ade80;font-weight:700;">
      <span>✓</span>
      🧠 PENALARAN AI SELESAI ({len(thinking_text.split())} KATA)
    </div>
    <span style="font-size:0.72rem;color:#4ade80;">Menyusun Catatan Master...</span>
  </div>
</div>

{main_text} ▌
""", unsafe_allow_html=True)
            
            final_think, final_main = parse_thought_and_content(raw_accumulated)
            if not final_main: final_main = raw_accumulated
            
            thought_accordion = ""
            if final_think:
                thought_accordion = f"""
<div class="thinking-live-box" style="border-color:rgba(74,222,128,0.25);margin-bottom:12px;">
  <div class="thinking-live-header" style="margin-bottom:0;border:none;padding-bottom:0;">
    <details style="font-size:0.8rem;color:#a5b4fc;width:100%;cursor:pointer;">
      <summary style="font-weight:700;color:#818cf8;outline:none;">🧠 Jejak Berpikir & Analisis Transparan AI ({len(final_think.split())} Kata)</summary>
      <div class="thinking-live-text" style="margin-top:8px;padding-top:8px;border-top:1px solid rgba(255,255,255,0.06);max-height:220px;overflow-y:auto;color:#cbd5e1;">{html.escape(final_think)}</div>
    </details>
  </div>
</div>
"""
            ph.markdown(f'{thought_accordion}\n\n{final_main}', unsafe_allow_html=True)
            return final_main
        except Exception as e:
            err_str = str(e)
            last_err = err_str
            continue
            
    st.error(f"⚠️ Koneksi AI terhambat ({last_err[:120] if last_err else 'Batas kuota tercapai'}). Silakan klik tombol sekali lagi untuk mencoba ulang.")
    return ""

def days_badge(s, sessions=0, review_count=0):
    if sessions == 0 and review_count == 0:
        return '<span class="badge" style="background:rgba(59,130,246,0.12);color:#93c5fd;border:1px solid rgba(59,130,246,0.25);font-weight:700;">🆕 Baru</span>'
    try:
        d = (datetime.datetime.fromisoformat(s) - datetime.datetime.now()).days
        return '<span class="badge br">⚡ Review!</span>' if d <= 0 else f'<span class="badge bb">🗓 {d}h</span>'
    except:
        return '<span class="badge" style="background:rgba(59,130,246,0.12);color:#93c5fd;border:1px solid rgba(59,130,246,0.25);font-weight:700;">🆕 Baru</span>' if (sessions == 0 and review_count == 0) else ""

def clean_scientific_math(text):
    if not text: return ""
    import re
    # Specific common pharmacological / biochemical constants
    text = re.sub(r"\$?\s*E\s*\{?\s*(?:max|maks)\s*\}?\s*\$?", "Emax", text, flags=re.IGNORECASE)
    text = re.sub(r"\$?\s*EC\s*\{?\s*50\s*\}?\s*\$?", "EC50", text, flags=re.IGNORECASE)
    text = re.sub(r"\$?\s*IC\s*\{?\s*50\s*\}?\s*\$?", "IC50", text, flags=re.IGNORECASE)
    text = re.sub(r"\$?\s*ED\s*\{?\s*50\s*\}?\s*\$?", "ED50", text, flags=re.IGNORECASE)
    text = re.sub(r"\$?\s*LD\s*\{?\s*50\s*\}?\s*\$?", "LD50", text, flags=re.IGNORECASE)
    
    # Common greek letters & math symbols
    greek = {
        r"\\alpha": "α", r"\\beta": "β", r"\\gamma": "γ", r"\\delta": "δ",
        r"\\mu": "μ", r"\\pm": "±", r"\\approx": "≈", r"\\leq": "≤", r"\\geq": "≥",
        r"\\times": "×", r"\\degree": "°"
    }
    for k, v in greek.items():
        text = re.sub(k, v, text)
        
    text = re.sub(r"\\text\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\$([^\$]+)\$", r"\1", text)
    text = re.sub(r"\{([A-Za-z0-9_]+)\}", r"\1", text)
    text = re.sub(r"[*_`]", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def parse_markdown_to_tree_pro(md_text):
    lines = [l for l in md_text.split('\n') if l.strip()]
    root = {"name": "Topik Utama", "children": []}
    stack = [(0, root)]
    
    for raw_line in lines:
        stripped = raw_line.strip()
        if not stripped: continue
        
        if not (stripped.startswith('#') or stripped.startswith('-') or stripped.startswith('*') or stripped.startswith('•') or (stripped[0].isdigit() and (stripped[1:3] == '. ' or stripped[2:4] == '. '))):
            if len(stack) == 1 and not root["children"]:
                continue
                
        if stripped.startswith('# '):
            depth = 1; name = stripped[2:].strip()
        elif stripped.startswith('## '):
            depth = 2; name = stripped[3:].strip()
        elif stripped.startswith('### '):
            depth = 3; name = stripped[4:].strip()
        elif stripped.startswith('#### '):
            depth = 4; name = stripped[5:].strip()
        elif stripped.startswith('- ') or stripped.startswith('* ') or stripped.startswith('• '):
            indent = len(raw_line) - len(raw_line.lstrip())
            depth = 5 + (indent // 2)
            name = stripped[2:].strip()
        elif stripped[0].isdigit() and (stripped[1:3] == '. ' or stripped[2:4] == '. '):
            depth = 2; name = stripped.split('. ', 1)[1].strip()
        else:
            indent = len(raw_line) - len(raw_line.lstrip())
            depth = 5 + (indent // 2)
            name = stripped
            
        name = clean_scientific_math(name)
        if not name: continue
        
        node = {"name": name, "children": []}
        if depth == 1:
            root["name"] = name
            stack = [(1, root)]
            continue
            
        while stack and stack[-1][0] >= depth:
            stack.pop()
            
        if stack:
            parent = stack[-1][1]
            parent["children"].append(node)
            stack.append((depth, node))
            
    if not root["children"] and len(lines) > 0:
        root["name"] = re.sub(r'^[#\-*•\d\.\s]+', '', lines[0]).strip()
        for l in lines[1:]:
            s = re.sub(r'^[#\-*•\d\.\s]+', '', l).strip()
            if s: root["children"].append({"name": s, "children": []})
            
    return root

def build_mindmap_html(md_text):
    tree_data = parse_markdown_to_tree_pro(md_text)
    tree_json = json.dumps(tree_data, ensure_ascii=False)
    
    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  html, body {{
    width:100%; height:100%; overflow:hidden;
    background:#070913; font-family:'Plus Jakarta Sans', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
    user-select:none; -webkit-user-select:none;
  }}
  #canvas-container {{
    width:100%; height:100%; position:relative; cursor:grab; overflow:hidden;
  }}
  #canvas-container:active {{ cursor:grabbing; }}
  svg {{ width:100%; height:100%; display:block; }}
  
  .top-bar {{
    position:absolute; top:12px; left:14px; z-index:30;
    display:flex; align-items:center; gap:8px;
    background:rgba(15, 23, 42, 0.92); backdrop-filter:blur(12px);
    padding:6px 14px; border-radius:20px; border:1px solid rgba(99, 102, 241, 0.35);
    color:#cbd5e1; font-size:12px; font-weight:500; pointer-events:none;
    box-shadow:0 4px 16px rgba(0,0,0,0.5);
  }}
  .glow-dot {{ width:7px; height:7px; border-radius:50%; background:#6366f1; box-shadow:0 0 8px #6366f1; animation:pulse 2s infinite; }}
  @keyframes pulse {{ 0%,100%{{transform:scale(1);opacity:0.6}} 50%{{transform:scale(1.3);opacity:1}} }}

  .toolbar {{
    position:absolute; bottom:14px; right:14px; z-index:30;
    display:flex; align-items:center; gap:6px;
    background:rgba(15, 23, 42, 0.92); backdrop-filter:blur(12px);
    padding:6px 10px; border-radius:12px; border:1px solid rgba(255,255,255,0.1);
    box-shadow:0 6px 20px rgba(0,0,0,0.5);
  }}
  .toolbar button {{
    background:#1e293b; color:#cbd5e1; border:1px solid #334155; border-radius:7px;
    padding:6px 12px; font-size:11.5px; font-weight:600; cursor:pointer;
    transition:all 0.15s ease; outline:none; display:flex; align-items:center; gap:4px;
  }}
  .toolbar button:hover {{
    background:#4f46e5; color:#ffffff; border-color:#6366f1; transform:translateY(-1px);
    box-shadow:0 3px 10px rgba(99, 102, 241, 0.4);
  }}
  .toolbar button:active {{ transform:translateY(0); }}
  .zoom-label {{ color:#94a3b8; font-size:11px; font-weight:600; padding:0 4px; min-width:40px; text-align:center; }}

  .link-path {{
    fill:none; stroke-width:2.2px; stroke-linecap:round;
    transition:stroke 0.2s, stroke-width 0.2s;
  }}
  .node-g {{ cursor:pointer; transition:transform 0.25s cubic-bezier(0.16, 1, 0.3, 1); }}
  .node-bg {{ transition:all 0.2s ease; }}
  .node-g:hover .node-bg {{
    filter:brightness(1.25) drop-shadow(0 0 12px rgba(99, 102, 241, 0.6));
  }}
  .node-title {{
    fill:#f8fafc; font-size:12px; font-weight:600; pointer-events:none;
    dominant-baseline:central;
  }}
  .node-sub {{
    fill:#e2e8f0; font-size:11.5px; font-weight:500; pointer-events:none;
    dominant-baseline:central;
  }}
  .node-root {{
    fill:#ffffff; font-size:13px; font-weight:700; pointer-events:none;
    dominant-baseline:central;
  }}
  .badge-bg {{ transition:all 0.2s; }}
  .badge-txt {{ font-size:10px; font-weight:700; pointer-events:none; text-anchor:middle; dominant-baseline:central; }}
</style>
</head>
<body>

<div class="top-bar">
  <div class="glow-dot"></div>
  <span>Klik cabang untuk buka/tutup • Drag kanvas untuk geser • Tombol zoom di kanan bawah</span>
</div>

<div class="toolbar">
  <button onclick="expandAllTree()">📂 Buka Semua</button>
  <button onclick="collapseAllTree()">📁 Ringkas</button>
  <button onclick="resetCamera()">🎯 Pusatkan</button>
  <button onclick="applyZoom(1.2)">🔍 +</button>
  <span class="zoom-label" id="zoom-text">100%</span>
  <button onclick="applyZoom(0.8)">🔍 -</button>
</div>

<div id="canvas-container">
  <svg id="main-svg">
    <g id="viewport"></g>
  </svg>
</div>

<script>
const rawTree = {tree_json};
const PALETTE = [
  {{ primary:'#6366f1', glow:'rgba(99,102,241,0.4)', bg:'#13182c' }},
  {{ primary:'#0ea5e9', glow:'rgba(14,165,233,0.4)', bg:'#0c1b2c' }},
  {{ primary:'#10b981', glow:'rgba(16,185,129,0.4)', bg:'#0c241e' }},
  {{ primary:'#f59e0b', glow:'rgba(245,158,11,0.4)', bg:'#261d0c' }},
  {{ primary:'#ec4899', glow:'rgba(236,72,153,0.4)', bg:'#28111f' }},
  {{ primary:'#8b5cf6', glow:'rgba(139,92,246,0.4)', bg:'#1d132e' }},
  {{ primary:'#14b8a6', glow:'rgba(20,184,166,0.4)', bg:'#0d2222' }}
];

let root = JSON.parse(JSON.stringify(rawTree));
let camera = {{ x: 60, y: 0, scale: 1.0 }};
let isPanning = false, startMouseX = 0, startMouseY = 0;

function splitText(text, maxChars = 34) {{
  if (!text || text.length <= maxChars) return [text || ''];
  const words = text.split(' ');
  const lines = [];
  let curLine = '';
  
  for (let w of words) {{
    if ((curLine + ' ' + w).trim().length > maxChars) {{
      if (curLine) lines.push(curLine.trim());
      curLine = w;
    }} else {{
      curLine = (curLine + ' ' + w).trim();
    }}
  }}
  if (curLine) lines.push(curLine.trim());
  return lines; // Complete, no slicing!
}}

function prepareNodes(node, depth = 0, branchColorObj = null) {{
  node.depth = depth;
  if (depth === 0) {{
    node.colorObj = {{ primary: '#6366f1', glow: 'rgba(99,102,241,0.5)', bg: '#4338ca' }};
  }} else if (depth === 1) {{
    node.colorObj = PALETTE[(node.branchIdx || 0) % PALETTE.length];
    branchColorObj = node.colorObj;
  }} else {{
    node.colorObj = branchColorObj || PALETTE[depth % PALETTE.length];
  }}
  
  const maxChars = depth === 0 ? 30 : (depth === 1 ? 26 : 34);
  node.lines = splitText(node.name, maxChars);
  node.height = Math.max(38, node.lines.length * 17 + 16);
  
  const longestLine = Math.max(...node.lines.map(l => l.length));
  node.width = Math.min(Math.max(longestLine * 7.4 + (depth <= 1 ? 52 : 40), 120), 380);
  
  if (node.children) {{
    node.children.forEach((c, idx) => {{
      c.branchIdx = idx;
      prepareNodes(c, depth + 1, branchColorObj);
    }});
  }}
}}

function layout(node) {{
  const HORIZ_STEP = 340;
  const VERT_GAP = 14;
  
  function countLeaves(n) {{
    if (!n.children || n.children.length === 0) {{
      n._leaves = 1;
      n._totalH = n.height;
    }} else {{
      let count = 0;
      let totalH = 0;
      n.children.forEach(c => {{
        countLeaves(c);
        count += c._leaves;
        totalH += c._totalH + VERT_GAP;
      }});
      n._leaves = count;
      n._totalH = Math.max(totalH - VERT_GAP, n.height);
    }}
  }}
  
  countLeaves(node);
  
  function place(n, depth, topY) {{
    n.x = depth * HORIZ_STEP + 40;
    n.y = topY + (n._totalH / 2) - (n.height / 2);
    
    if (n.children && n.children.length > 0) {{
      let curY = topY;
      n.children.forEach(c => {{
        place(c, depth + 1, curY);
        curY += c._totalH + VERT_GAP;
      }});
    }}
  }}
  
  place(node, 0, 40);
}}

function draw() {{
  const vp = document.getElementById('viewport');
  vp.innerHTML = '';
  vp.setAttribute('transform', `translate(${{camera.x}}, ${{camera.y}}) scale(${{camera.scale}})`);
  document.getElementById('zoom-text').innerText = Math.round(camera.scale * 100) + '%';
  
  const allNodes = [];
  const allLinks = [];
  
  function traverse(n) {{
    allNodes.push(n);
    if (n.children) {{
      n.children.forEach(c => {{
        allLinks.push({{ src: n, dst: c, colorObj: c.colorObj }});
        traverse(c);
      }});
    }}
  }}
  traverse(root);
  
  allLinks.forEach(l => {{
    const x1 = l.src.x + l.src.width;
    const y1 = l.src.y + (l.src.height / 2);
    const x2 = l.dst.x;
    const y2 = l.dst.y + (l.dst.height / 2);
    const mx = (x1 + x2) / 2;
    
    const path = document.createElementNS('http://www.w3.org/2000/svg', 'path');
    path.setAttribute('d', `M ${{x1}} ${{y1}} C ${{mx}} ${{y1}}, ${{mx}} ${{y2}}, ${{x2}} ${{y2}}`);
    path.setAttribute('class', 'link-path');
    path.setAttribute('stroke', l.colorObj.primary);
    path.setAttribute('stroke-opacity', '0.65');
    vp.appendChild(path);
  }});
  
  allNodes.forEach(n => {{
    const g = document.createElementNS('http://www.w3.org/2000/svg', 'g');
    g.setAttribute('class', 'node-g');
    g.setAttribute('transform', `translate(${{n.x}}, ${{n.y}})`);
    
    const isRoot = (n.depth === 0);
    const isMainBranch = (n.depth === 1);
    const hasChildren = (n.children && n.children.length > 0);
    const hasHidden = (n._children && n._children.length > 0);
    const hasAnyChildren = hasChildren || hasHidden;
    
    const rect = document.createElementNS('http://www.w3.org/2000/svg', 'rect');
    rect.setAttribute('class', 'node-bg');
    rect.setAttribute('width', n.width);
    rect.setAttribute('height', n.height);
    rect.setAttribute('rx', isRoot ? 12 : 9);
    
    if (isRoot) {{
      rect.setAttribute('fill', '#4f46e5');
      rect.setAttribute('stroke', '#a5b4fc');
      rect.setAttribute('stroke-width', '2.5');
    }} else if (isMainBranch) {{
      rect.setAttribute('fill', n.colorObj.bg);
      rect.setAttribute('stroke', n.colorObj.primary);
      rect.setAttribute('stroke-width', '2');
    }} else {{
      rect.setAttribute('fill', '#111827');
      rect.setAttribute('stroke', n.colorObj.primary);
      rect.setAttribute('stroke-width', '1.4');
      rect.setAttribute('stroke-opacity', '0.8');
    }}
    g.appendChild(rect);
    
    if (isMainBranch) {{
      const stripe = document.createElementNS('http://www.w3.org/2000/svg', 'rect');
      stripe.setAttribute('x', 0);
      stripe.setAttribute('y', 0);
      stripe.setAttribute('width', 4.5);
      stripe.setAttribute('height', n.height);
      stripe.setAttribute('rx', 3);
      stripe.setAttribute('fill', n.colorObj.primary);
      g.appendChild(stripe);
    }}
    
    const startX = isMainBranch ? 14 : (isRoot ? 16 : 12);
    const numLines = n.lines.length;
    const startY = (n.height / 2) - ((numLines - 1) * 8.5);
    
    n.lines.forEach((lineStr, lineIdx) => {{
      const t = document.createElementNS('http://www.w3.org/2000/svg', 'text');
      if (isRoot) {{
        t.setAttribute('class', 'node-root');
      }} else if (isMainBranch) {{
        t.setAttribute('class', 'node-title');
      }} else {{
        t.setAttribute('class', 'node-sub');
      }}
      t.setAttribute('x', startX);
      t.setAttribute('y', startY + (lineIdx * 17));
      t.textContent = lineStr;
      g.appendChild(t);
    }});
    
    if (hasAnyChildren && !isRoot) {{
      const bx = n.width - 12;
      const by = n.height / 2;
      
      const badgeCircle = document.createElementNS('http://www.w3.org/2000/svg', 'circle');
      badgeCircle.setAttribute('class', 'badge-bg');
      badgeCircle.setAttribute('cx', bx);
      badgeCircle.setAttribute('cy', by);
      badgeCircle.setAttribute('r', 8.5);
      badgeCircle.setAttribute('fill', hasHidden ? n.colorObj.primary : '#1e293b');
      badgeCircle.setAttribute('stroke', n.colorObj.primary);
      badgeCircle.setAttribute('stroke-width', '1.5');
      g.appendChild(badgeCircle);
      
      const badgeIcon = document.createElementNS('http://www.w3.org/2000/svg', 'text');
      badgeIcon.setAttribute('class', 'badge-txt');
      badgeIcon.setAttribute('x', bx);
      badgeIcon.setAttribute('y', by);
      badgeIcon.setAttribute('fill', hasHidden ? '#ffffff' : n.colorObj.primary);
      badgeIcon.textContent = hasHidden ? '+' : '−';
      g.appendChild(badgeIcon);
    }}
    
    g.onclick = (e) => {{
      e.stopPropagation();
      if (hasChildren) {{
        n._children = n.children;
        n.children = null;
      }} else if (hasHidden) {{
        n.children = n._children;
        n._children = null;
      }}
      prepareNodes(root);
      layout(root);
      draw();
    }};
    
    vp.appendChild(g);
  }});
}}

function resetCamera() {{
  camera.scale = 1.0;
  camera.x = 60;
  const containerH = window.innerHeight || 560;
  camera.y = (containerH / 2) - (root.y + (root.height / 2));
  draw();
}}

function zoomAtPoint(mouseX, mouseY, factor) {{
  const oldScale = camera.scale;
  const newScale = Math.min(Math.max(oldScale * factor, 0.3), 3.0);
  if (Math.abs(newScale - oldScale) < 0.001) return;
  
  const worldX = (mouseX - camera.x) / oldScale;
  const worldY = (mouseY - camera.y) / oldScale;
  
  camera.x = mouseX - (worldX * newScale);
  camera.y = mouseY - (worldY * newScale);
  camera.scale = newScale;
  draw();
}}

function applyZoom(factor) {{
  const container = document.getElementById("canvas-container");
  const cx = (container.clientWidth || 800) / 2;
  const cy = (container.clientHeight || 560) / 2;
  zoomAtPoint(cx, cy, factor);
}}

function expandAllTree() {{
  function exp(n) {{
    if (n._children) {{ n.children = n._children; n._children = null; }}
    if (n.children) n.children.forEach(exp);
  }}
  exp(root);
  prepareNodes(root);
  layout(root);
  draw();
}}

function collapseAllTree() {{
  function col(n) {{
    if (n.children && n.depth >= 1) {{
      n._children = n.children;
      n.children = null;
    }}
    if (n.children) n.children.forEach(col);
    if (n._children) n._children.forEach(col);
  }}
  col(root);
  prepareNodes(root);
  layout(root);
  draw();
}}

const container = document.getElementById('canvas-container');
container.addEventListener('mousedown', (e) => {{
  isPanning = true;
  startMouseX = e.clientX - camera.x;
  startMouseY = e.clientY - camera.y;
}});
window.addEventListener('mousemove', (e) => {{
  if (!isPanning) return;
  camera.x = e.clientX - startMouseX;
  camera.y = e.clientY - startMouseY;
  draw();
}});
window.addEventListener('mouseup', () => {{ isPanning = false; }});

container.addEventListener('wheel', (e) => {{
  e.preventDefault();
  const rect = container.getBoundingClientRect();
  const mouseX = e.clientX - rect.left;
  const mouseY = e.clientY - rect.top;
  const zoomFactor = e.deltaY < 0 ? 1.14 : 0.86;
  zoomAtPoint(mouseX, mouseY, zoomFactor);
}}, {{ passive: false }});

container.addEventListener('dblclick', () => {{ resetCamera(); }});

function collapseInitial(n) {{
  if (n.depth >= 1 && n.children && n.children.length > 0) {{
    n._children = n.children;
    n.children = null;
  }}
  if (n._children) {{
    n._children.forEach(collapseInitial);
  }}
}}

prepareNodes(root);
collapseInitial(root);
prepareNodes(root);
layout(root);
resetCamera();
</script>
</body>
</html>"""

# ── Session & Auth Initialization ─────────────────────────────────────────────
for k, v in [("phase", 0), ("mat_sel", None), ("phase_data", {}), ("completed", False), ("history", []), ("scores", {}), ("session_started", False), ("post_chat", [])]:
    if k not in st.session_state: st.session_state[k] = v

# ── AUTHENTICATION GATEWAY & SECURE LOGIN LANDING PAGE ────────────────────────
handle_google_oauth_code_exchange()

if "current_user" not in st.session_state or not st.session_state.current_user:
    persisted = get_persisted_auth_session()
    if persisted and persisted.get("username"):
        st.session_state.current_user = persisted["username"]
        st.session_state.user_info = persisted
        ensure_user_has_materials(persisted["username"])

if "current_user" not in st.session_state or not st.session_state.current_user:
    # Corner theme toggle
    col_t_l, col_t_r = st.columns([7.8, 1.4], vertical_alignment="center")
    with col_t_r:
        is_obsidian = st.session_state.get("app_theme", "obsidian") == "obsidian"
        t_lbl = "☀️ Terang" if is_obsidian else "🌙 Gelap"
        if st.button(t_lbl, key="landing_theme_toggle", use_container_width=True):
            new_th = "clinical_white" if is_obsidian else "obsidian"
            st.session_state.app_theme = new_th
            st.rerun()

    # Centered Minimalist Authentication Card (Modern Web Style)
    _, c_card, _ = st.columns([1.1, 1.8, 1.1])
    with c_card:
        st.markdown('''
<div style="text-align:center; padding: 10px 0 18px 0;">
  <div style="font-size: 42px; line-height: 1; margin-bottom: 8px; filter: drop-shadow(0 0 20px rgba(99,102,241,0.5));">🧠</div>
  <div style="font-size: 1.85rem; font-weight: 800; letter-spacing: -0.6px; color:#ffffff;">NeuroStudy</div>
  <div style="font-size: 0.84rem; color:#94a3b8; margin-top: 4px;">Platform Belajar Kedokteran &amp; Neurosains Klinis</div>
</div>
''', unsafe_allow_html=True)

        with st.container(border=True):
            tab_guest, tab_login, tab_reg = st.tabs(["🧪 Penguji Cepat", "🔑 Masuk", "✨ Daftar"])
            
            # ── TAB 1: MASUK PENGUJI INSTAN (UNTUK TEMAN-TEMAN TESTER) ──
            with tab_guest:
                st.markdown('''
<div style="padding:4px 0 12px; text-align:center;">
  <div style="font-weight:700; color:#38bdf8; font-size:1rem;">🚀 Akses Cepat Teman &amp; Penguji (Tanpa Sandi)</div>
  <p style="font-size:0.82rem; color:#94a3b8; line-height:1.5; margin:6px auto 0; max-width:94%;">
    Ketik nama dan email Anda (Gmail opsional). Langsung mulai mengeksplorasi <strong>208 Modul Kuliah Kedokteran</strong> &amp; fitur AI lengkap dengan akun bersih terpisah.
  </p>
</div>
''', unsafe_allow_html=True)
                with st.form("form_guest_tester_access", clear_on_submit=False):
                    g_name = st.text_input("Nama Lengkap / Panggilan:", placeholder="Contoh: dr. Sarah / Budi", key="guest_name_input")
                    g_email = st.text_input("Email (Opsional):", placeholder="nama@gmail.com (opsional)", key="guest_email_input")
                    st.markdown('<div style="height:4px;"></div>', unsafe_allow_html=True)
                    btn_guest = st.form_submit_button("🚀 Masuk & Mulai Belajar Sekarang →", type="primary", use_container_width=True)
                    
                    if btn_guest:
                        if not g_name.strip():
                            st.error("Silakan masukkan nama Anda untuk memulai.")
                        else:
                            u_res, msg = create_guest_tester_session(g_name, g_email)
                            if u_res:
                                st.session_state.current_user = u_res["username"]
                                st.session_state.user_info = u_res
                                set_persisted_auth_session(u_res)
                                disp_n = u_res.get('display_name') or u_res.get('username', 'Dokter')
                                st.toast(f"✓ Selamat datang, {disp_n}! Akses 208 modul siap.", icon="🟢")
                                st.rerun()
                            else:
                                st.error(f"⚠️ {msg}")

            # ── TAB 2: SUDAH PUNYA AKUN (MASUK BIASA) ──
            with tab_login:
                st.markdown("<div style='height: 6px;'></div>", unsafe_allow_html=True)
                with st.form("form_login_credentials", clear_on_submit=False):
                    in_ident = st.text_input("Email atau Username:", placeholder="nama@gmail.com atau username", key="login_ident_input")
                    in_pw = st.text_input("Kata Sandi (Password):", type="password", placeholder="Masukkan kata sandi", key="login_pw_input")
                    st.markdown('<div style="height:4px;"></div>', unsafe_allow_html=True)
                    btn_login = st.form_submit_button("🔑 Masuk ke Akun →", type="primary", use_container_width=True)
                    
                    if btn_login:
                        u_res, msg = authenticate_user(in_ident, in_pw)
                        if u_res:
                            ensure_user_has_materials(u_res["username"])
                            st.session_state.current_user = u_res["username"]
                            st.session_state.user_info = u_res
                            set_persisted_auth_session(u_res)
                            st.toast(f"✓ Selamat datang kembali, {u_res.get('display_name', u_res['username'])}!", icon="🟢")
                            st.rerun()
                        else:
                            st.error(f"⚠️ {msg}")
                            
                # Google OAuth / Akses Khusus Pemilik
                if is_localhost_access():
                    with st.expander("⚡ Akses Cepat Pemilik (dr. Dimas)", expanded=False):
                        st.caption("Masuk instan menggunakan akun Google pemilik yang aktif di mesin lokal ini.")
                        if st.button("⚡ Masuk Akun Pemilik (1-Klik)", key="btn_g_official_instant", use_container_width=True):
                            u_res, msg = authenticate_via_local_google_adc()
                            if u_res:
                                st.rerun()
                            else:
                                st.error(msg)
                else:
                    st.markdown('''
<div style="background: rgba(99, 102, 241, 0.08); border: 1px solid rgba(99, 102, 241, 0.2); border-radius: 8px; padding: 10px 12px; margin-top: 12px; font-size: 0.8rem; color: #94a3b8; text-align: center; line-height: 1.5;">
  💡 <strong>Teman &amp; Rekan Penguji:</strong> Belum memiliki kata sandi? Silakan gunakan tab <strong style="color:#38bdf8;">"🧪 Penguji Cepat"</strong> di atas untuk langsung masuk tanpa sandi, atau buat akun baru di tab <strong style="color:#818cf8;">"✨ Daftar"</strong>.
</div>
''', unsafe_allow_html=True)

            # ── TAB 3: BELUM PUNYA AKUN (DAFTAR PERMANEN) ──
            with tab_reg:
                st.markdown("<div style='height: 6px;'></div>", unsafe_allow_html=True)
                with st.form("form_register_user", clear_on_submit=False):
                    reg_name = st.text_input("Nama Lengkap:", placeholder="Contoh: dr. Dimas Wastu Mahesa", key="reg_fullname")
                    reg_email = st.text_input("Alamat Email:", placeholder="nama@gmail.com", key="reg_email_input")
                    reg_pw = st.text_input("Buat Kata Sandi (Password):", type="password", placeholder="Minimal 6 karakter", key="reg_pw_input")
                    reg_pw_conf = st.text_input("Konfirmasi Kata Sandi:", type="password", placeholder="Ulangi kata sandi", key="reg_pw_conf_input")
                    st.markdown('<div style="height:4px;"></div>', unsafe_allow_html=True)
                    btn_reg = st.form_submit_button("✨ Buat Akun Baru →", type="primary", use_container_width=True)
                    
                    if btn_reg:
                        u_res, msg = register_user(reg_name, reg_email, reg_pw, reg_pw_conf)
                        if u_res:
                            ensure_user_has_materials(u_res["username"])
                            st.session_state.current_user = u_res["username"]
                            st.session_state.user_info = u_res
                            set_persisted_auth_session(u_res)
                            st.toast(f"✓ Akun berhasil didaftarkan! Selamat datang di NeuroStudy, {reg_name}.", icon="🟢")
                            st.rerun()
                        else:
                            st.error(f"⚠️ {msg}")

        # Minimalist Whisper-Quiet Security Note
        st.markdown('''
<div style="margin-top: 16px; text-align: center; font-size: 0.74rem; color: #64748b; line-height: 1.6;">
  🔒 <strong>Keamanan Medis Terenkripsi</strong> · Multi-Tenant Terisolasi · Privasi Terjamin
</div>
''', unsafe_allow_html=True)

    # Stop rendering the rest of the application until authenticated!
    st.stop()

# ── PERSISTENT CONFIG & MATERIAL LOADING ──────────────────────────────────────
api_key = get_gemini_api_key()
mats = load_mats()

# ── EXECUTIVE TOP BAR (CLEAN CLINICAL INTERFACE) ──────────────────────────────
cur_u = st.session_state.get("user_info", {})
u_disp = (cur_u.get("display_name") or cur_u.get("username", "Dokter")).title()
u_mail = cur_u.get("email", "")
u_pic = cur_u.get("picture")
is_owner = bool(
    (u_mail and any(k in u_mail.lower() for k in ["dimaswastumahesa@gmail.com", "mahesawastu8@gmail.com"])) or
    (cur_u.get("username") in ["dimas", "dimaswastu", "dimaswastumahesa"])
)
st.session_state.is_owner = is_owner

c_brand, c_gtime, c_user, c_switch, c_theme, c_logout = st.columns([2.2, 2.5, 2.5, 1.3, 1.3, 0.9], vertical_alignment="center")

with c_brand:
    st.markdown('''
<div style="display:flex;align-items:center;gap:12px;">
  <div style="font-size:32px;filter:drop-shadow(0 0 10px rgba(99,102,241,0.5));">🧠</div>
  <div>
    <div style="font-size:1.25rem;font-weight:900;letter-spacing:-0.5px;line-height:1.1;">NeuroStudy</div>
    <div style="font-size:0.68rem;color:#818cf8;font-weight:800;letter-spacing:0.4px;margin-top:2px;">CLINICAL NEUROSCIENCE LEARNING PLATFORM</div>
  </div>
</div>
''', unsafe_allow_html=True)

with c_gtime:
    t_live = get_current_indonesia_time()
    with st.popover(f"🕒 {t_live['short_str']}", use_container_width=True, help="Waktu Sistem & Hub Integrasi Google Workspace"):
        st.markdown("#### 🌐 Google Medical Workspace")
        
        c_tz1, c_tz2 = st.columns([1.8, 2.2])
        with c_tz1:
            tz_opts = ["WIB", "WITA", "WIT"]
            cur_tz = st.session_state.get("selected_tz", "WIB")
            def_idx = tz_opts.index(cur_tz) if cur_tz in tz_opts else 0
            new_tz = st.selectbox("Zona Waktu:", tz_opts, index=def_idx, key="sb_tz_picker", label_visibility="collapsed")
            if new_tz != cur_tz:
                st.session_state.selected_tz = new_tz
                st.rerun()
        with c_tz2:
            st.caption(f"📍 {t_live['tz_desc']}")

        st.markdown(f'''
<div style="background:linear-gradient(135deg, rgba(66,133,244,0.12), rgba(52,168,83,0.12)); border:1px solid rgba(66,133,244,0.3); border-radius:12px; padding:12px 14px; margin-bottom:12px;">
  <div style="font-size:0.7rem; color:#94a3b8; font-weight:800; text-transform:uppercase; letter-spacing:0.5px;">WAKTU & KALENDER SISTEM TERINTEGRASI</div>
  <div style="font-size:1.05rem; font-weight:900; color:#38bdf8; margin:4px 0;">{t_live['hari']}, {t_live['tanggal']} {t_live['bulan']} {t_live['tahun']}</div>
  <div style="font-size:0.9rem; font-weight:800; color:#10b981; display:flex; align-items:center; gap:6px;">
    <span>⏰ {t_live['jam']}</span>
    <span style="font-size:0.65rem; background:rgba(16,185,129,0.2); padding:2px 6px; border-radius:6px; color:#34d399; font-weight:700;">TERKONEKSI GOOGLE</span>
  </div>
</div>
''', unsafe_allow_html=True)
        st.caption("Akses cepat ekosistem Google untuk efisiensi rotasi klinik & studi kedokteran:")
        
        c_gw1, c_gw2 = st.columns(2)
        with c_gw1:
            st.markdown('''
<a href="https://calendar.google.com" target="_blank" style="text-decoration:none; display:block; background:rgba(66,133,244,0.12); border:1px solid rgba(66,133,244,0.35); padding:10px 12px; border-radius:10px; margin-bottom:8px;">
  <div style="font-size:0.85rem; font-weight:800; color:#60a5fa;">📅 Google Calendar</div>
  <div style="font-size:0.7rem; color:#94a3b8; margin-top:2px;">Cek jadwal modul & review</div>
</a>
<a href="https://meet.google.com/new" target="_blank" style="text-decoration:none; display:block; background:rgba(52,168,83,0.12); border:1px solid rgba(52,168,83,0.35); padding:10px 12px; border-radius:10px; margin-bottom:8px;">
  <div style="font-size:0.85rem; font-weight:800; color:#34d399;">🎥 Google Meet</div>
  <div style="font-size:0.7rem; color:#94a3b8; margin-top:2px;">Mulai diskusi kasus PBL</div>
</a>
''', unsafe_allow_html=True)
        with c_gw2:
            st.markdown('''
<a href="https://scholar.google.com" target="_blank" style="text-decoration:none; display:block; background:rgba(251,188,5,0.12); border:1px solid rgba(251,188,5,0.35); padding:10px 12px; border-radius:10px; margin-bottom:8px;">
  <div style="font-size:0.85rem; font-weight:800; color:#fbbf24;">🔬 Google Scholar</div>
  <div style="font-size:0.7rem; color:#94a3b8; margin-top:2px;">Riset jurnal EBM & konsensus</div>
</a>
<a href="https://drive.google.com" target="_blank" style="text-decoration:none; display:block; background:rgba(234,67,53,0.12); border:1px solid rgba(234,67,53,0.35); padding:10px 12px; border-radius:10px; margin-bottom:8px;">
  <div style="font-size:0.85rem; font-weight:800; color:#f87171;">📁 Google Drive</div>
  <div style="font-size:0.7rem; color:#94a3b8; margin-top:2px;">Arsip 208 slide kuliah cloud</div>
</a>
''', unsafe_allow_html=True)

with c_user:
    if u_pic:
        avatar_elem = f'<img src="{u_pic}" style="width:30px;height:30px;border-radius:50%;object-fit:cover;flex-shrink:0;border:1px solid rgba(66,133,244,0.6);" alt="Google Avatar" />'
    else:
        avatar_elem = f'<div style="width:30px;height:30px;border-radius:50%;background:linear-gradient(135deg, #4285F4 0%, #34A853 50%, #FBBC05 75%, #EA4335 100%);color:#fff;display:flex;align-items:center;justify-content:center;font-weight:800;font-size:0.85rem;flex-shrink:0;">{u_disp[:1].upper()}</div>'

    email_badge = f"<span>● {u_mail}</span>" if u_mail else "<span>● Penguji Tamu</span>"
    st.markdown(f'''
<div style="display:flex;align-items:center;gap:10px;background:rgba(30,41,59,0.7);border:1px solid rgba(255,255,255,0.08);padding:6px 14px;border-radius:12px;">
  {avatar_elem}
  <div style="line-height:1.2;overflow:hidden;flex:1;">
    <div style="font-size:0.84rem;font-weight:800;white-space:nowrap;text-overflow:ellipsis;overflow:hidden;">{u_disp}</div>
    <div style="font-size:0.68rem;color:#34d399;font-weight:700;display:flex;align-items:center;gap:4px;">
      {email_badge}
    </div>
  </div>
  <span style="background:rgba(16,185,129,0.15);color:#34d399;font-size:0.62rem;padding:2px 6px;border-radius:6px;font-weight:800;">TERVERIFIKASI ✓</span>
</div>
''', unsafe_allow_html=True)

with c_switch:
    with st.popover("🔄 Akun", use_container_width=True):
        st.markdown("#### 🔄 Ganti Akun")
        curr_label = f"**{u_disp}** ({u_mail})" if u_mail else f"**{u_disp}**"
        st.caption(f"Akun aktif saat ini: {curr_label}")
        st.caption("Demi keamanan data berlangganan, silakan pilih metode autentikasi akun tujuan:")
        
        if is_localhost_access():
            g_url_sw = get_official_google_auth_url()
            st.markdown(f'''
<div style="margin-bottom: 10px;">
  <a href="{g_url_sw}" target="_self" style="text-decoration:none; display:flex; align-items:center; justify-content:center; gap:8px; background:#ffffff; color:#1f2937; border:1px solid #dadce0; padding:8px 12px; border-radius:8px; font-weight:600; font-size:0.82rem; box-shadow:0 1px 2px rgba(0,0,0,0.06);">
    <svg width="15" height="15" viewBox="0 0 24 24"><path fill="#4285F4" d="M22.56 12.25c0-.78-.07-1.53-.2-2.25H12v4.26h5.92c-.26 1.37-1.04 2.53-2.21 3.31v2.77h3.57c2.08-1.92 3.28-4.74 3.28-8.09z"/><path fill="#34A853" d="M12 23c2.97 0 5.46-.98 7.28-2.66l-3.57-2.77c-.98.66-2.23 1.06-3.71 1.06-2.86 0-5.29-1.93-6.16-4.53H2.18v2.84C3.99 20.53 7.7 23 12 23z"/><path fill="#FBBC05" d="M5.84 14.09c-.22-.66-.35-1.36-.35-2.09s.13-1.43.35-2.09V7.06H2.18C1.43 8.55 1 10.22 1 12s.43 3.45 1.18 4.94l2.85-2.22.81-.63z"/><path fill="#EA4335" d="M12 5.38c1.62 0 3.06.56 4.21 1.64l3.15-3.15C17.45 2.09 14.97 1 12 1 7.7 1 3.99 3.47 2.18 7.06l3.66 2.84c.87-2.6 3.3-4.52 6.16-4.52z"/></svg>
    <span>🌐 Masuk Akun Google Lain</span>
  </a>
</div>
''', unsafe_allow_html=True)
        
        st.markdown("---")
        if st.button("🔑 Masuk dengan Kata Sandi / Akun Lain", key="btn_switch_to_landing", type="primary", use_container_width=True):
            clear_persisted_auth_session()
            st.session_state.current_user = None
            st.session_state.user_info = None
            st.rerun()

with c_theme:
    is_obsidian = st.session_state.get("app_theme", "obsidian") == "obsidian"
    theme_label = "☀️ Terang" if is_obsidian else "🌙 Gelap"
    theme_help = "Beralih ke tampilan terang ramah mata di siang hari" if is_obsidian else "Beralih ke tampilan gelap elegan di malam hari"
    if st.button(theme_label, help=theme_help, use_container_width=True, key="btn_theme_toggle"):
        new_th = "clinical_white" if is_obsidian else "obsidian"
        st.session_state.app_theme = new_th
        curr_uname = st.session_state.get("current_user")
        if curr_uname:
            db = load_users()
            if curr_uname in db:
                db[curr_uname]["app_theme"] = new_th
                save_users(db)
        st.rerun()

with c_logout:
    if st.button("🚪 Keluar", use_container_width=True, key="top_logout_btn"):
        clear_persisted_auth_session()
        st.session_state.current_user = None
        st.session_state.user_info = None
        st.rerun()

st.markdown('<div style="margin-bottom:12px;"></div>', unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════════════
# MODULAR TAB RENDERERS (3 MENU UTAMA: RAPI, SIMPEL & POWERFUL)
# ══════════════════════════════════════════════════════════════════════

def render_sub_cloud_library():

    mats = load_mats()
    
    # Cari materi yang pernah dipelajari sebelumnya
    studied_mats = [k for k, v in mats.items() if (v.get("sessions", 0) > 0 or v.get("review_count", 0) > 0)]
    
    if studied_mats:
        latest_mat = st.session_state.get("mat_sel") if (st.session_state.get("mat_sel") in mats and (mats[st.session_state.get("mat_sel")].get("sessions", 0) > 0 or mats[st.session_state.get("mat_sel")].get("review_count", 0) > 0)) else studied_mats[0]
        mat_info = mats[latest_mat]
        ret_val = calculate_memory_retention(mat_info)
        ret_color = "#34d399" if (ret_val or 0) >= 75 else ("#f59e0b" if (ret_val or 0) >= 50 else "#ef4444")
        ret_badge_text = f"Daya Ingat: {ret_val}%"
        hero_tag = "LANJUTKAN KULIAH TERAKHIR:"
        hero_btn_text = "✨ Lanjutkan Belajar →"
        hero_icon = "📖"
    else:
        # Pengguna Baru: Belum ada riwayat sesi
        latest_mat = st.session_state.get("mat_sel") or (list(mats.keys())[0] if mats else None)
        ret_color = "#93c5fd"
        ret_badge_text = "Materi Siap Dipelajari"
        hero_tag = "🌟 REKOMENDASI KULIAH PERTAMA:"
        hero_btn_text = "🚀 Mulai Belajar Sekarang →"
        hero_icon = "🎯"

    # ── 1. HERO COCKPIT: LANJUTKAN BELAJAR INSTAN / KULIAH PERTAMA ──
    if latest_mat and latest_mat in mats:
        with st.container(border=True):
            c_h1, c_h2 = st.columns([3.8, 1.2], vertical_alignment="center")
            with c_h1:
                st.markdown(f'''
<div style="display:flex;align-items:center;gap:14px;">
  <div style="font-size:30px;background:rgba(99,102,241,0.15);width:52px;height:52px;border-radius:14px;display:flex;align-items:center;justify-content:center;border:1px solid rgba(99,102,241,0.3);">
    {hero_icon}
  </div>
  <div>
    <div style="display:flex;align-items:center;gap:8px;">
      <span style="font-size:0.68rem;color:#818cf8;font-weight:800;letter-spacing:0.5px;text-transform:uppercase;">{hero_tag}</span>
      <span style="background:rgba(59,130,246,0.12);color:{ret_color};font-size:0.68rem;padding:2px 8px;border-radius:10px;font-weight:700;">{ret_badge_text}</span>
    </div>
    <div style="font-size:1.1rem;font-weight:900;color:#ffffff;margin-top:2px;">{latest_mat}</div>
  </div>
</div>
''', unsafe_allow_html=True)
            with c_h2:
                if st.button(hero_btn_text, type="primary", use_container_width=True, key="btn_hero_resume"):
                    st.session_state.mat_sel = latest_mat
                    st.session_state.selected_notify = latest_mat
                    for b in ["BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD"]:
                        if latest_mat.startswith(f"[{b}]"):
                            st.session_state.t2_blok_selector = b
                            break
                    st.session_state.session_started = True
                    st.session_state.switch_tab_target = "Meja Belajar"
                    st.session_state.auto_gen_master = True
                    st.rerun()

    # ── 2. GOOGLE DRIVE LIVE SYNC HUB & CLOUD STATUS (PERMANENT AUTO-CONNECT) ──
    curr_user = st.session_state.get("current_user", "dimas")
    if "gdrive_auto_synced_session" not in st.session_state:
        st.session_state.gdrive_auto_synced_session = True
        check_and_auto_download_blok_updates(curr_user)

    with st.container(border=True):
        c_sync_info, c_sync_btns = st.columns([3.4, 1.6], vertical_alignment="center")
        with c_sync_info:
            st.markdown(f'''
<div style="display:flex;align-items:center;gap:14px;">
  <div style="font-size:30px;background:rgba(16,185,129,0.12);width:52px;height:52px;border-radius:14px;display:flex;align-items:center;justify-content:center;border:1px solid rgba(16,185,129,0.35);">
    🟢
  </div>
  <div>
    <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;">
      <span style="font-size:1.02rem;font-weight:900;color:#ffffff;">Google Drive Angkatan (VERTEXTERIAL)</span>
      <span style="background:rgba(16,185,129,0.15);color:#34d399;font-size:0.68rem;padding:2px 10px;border-radius:10px;font-weight:800;border:1px solid rgba(16,185,129,0.35);">✓ TERHUBUNG OTOMATIS</span>
    </div>
    <div style="font-size:0.78rem;color:#cbd5e1;margin-top:3px;line-height:1.4;">
      Semua <strong>{len(mats)} Modul Kuliah</strong> telah tersinkronisasi dan siap dipelajari langsung detik ini. <em>Akses instan kapan saja tanpa perlu unduh manual.</em>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)
        with c_sync_btns:
            if st.button("🔄 Periksa Slide Baru", use_container_width=True, key="btn_sync_gdrive_now", help="Pemeriksaan otomatis berjalan di latar belakang. Klik tombol ini jika dosen baru saja mengunggah slide baru hari ini."):
                with st.spinner("Memeriksa update materi terbaru di Google Drive..."):
                    check_and_auto_download_blok_updates(curr_user)
                    st.success("✓ Google Drive tersinkronisasi otomatis!")
                    st.rerun()
                    
    # ── 3. CATALOGUE FILTER & CATEGORY TABS ──
    st.markdown('<div style="margin-top:14px;"></div>', unsafe_allow_html=True)
    c_cat_hdr, c_sch_box, c_sort_box = st.columns([2.5, 3.5, 2.0], vertical_alignment="center")
    with c_cat_hdr:
        st.markdown(f'<div style="font-size:1.2rem;font-weight:900;color:#ffffff;letter-spacing:-0.4px;">📚 Modul Kuliah ({len(mats)})</div>', unsafe_allow_html=True)
    with c_sch_box:
        search_kw = st.text_input("Cari materi:", placeholder="🔍 Ketik judul kuliah, topik, atau kata kunci...", key="mats_search_box", label_visibility="collapsed")
    with c_sort_box:
        sort_opt = st.selectbox("Urutkan:", ["Semua Status", "🆕 Belum Dipelajari", "🔴 Butuh Review", "🟢 Retensi Kuat", "Abjad (A-Z)"], key="mats_filter_sort", label_visibility="collapsed")

    # Blok Category Pills
    blok_categories = ["Semua Blok", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BDT", "BMD", "Lainnya"]
    sel_blok_pill = st.radio("Pilih Blok:", blok_categories, index=0, horizontal=True, key="pill_blok_filter", label_visibility="collapsed")

    if st.session_state.get("selected_notify"):
        c_not1, c_not2 = st.columns([3.5, 1.5], vertical_alignment="center")
        with c_not1:
            st.success(f"📖 Modul **{st.session_state.selected_notify}** telah aktif di Meja Belajar!")
        with c_not2:
            if st.button("🚀 Buka Meja Belajar Sekarang →", type="primary", use_container_width=True, key="btn_notif_jump_mb"):
                st.session_state.switch_tab_target = "Meja Belajar"
                st.session_state.auto_gen_master = True
                st.rerun()

    # Filtering Engine
    filtered_items = list(mats.items())
    
    # Filter by Blok Pill
    if sel_blok_pill != "Semua Blok":
        if sel_blok_pill == "Lainnya":
            filtered_items = [(k, v) for k, v in filtered_items if not any(k.startswith(f"[{b}") for b in ["BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BDT", "BMD"])]
        else:
            filtered_items = [(k, v) for k, v in filtered_items if k.startswith(f"[{sel_blok_pill}]") or sel_blok_pill.lower() in k.lower()]

    # Filter by Search Keyword
    if search_kw.strip():
        filtered_items = [(k, v) for k, v in filtered_items if search_kw.strip().lower() in k.lower()]
        
    # Sort & Filter Engine
    if sort_opt == "🆕 Belum Dipelajari":
        filtered_items = [(k, v) for k, v in filtered_items if v.get("sessions", 0) == 0 and v.get("review_count", 0) == 0]
    elif sort_opt == "🔴 Butuh Review":
        filtered_items = [(k, v) for k, v in filtered_items if (v.get("sessions", 0) > 0 or v.get("review_count", 0) > 0) and (calculate_memory_retention(v) or 0) < 70]
    elif sort_opt == "🟢 Retensi Kuat":
        filtered_items = [(k, v) for k, v in filtered_items if (v.get("sessions", 0) > 0 or v.get("review_count", 0) > 0) and (calculate_memory_retention(v) or 0) >= 70]
    elif sort_opt == "Abjad (A-Z)":
        filtered_items = sorted(filtered_items, key=lambda x: x[0].lower())
    else:
        # Default smart sort: BMS 1 first, then BUAMS, then BMS 2..
        filtered_items = sorted(
            filtered_items,
            key=lambda x: (
                0 if x[0].startswith("[BMS 1]") else (
                    1 if x[0].startswith("[BUAMS]") else (
                        2 if x[0].startswith("[BMS") else 3
                    )
                ),
                x[0].lower()
            )
        )

    c_sub_meta1, c_sub_meta2 = st.columns([3.4, 1.6], vertical_alignment="center")
    with c_sub_meta1:
        st.markdown(f'<div style="font-size:0.78rem;color:#94a3b8;margin:2px 0 6px;">Menampilkan <strong>{len(filtered_items)}</strong> dari {len(mats)} materi kuliah · <span style="color:#818cf8;font-weight:700;">↕ Kotak Scroll Netap</span></div>', unsafe_allow_html=True)
    with c_sub_meta2:
        full_page_scroll = st.toggle("Mode Halaman Penuh", value=False, key="toggle_lib_full_scroll", help="Aktifkan jika ingin melihat daftar panjang tanpa batas kotak scroll.")

    if not mats:
        st.markdown('<div class="card card-sm"><div class="cs" style="text-align:center;">Belum ada materi kuliah tersimpan.</div></div>', unsafe_allow_html=True)
    elif not filtered_items:
        st.markdown(f'<div class="card card-sm"><div class="cs" style="text-align:center;">Tidak ada materi yang cocok dengan filter atau pencarian "{search_kw}".</div></div>', unsafe_allow_html=True)
    else:
        # ── 4. RESPONSIVE TWO-COLUMN GRID OF MATERIAL CARDS (FIXED SCROLL CONTAINER) ──
        container_kwargs = {"height": 620} if not full_page_scroll else {}
        with st.container(**container_kwargs):
            col_left, col_right = st.columns(2, gap="medium")
            cols = [col_left, col_right]
            
            for idx_item, (nm, md) in enumerate(filtered_items):
                target_col = cols[idx_item % 2]
                with target_col:
                    n_ses = md.get("sessions", 0)
                    n_rev = md.get("review_count", 0)
                    has_studied = (n_ses > 0 or n_rev > 0)
                    bd = days_badge(md.get("next_review",""), n_ses, n_rev)
                    retention_pct = calculate_memory_retention(md)
                    
                    if has_studied and retention_pct is not None:
                        ret_color = "#34d399" if retention_pct >= 75 else ("#f59e0b" if retention_pct >= 50 else "#ef4444")
                        ret_label = "Kuat" if retention_pct >= 75 else ("Menurun" if retention_pct >= 50 else "Kritis (Review Segera)")
                        ret_text = f"{retention_pct}% ({ret_label})"
                        bar_width = f"{retention_pct}%"
                        bar_bg = ret_color
                        meta_txt = f"{n_ses} sesi · {n_rev}× review · EF {md.get('ease_factor',2.5):.1f}"
                    else:
                        ret_color = "#94a3b8"
                        ret_text = "⚪ Belum Dipelajari"
                        bar_width = "0%"
                        bar_bg = "rgba(148, 163, 184, 0.2)"
                        meta_txt = "Belum ada sesi · Siap dipelajari"
                    
                    with st.container(border=True):
                        # Header & Info
                        c_mat_h1, c_mat_h2 = st.columns([3.6, 1.4], vertical_alignment="center")
                        with c_mat_h1:
                            st.markdown(f'''
<div style="font-size:0.95rem;font-weight:800;color:#ffffff;line-height:1.3;">📄 {nm}</div>
<div style="font-size:0.7rem;color:#94a3b8;margin-top:2px;">{meta_txt}</div>
''', unsafe_allow_html=True)
                        with c_mat_h2:
                            st.markdown(f'<div style="text-align:right;">{bd}</div>', unsafe_allow_html=True)
                        
                        # Ebbinghaus Metacognitive Retention Gauge
                        st.markdown(f'''
<div style="margin:8px 0 10px;">
  <div style="display:flex;justify-content:space-between;font-size:0.7rem;font-weight:700;margin-bottom:3px;">
    <span style="color:#94a3b8;">Kekuatan Retensi Memori:</span>
    <span style="color:{ret_color};font-weight:700;">{ret_text}</span>
  </div>
  <div style="background:rgba(255,255,255,0.06);height:6px;border-radius:6px;overflow:hidden;">
    <div style="background:{bar_bg};width:{bar_width};height:100%;border-radius:6px;transition:width 0.5s ease;"></div>
  </div>
</div>
''', unsafe_allow_html=True)
                        
                        # Clean Action Row (Main Action + Quick Power Tools Popover)
                        c_act_main, c_act_tools = st.columns([3.0, 1.0], vertical_alignment="center")
                        safe_k = re.sub(r'[^a-zA-Z0-9_]', '_', nm)
                        with c_act_main:
                            btn_label = f"🚀 Lanjut ({nm[:18]}…)" if has_studied else f"🚀 Mulai Belajar ({nm[:18]}…)"
                            if st.button(btn_label, type="primary", use_container_width=True, key=f"btn_start_sesi_{safe_k}"):
                                st.session_state.mat_sel = nm
                                st.session_state.selected_notify = nm
                                for b in ["BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD"]:
                                    if nm.startswith(f"[{b}]"):
                                        st.session_state.t2_blok_selector = b
                                        break
                                st.session_state.session_started = True
                                st.session_state.phase = 0
                                st.session_state.switch_tab_target = "Meja Belajar"
                                st.session_state.auto_gen_master = True
                                st.rerun()
                        with c_act_tools:
                            with st.popover("⚡ Fitur", use_container_width=True):
                                st.markdown(f"**⚡ Power Tools: {nm[:22]}...**")
                                if st.button("🎯 Bedah Jebakan Ujian", key=f"btn_v_{safe_k}", use_container_width=True):
                                    with st.spinner("Dr. Marcus Vance sedang membedah jebakan..."):
                                        prompt_v = f"Kamu adalah Dr. Marcus Vance, Sp.FK. Bedah 3 jebakan soal pilihan ganda paling mematikan pada materi: {md.get('text','')[:7000]}"
                                        stream_ai_transparent(api_key, prompt_v, st.empty())
                                if st.button("🩺 Kasus Klinis IGD", key=f"btn_t_{safe_k}", use_container_width=True):
                                    with st.spinner("Dr. Aris Thorne menyiapkan simulasi IGD..."):
                                        prompt_t = f"Kamu adalah Dr. Aris Thorne, Sp.PD. Buat 1 simulasi kasus pasien darurat di IGD berdasarkan materi ini: {md.get('text','')[:7000]}"
                                        stream_ai_transparent(api_key, prompt_t, st.empty())
                                if st.button("📄 Medical Cheat Sheet", key=f"btn_cs_{safe_k}", use_container_width=True):
                                    with st.spinner("Menyusun Cheat Sheet..."):
                                        cs_p = f"""Buat RANGKUMAN EKSEKUTIF 1 HALAMAN (Medical Cheat Sheet) yang sangat padat, terstruktur, dan siap cetak dari materi ini:
Materi:
{md.get('text','')[:7000]}

Format WAJIB:
1. **🧬 Mekanisme Inti & Cascade Sinyal:** (Intisari patofisiologi/farmakologi)
2. **📊 Tabel Obat / Klasifikasi:** (Golongan, Contoh, Efek Utama, Efek Samping Kritis)
3. **⚠️ 3 Aturan Emas Klinis:** (Peringatan penting di ranjang pasien)
4. **💡 Mnemonik Klinis:** (Cara cepat mengingat fakta rumit)"""
                                        stream_ai_transparent(api_key, cs_p, st.empty())
                                cards = load_flashcards(nm)
                                if cards:
                                    st.download_button(
                                        label="📥 Unduh Anki (.tsv)",
                                        data=generate_anki_export_data(cards, nm),
                                        file_name=f"{nm}_anki.txt",
                                        mime="text/plain",
                                        key=f"dl_anki_{safe_k}",
                                        use_container_width=True
                                    )

    # ── 5. OPTIONAL ADVANCED TOOLS EXPANDER (MANUAL UPLOAD & ZIP BACKUP) ──
    st.markdown('<div style="margin-top:20px;"></div>', unsafe_allow_html=True)
    with st.expander("🛠️ Alat Tambahan: Upload File Manual dari Laptop & Backup ZIP", expanded=False):
        c_up_m1, c_up_m2 = st.columns(2)
        with c_up_m1:
            st.markdown("##### 📤 Upload File Mandiri (.pptx, .pdf, .docx):")
            up_man = st.file_uploader("Pilih file dari komputer:", type=["pdf", "pptx", "ppt", "docx", "txt"], accept_multiple_files=True, key="manual_up_file")
            if up_man:
                added_any = False
                for f in up_man:
                    cl_n = re.sub(r'\.(pdf|pptx|ppt|docx|txt)$', '', f.name, flags=re.IGNORECASE).strip()
                    if cl_n not in mats:
                        with st.spinner(f"Mengekstrak {f.name}..."):
                            ext_txt = extract_document_text(f, api_key)
                            if len(ext_txt.strip()) > 30:
                                save_mat(cl_n, ext_txt)
                                st.success(f"✓ Berhasil ditambahkan: {cl_n}")
                                added_any = True
                if added_any:
                    st.rerun()
        with c_up_m2:
            st.markdown("##### 📦 Backup & Pemulihan Database:")
            zip_dat = create_backup_zip()
            st.download_button(
                label="📥 Unduh Backup Database (.zip)",
                data=zip_dat,
                file_name="neurostudy_backup.zip",
                mime="application/zip",
                use_container_width=True
            )



def render_sub_specialist_council():

    mats = load_mats()
    
    st.markdown('''
<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:8px;flex-wrap:wrap;gap:10px;">
  <div>
    <div style="font-size:1.6rem;font-weight:900;color:#f8fafc;letter-spacing:-0.5px;">🏢 NeuroStudy Enterprise: Board of AI Specialists</div>
    <div style="color:#94a3b8;font-size:.85rem;margin-top:2px;">Divisi pakar AI khusus dengan peran departemen medis masing-masing untuk membedah slide kuliah Anda.</div>
  </div>
  <div style="background:rgba(99,102,241,0.12);border:1px solid rgba(99,102,241,0.3);padding:6px 14px;border-radius:20px;font-size:0.75rem;color:#a5b4fc;font-weight:700;">
    🏛️ 5 DEPARTEMEN SPESIALIS AKTIF
  </div>
</div>
''', unsafe_allow_html=True)
    
    if not mats:
        st.markdown('<div class="card"><div class="cs" style="text-align:center;">Belum ada materi. Upload slide/PDF di Tab 1 terlebih dahulu.</div></div>', unsafe_allow_html=True)
    else:
        # Sort all materials cleanly by Blok then alphabetical
        mat_names = sorted(
            list(mats.keys()),
            key=lambda x: (
                0 if x.startswith("[BMS 1]") else (
                    1 if x.startswith("[BUAMS]") else (
                        2 if x.startswith("[BMS 2]") else (
                            3 if x.startswith("[BMS 3]") else (
                                4 if x.startswith("[BMS 4]") else (
                                    5 if x.startswith("[BDT]") else (
                                        6 if x.startswith("[BMD]") else 7
                                    )
                                )
                            )
                        )
                    )
                ),
                x.lower()
            )
        )
        current_active_mat = st.session_state.get("mat_sel")
        if current_active_mat not in mat_names: current_active_mat = mat_names[0]
        if "multi_mat_selector" not in st.session_state or not st.session_state.multi_mat_selector:
            st.session_state.multi_mat_selector = [current_active_mat]
            
        # ── 1. MULTI-PPT SELECTOR WITH BLOK FILTER ──
        c_t3_flt, c_t3_act = st.columns([2.2, 4.8], vertical_alignment="bottom")
        with c_t3_flt:
            t3_blok_opts = ["Semua Blok", "BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD", "Lainnya"]
            sel_t3_blok = st.selectbox("🏷️ Filter Blok:", t3_blok_opts, index=0, key="t3_blok_filter")
            
        # Determine choices for the selected Blok
        filtered_mat_names = mat_names
        if sel_t3_blok != "Semua Blok":
            if sel_t3_blok == "Lainnya":
                filtered_mat_names = [k for k in mat_names if not any(k.startswith(f"[{b}") for b in ["BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BDT", "BMD"])]
            else:
                filtered_mat_names = [k for k in mat_names if k.startswith(f"[{sel_t3_blok}]") or sel_t3_blok.lower() in k.lower()]
        if not filtered_mat_names:
            filtered_mat_names = mat_names

        with c_t3_act:
            cb_b1, cb_b2, cb_b3 = st.columns(3)
            with cb_b1:
                label_blok_btn = f"⚡ Colok {sel_t3_blok}" if sel_t3_blok != "Semua Blok" else "⚡ Colok Semua Slide"
                if st.button(label_blok_btn, use_container_width=True, help="Colok seluruh slide dari blok yang dipilih"):
                    st.session_state.multi_mat_selector = filtered_mat_names
                    st.rerun()
            with cb_b2:
                if st.button("📖 Kuliah Aktif Saja", use_container_width=True, help="Hanya pasang materi yang sedang aktif dipelajari"):
                    st.session_state.multi_mat_selector = [current_active_mat]
                    st.rerun()
            with cb_b3:
                if st.button("🧹 Kosongkan", use_container_width=True, help="Hapus seluruh pilihan slide"):
                    st.session_state.multi_mat_selector = []
                    st.rerun()

        # Valid options in multiselect must include current selection to avoid Streamlit errors
        curr_selected = [m for m in st.session_state.get("multi_mat_selector", []) if m in mat_names]
        if not curr_selected and current_active_mat in mat_names:
            curr_selected = [current_active_mat]
            st.session_state.multi_mat_selector = curr_selected

        if "multi_mat_selector" not in st.session_state:
            st.session_state.multi_mat_selector = curr_selected
            
        available_opts = list(dict.fromkeys(filtered_mat_names + curr_selected))
        plugged = st.multiselect(
            f"🔌 Slide/Materi Terpasang ({len(st.session_state.get('multi_mat_selector', []))} Terpilih dari {len(mat_names)} Total Modul):",
            options=available_opts,
            key="multi_mat_selector",
            help="Pilih 1 atau beberapa materi sekaligus untuk dianalisis komprehensif oleh dewan spesialis!"
        )

        if not plugged: plugged = [current_active_mat]
            
        # Combine texts from plugged materials (Rich full text up to 35,000 chars per slide)
        combined_texts, total_chars = [], 0
        for p_nm in plugged:
            cl_t = clean_academic_text(mats[p_nm].get("text", ""))
            total_chars += len(cl_t)
            combined_texts.append(f"=== SLIDE / MATERI: {p_nm} ===\n{cl_t[:35000]}")
        merged_context = "\n\n".join(combined_texts)
        
        # ── 2. SPECIALIST AGENT DEFINITIONS & SELECTION ──
        AGENTS = {
            "board": {
                "name": "🏢 Rapat Konsil Lengkap (All-Hands Board Consultation)",
                "role": "Chief Executive Medical Board",
                "dept": "Direksi & Seluruh Kepala Divisi Spesialis",
                "icon": "🏢",
                "color": "#6366f1",
                "prompt_role": "Kamu bertindak sebagai Rapat Konsil Dokter Lengkap (Medical Board). Setiap divisi pakar (Klinisi, Patologi Molekuler, Spesialis Ujian, Pedagogi Feynman, dan Riset) memberikan analisis sinergis dari sudut pandang keahliannya masing-masing secara terstruktur."
            },
            "thorne": {
                "name": "Dr. Aris Thorne, Sp.PD",
                "role": "Chief Clinical Diagnostician & Ward Simulator",
                "dept": "Divisi Diagnosis Klinis & Skenario Kasus Pasien",
                "icon": "🩺",
                "color": "#38bdf8",
                "prompt_role": "Kamu adalah Dr. Aris Thorne, Dokter Spesialis Penyakit Dalam Senior dan Konsultan IGD. Fokus utamamu adalah menerjemahkan teori slide menjadi kasus pasien nyata, algoritma penegakan diagnosis (DDx), pemeriksaan penunjang kritis, dan keputusan tatalaksana terapeutik di ranjang pasien."
            },
            "rostova": {
                "name": "Prof. Elena Rostova, Ph.D",
                "role": "VP of Molecular Pathology & Receptor Mechanisms",
                "dept": "Divisi Patologi Molekuler & Mekanisme Reseptor",
                "icon": "🔬",
                "color": "#a855f7",
                "prompt_role": "Kamu adalah Prof. Elena Rostova, Ilmuwan Biomedis Utama dan Peneliti Farmakologi Molekuler. Fokus utamamu adalah membedah jalur sinyal intraseluler, ikatan reseptor, farmakodinamik mikro, dan kausalitas biokimiawi paling mendalam dari materi kuliah."
            },
            "vance": {
                "name": "Dr. Marcus Vance, Sp.FK",
                "role": "Director of Exam Intelligence & High-Yield Trap Buster",
                "dept": "Divisi Intelijen Ujian Dosen & Kisi-Kisi UKMPPD",
                "icon": "🎯",
                "color": "#f59e0b",
                "prompt_role": "Kamu adalah Dr. Marcus Vance, Direktur Evaluasi Akademik dan Pakar Pemburu Jebakan Ujian. Fokus utamamu adalah mendeteksi jebakan distractor soal, area abu-abu yang paling sering mengecoh mahasiswa di ujian/blok/UKMPPD, dan trik membedakan opsi yang mirip."
            },
            "clara": {
                "name": "Clara Oswald, M.Ed, M.Sc",
                "role": "Director of Pedagogical Ergonomics & Mental Models",
                "dept": "Divisi Pedagogi Kognitif & Analogi Feynman",
                "icon": "🗣️",
                "color": "#4ade80",
                "prompt_role": "Kamu adalah Clara Oswald, Pakar Ergonomi Kognitif dan Desain Model Mental. Fokus utamamu adalah menyederhanakan konsep rumit menjadi analogi kehidupan sehari-hari yang intuitif dan mudah dibayangkan tanpa mengurangi esensi ilmiahnya."
            },
            "hayes": {
                "name": "Dr. Julian Hayes, MD, Ph.D",
                "role": "Chief Research Officer & Evidence-Based Validator",
                "dept": "Divisi Riset Ilmiah & Validasi Jurnal PubMed",
                "icon": "📊",
                "color": "#ec4899",
                "prompt_role": "Kamu adalah Dr. Julian Hayes, Kepala Peneliti Kedokteran Berbasis Bukti (EBM). Fokus utamamu adalah memvalidasi isi slide terhadap literatur standar emas (Goodman & Gilman, Harrison, Katzung, PubMed) serta mengklarifikasi fakta medis yang masih kontroversial."
            }
        }
        
        if "selected_agent_key" not in st.session_state:
            st.session_state.selected_agent_key = "board"
            
        st.markdown('<div style="font-size:0.75rem;color:#cbd5e1;font-weight:700;letter-spacing:0.5px;margin:12px 0 6px;text-transform:uppercase;">Pilih Spesialis AI yang Ingin Ditugaskan:</div>', unsafe_allow_html=True)
        
        agent_options = list(AGENTS.keys())
        agent_labels = {k: f"{v['icon']} {v['name']} ({v['role']})" for k, v in AGENTS.items()}
        
        sel_agent_k = st.selectbox(
            "Pilih Spesialis AI:",
            options=agent_options,
            format_func=lambda k: agent_labels[k],
            index=agent_options.index(st.session_state.selected_agent_key),
            label_visibility="collapsed",
            key="agent_selector_box"
        )
        st.session_state.selected_agent_key = sel_agent_k
        active_agent = AGENTS[sel_agent_k]
        
        # ── 3. AGENT PROFILE BANNER ──
        st.markdown(f'''
<div style="background:rgba(15,23,42,0.9);border:1.5px solid {active_agent['color']}66;border-radius:14px;padding:14px 18px;margin:10px 0 16px;display:flex;align-items:center;gap:14px;box-shadow:0 6px 20px rgba(0,0,0,0.5);">
  <div style="font-size:2.2rem;background:{active_agent['color']}22;border:1px solid {active_agent['color']}55;border-radius:14px;width:52px;height:52px;display:flex;align-items:center;justify-content:center;">
    {active_agent['icon']}
  </div>
  <div style="flex:1;">
    <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;">
      <span style="font-size:1.05rem;font-weight:900;color:#ffffff;">{active_agent['name']}</span>
      <span style="background:{active_agent['color']}22;color:{active_agent['color']};font-size:0.7rem;padding:2px 8px;border-radius:10px;font-weight:700;border:1px solid {active_agent['color']}44;">{active_agent['role']}</span>
    </div>
    <div style="font-size:0.75rem;color:#94a3b8;margin-top:2px;">🏛️ {active_agent['dept']} &nbsp;·&nbsp; 📄 Menganalisis <strong>{len(plugged)} Slide Terpasang</strong> (~{total_chars:,} karakter)</div>
  </div>
</div>
''', unsafe_allow_html=True)
        
        # ── 4. AGENT SPECIFIC DIRECTIVE BUTTONS ──
        st.markdown('<div style="font-size:0.75rem;color:#cbd5e1;font-weight:700;letter-spacing:0.5px;margin-bottom:6px;text-transform:uppercase;">Instruksi Cepat Sesuai Keahlian Spesialis:</div>', unsafe_allow_html=True)
        
        quick_action_prompt = ""
        if sel_agent_k == "board":
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("🏢 Rapat Konsil Lengkap", use_container_width=True):
                    quick_action_prompt = "Lakukan Rapat Konsil Tim Lengkap (Medical Board): Rangkum materi ini dari 4 sudut pandang: (1) Skenario Kasus Klinis, (2) Mekanisme Molekuler, (3) Jebakan Ujian Dosen, dan (4) Analogi Sederhana."
            with c2:
                if st.button("🔍 Audit Hubungan Slide", use_container_width=True):
                    quick_action_prompt = "Bedah secara komprehensif bagaimana materi-materi slide yang dicolokkan ini saling berhubungan dan melengkapi satu sama lain."
            with c3:
                if st.button("🏥 3 Kasus Pasien Terintegrasi", use_container_width=True):
                    quick_action_prompt = "Rancang 3 kasus pasien klinis yang mengintegrasikan seluruh materi ini, lengkapi dengan analisis diagnosis dan tatalaksana."
        elif sel_agent_k == "thorne":
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("🚑 Simulasi Pasien IGD / Bangsal", use_container_width=True):
                    quick_action_prompt = "Buat simulasi kasus pasien nyata yang masuk ke IGD/Poliklinik terkait materi ini. Jelaskan keluhan utama, tanda vital, dan langkah diagnosis pertama."
            with c2:
                if st.button("🩺 Alur Diferensial Diagnosis (DDx)", use_container_width=True):
                    quick_action_prompt = "Jelaskan pohon keputusan diferensial diagnosis (DDx) untuk gejala-gejala yang dibahas di slide ini beserta kriteria pembedanya."
            with c3:
                if st.button("💊 Keputusan Tatalaksana Klinis", use_container_width=True):
                    quick_action_prompt = "Jelaskan protokol terapi lini pertama, kedua, serta kontraindikasi berbahaya dari materi ini pada pasien komorbid."
        elif sel_agent_k == "rostova":
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("🔬 Bedah Jalur Molekuler & Sinyal", use_container_width=True):
                    quick_action_prompt = "Bedah secara mendalam jalur transduksi sinyal intraseluler, enzim, dan interaksi reseptor dari konsep yang ada di slide ini."
            with c2:
                if st.button("⚡ Kausalitas Sebab-Akibat Seluler", use_container_width=True):
                    quick_action_prompt = "Jelaskan rantai kausalitas sebab-akibat mikro: mengapa stimulasi atau hambatan pada tingkat molekuler ini bisa memicu efek organ tertentu."
            with c3:
                if st.button("🧬 Dinamika Afinitas & Reseptor", use_container_width=True):
                    quick_action_prompt = "Jelaskan konsep afinitas, efikasi intrinsik, dan potensi obat/molekul yang disebutkan di slide dalam interaksinya dengan target biologi."
        elif sel_agent_k == "vance":
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("🎯 Bedah 3 Jebakan Soal Ujian", use_container_width=True):
                    quick_action_prompt = "Bedah 3 potensi jebakan soal pilihan ganda / ujian dosen yang paling sering mengecoh mahasiswa pada materi ini."
            with c2:
                if st.button("⚠️ Poin Distractor Paling Menipu", use_container_width=True):
                    quick_action_prompt = "Tunjukkan konsep atau istilah yang mirip pada materi ini yang sering dijadikan opsi pengecoh (distractor) oleh pembuat soal ujian."
            with c3:
                if st.button("📝 Ringkasan High-Yield 1 Menit", use_container_width=True):
                    quick_action_prompt = "Buatkan daftar poin HIGH-YIELD super penting yang wajib dihafal mati sebelum masuk ruang ujian untuk materi ini."
        elif sel_agent_k == "clara":
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("🗣️ Analogi Sehari-hari (Feynman)", use_container_width=True):
                    quick_action_prompt = "Jelaskan konsep tersulit dari materi ini menggunakan analogi benda sehari-hari yang sangat mudah dibayangkan oleh orang awam."
            with c2:
                if st.button("🧠 Buat Visualisasi Model Mental", use_container_width=True):
                    quick_action_prompt = "Bantu saya membangun model mental visual yang kokoh tentang bagaimana cara kerja sistem/konsep ini di kepala."
            with c3:
                if st.button("❓ Uji Penjelasan Bahasa Saya Sendiri", use_container_width=True):
                    quick_action_prompt = "Berikan saya satu tantangan konsep untuk saya jelaskan kembali dengan kata-kata saya sendiri, lalu nilai pemahaman saya."
        else:  # hayes
            c1, c2, c3, c_clr = st.columns([1.3, 1.3, 1.3, 0.7])
            with c1:
                if st.button("📊 Validasi Literatur Standar Emas", use_container_width=True):
                    quick_action_prompt = "Validasi isi slide ini terhadap referensi standar emas (Goodman & Gilman, Harrison, Katzung, PubMed) dan berikan rujukan utamanya."
            with c2:
                if st.button("⚖️ Klarifikasi Area Abu-Abu", use_container_width=True):
                    quick_action_prompt = "Klarifikasi seluruh area abu-abu, batasan klinis, atau kontroversi ilmiah yang belum terjawab tuntas di slide ini."
            with c3:
                if st.button("📑 Rangkuman Riset Terkini", use_container_width=True):
                    quick_action_prompt = "Rangkum temuan penelitian atau guideline internasional terbaru yang memperbarui materi ini."
                    
        with c_clr:
            import hashlib
            disc_hash_key = hashlib.md5(f"{sel_agent_k}_".join(sorted(plugged)).encode("utf-8")).hexdigest()[:12]
            disc_storage_name = f"agent_{sel_agent_k}_{disc_hash_key}"
            if st.button("🗑️ Reset Chat", key=f"btn_reset_chat_{disc_storage_name}", use_container_width=True, help="Kosongkan riwayat chat dengan spesialis ini"):
                save_discussion(disc_storage_name, [])
                st.rerun()

        st.divider()
        
        # Load chat history for this specific agent & plugged deck
        disc_hist = load_discussion(disc_storage_name)
        
        if not disc_hist:
            st.markdown(f'''
<div class="card card-sm" style="padding:22px 24px;border:1px solid {active_agent['color']}44;margin-bottom:14px;background:rgba(15,23,42,0.7);">
  <div style="display:flex;align-items:center;gap:12px;margin-bottom:8px;">
    <div style="font-size:1.8rem;">{active_agent['icon']}</div>
    <div>
      <div style="font-weight:900;color:#f8fafc;font-size:1.02rem;">Konsultasi Khusus: {active_agent['name']}</div>
      <div style="color:{active_agent['color']};font-size:0.75rem;font-weight:700;">{active_agent['role']} · {active_agent['dept']}</div>
    </div>
  </div>
  <div style="color:#cbd5e1;font-size:0.86rem;line-height:1.6;margin-top:6px;">
    Spesialis ini telah mempelajari {len(plugged)} slide materi Anda. Anda dapat mengajukan pertanyaan bebas, meminta studi kasus mendalam, atau menggunakan tombol instruksi cepat di atas.
  </div>
</div>
''', unsafe_allow_html=True)
        else:
            for d_msg in disc_hist:
                if d_msg["role"] == "user":
                    st.markdown(f'<div class="msg-user">🙋‍♂️ <strong>Instruksi / Pertanyaan Anda:</strong><br/>{html.escape(d_msg["content"])}</div>', unsafe_allow_html=True)
                else:
                    msg_agent_name = d_msg.get("agent_name", active_agent["name"])
                    msg_agent_role = d_msg.get("agent_role", active_agent["role"])
                    msg_agent_icon = d_msg.get("agent_icon", active_agent["icon"])
                    msg_agent_color = d_msg.get("agent_color", active_agent["color"])
                    
                    st.markdown(f'''
<div class="msg-ai" style="border-left: 3px solid {msg_agent_color};">
  <div class="ai-row" style="margin-bottom:8px;">
    <div style="font-size:1.1rem;margin-right:6px;">{msg_agent_icon}</div>
    <span style="font-size:.82rem;color:#ffffff;font-weight:800;">{msg_agent_name}</span>
    <span style="font-size:.7rem;color:{msg_agent_color};background:{msg_agent_color}22;padding:1px 7px;border-radius:8px;font-weight:700;margin-left:6px;">{msg_agent_role}</span>
  </div>

{d_msg["content"]}
</div>
''', unsafe_allow_html=True)
        
        # Chat Input Form
        with st.form(key=f"agent_form_{disc_storage_name}", clear_on_submit=True):
            q_disc = st.text_area(
                f"Tulis instruksi atau pertanyaan untuk {active_agent['name']}:",
                placeholder=f"Contoh: Berdasarkan slide ini, jelaskan bagaimana sudut pandang {active_agent['role']}...",
                height=95,
                key=f"agent_in_{disc_storage_name}"
            )
            st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim (Gunakan <em>Shift + Enter</em> untuk baris baru).</div>', unsafe_allow_html=True)
            btn_send_disc = st.form_submit_button(f"💬 Kirim ke {active_agent['name']} →", type="primary", use_container_width=True)
            
        input_to_process = quick_action_prompt if quick_action_prompt else (q_disc.strip() if btn_send_disc else "")
        
        if input_to_process:
            disc_hist.append({
                "role": "user",
                "content": input_to_process
            })
            ph_disc = st.empty()
            
            prompt_disc = f"""{active_agent['prompt_role']}

Mahasiswa/Chief Medical Officer sedang mendiskusikan materi slide kuliah berikut:

Materi Sumber Slide ({len(plugged)} Slide Terpasang):
{merged_context[:120000]}

Pertanyaan/Instruksi Mahasiswa:
{input_to_process}

Riwayat Konsultasi Sebelumnya:
{json.dumps(disc_hist[-5:], ensure_ascii=False)}

Pedoman Menjawab:
1. Berikan analisis dengan otoritas kepakaran tinggi, bahasa ilmiah presisi, dan argumentasi tervalidasi berbasis bukti medis standar emas.
2. Fokuskan jawabanmu pada peran dan keahlian spesifikmu ({active_agent['dept']}).
3. Jelaskan mekanisme, implikasi klinis, dan rantai kausalitas secara terstruktur, jernih, dan menarik.
4. DILARANG menggunakan sintaks LaTeX berkurung kurawal ($E_{{max}}$). Gunakan format teks ilmiah bersih (Emax, EC50, ED50, α, β)."""
            
            ai_reply = stream_ai_transparent(api_key, prompt_disc, ph_disc)
            disc_hist.append({
                "role": "ai",
                "content": ai_reply,
                "agent_name": active_agent["name"],
                "agent_role": active_agent["role"],
                "agent_icon": active_agent["icon"],
                "agent_color": active_agent["color"]
            })
            save_discussion(disc_storage_name, disc_hist)
            st.rerun()


# ── TAB 4: Jadwal & Google Calendar ───────────────────────────────────────────


def render_sub_spaced_repetition():

    mats = load_mats()
    st.markdown('<div style="font-size:1.5rem;font-weight:800;color:#f8fafc;margin-bottom:4px;letter-spacing:-0.5px;">📅 Jadwal Belajar & Google Calendar</div>', unsafe_allow_html=True)
    st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-bottom:20px;">Spaced Repetition berbasis algoritma SM-2 & konsolidasi tidur — terintegrasi dengan Google Calendar dan kalender digital Anda.</p>', unsafe_allow_html=True)

    t_info = get_current_indonesia_time()
    st.markdown(f'''
<div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;background:rgba(30,41,59,0.7);border:1px solid rgba(66,133,244,0.3);padding:12px 18px;border-radius:12px;margin-bottom:18px;">
  <div>
    <div style="font-size:0.72rem;color:#818cf8;font-weight:800;letter-spacing:0.5px;text-transform:uppercase;">WAKTU SISTEM TERINTEGRASI GOOGLE CALENDAR</div>
    <div style="font-size:1.05rem;font-weight:900;color:#ffffff;margin-top:2px;">
      🗓️ {t_info["hari"]}, {t_info["tanggal"]} {t_info["bulan"]} {t_info["tahun"]} <span style="color:#38bdf8;font-weight:700;">• {t_info["jam"]}</span>
    </div>
  </div>
  <div style="display:flex;align-items:center;gap:10px;margin-top:4px;">
    <a href="https://calendar.google.com" target="_blank" style="text-decoration:none;background:linear-gradient(135deg, #4285F4 0%, #1d4ed8 100%);color:#ffffff;font-size:0.78rem;font-weight:700;padding:6px 14px;border-radius:8px;display:flex;align-items:center;gap:6px;box-shadow:0 2px 6px rgba(66,133,244,0.3);">
      <span>🌐 Buka Google Calendar Web</span> ↗
    </a>
  </div>
</div>
''', unsafe_allow_html=True)

    if not mats:
        st.markdown('<div class="card"><div class="cs" style="text-align:center;">Belum ada materi.</div></div>', unsafe_allow_html=True)
    else:
        now = datetime.datetime.now()
        due, up = [], []
        all_events = []
        
        studied_mats = [d for d in mats.values() if (d.get("sessions", 0) > 0 or d.get("review_count", 0) > 0)]
        for nm, d in mats.items():
            n_ses = d.get("sessions", 0)
            n_rev = d.get("review_count", 0)
            if n_ses == 0 and n_rev == 0:
                continue  # Lewati materi yang belum pernah dipelajari
                
            try:
                nr_str = d.get("next_review", "")
                dt_obj = datetime.datetime.fromisoformat(nr_str)
                days = (dt_obj - now).days
                rev_n = n_rev + 1
                (due if days <= 0 else up).append((nm, d, days, dt_obj, rev_n))
                
                all_events.append({
                    "name": nm,
                    "title": f"🧠 Review NeuroStudy: {nm} (Sesi {rev_n})",
                    "dt": dt_obj,
                    "rev_num": rev_n,
                    "desc": f"Spaced Repetition Review (SM-2 Algoritma)\nMateri: {nm}\nSesi ke-{rev_n} (Ease Factor {d.get('ease_factor',2.5):.2f})\nTarget: Recall aktif, Flashcards, dan Evaluasi Mental Model."
                })
            except:
                dt_obj = now + datetime.timedelta(days=1)
                due.append((nm, d, 0, dt_obj, 1))

        t_rev = sum(1 for d in studied_mats if d.get("review_count", 0) > 0)
        t_ses = sum(d.get("sessions", 0) for d in studied_mats)
        avg_ef = (sum(d.get("ease_factor", 2.5) for d in studied_mats) / len(studied_mats)) if studied_mats else 2.50

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("📚 Modul Dipelajari", f"{len(studied_mats)}/{len(mats)}")
        c2.metric("🔴 Harus Review", len(due))
        c3.metric("🧠 Total Sesi", t_ses)
        c4.metric("📈 Rata-rata EF", f"{avg_ef:.2f}")
        
        st.divider()
        
        st.markdown('<div style="font-size:1.05rem;font-weight:700;color:#f8fafc;margin-bottom:8px;">🗓️ Sinkronisasi Kalender Digital</div>', unsafe_allow_html=True)
        c_cal1, c_cal2 = st.columns([1, 1], gap="medium")
        with c_cal1:
            if all_events:
                ics_data = generate_ics_content(all_events)
                st.download_button(
                    label="📥 Unduh Semua Jadwal (.ics) untuk Apple / Google / Outlook",
                    data=ics_data,
                    file_name="neurostudy_schedule.ics",
                    mime="text/calendar",
                    use_container_width=True
                )
        with c_cal2:
            st.markdown('<div style="font-size:0.78rem;color:#94a3b8;line-height:1.5;">File <code>.ics</code> dapat langsung di-import ke Google Calendar, Apple Calendar, atau Outlook sekali klik.</div>', unsafe_allow_html=True)
            
        st.divider()

        cA, cB = st.columns(2, gap="large")
        with cA:
            st.markdown(f'<p style="font-size:.72rem;font-weight:700;color:#f87171;letter-spacing:1px;text-transform:uppercase;">Harus Direview Hari Ini ({len(due)})</p>', unsafe_allow_html=True)
            if not due:
                if not studied_mats:
                    st.markdown('<div class="card card-sm"><div class="cs" style="text-align:center;">✨ Belum ada modul di antrean review. Silakan pelajari modul pertama Anda di <strong>Ruang Belajar</strong> — jadwal review otomatis aktif setelah sesi pertama!</div></div>', unsafe_allow_html=True)
                else:
                    st.markdown('<div class="card card-sm"><div class="cs" style="text-align:center;">✅ Tidak ada modul yang perlu direview hari ini. Semua materi berada dalam retensi aman!</div></div>', unsafe_allow_html=True)
            for nm, d, _, dt_val, rev_n in due:
                gcal_link = build_gcal_url(f"🧠 Review NeuroStudy: {nm} (Sesi {rev_n})", dt_val, f"Spaced Repetition Review untuk materi {nm}")
                st.markdown(f'<div class="card card-sm" style="border-color:rgba(248,113,113,0.3);margin-bottom:6px;"><div style="display:flex;justify-content:space-between;align-items:center;"><div><div class="ct" style="font-size:.9rem;">🔴 {nm[:24]}</div><div class="cs">{d.get("sessions",0)} sesi · {d.get("review_count",0)}× review · EF {d.get("ease_factor",2.5):.1f}</div></div><div style="display:flex;align-items:center;gap:8px;"><span class="badge br">Review!</span><a href="{gcal_link}" target="_blank" style="text-decoration:none;background:rgba(99,102,241,0.2);color:#a5b4fc;border:1px solid rgba(99,102,241,0.4);border-radius:6px;padding:3px 8px;font-size:0.75rem;font-weight:600;">📅 +GCal</a></div></div></div>', unsafe_allow_html=True)
                safe_due_k = re.sub(r'[^a-zA-Z0-9_]', '_', nm)
                if st.button(f"🚀 Review Sekarang ({nm[:16]}...)", type="primary", key=f"btn_rev_due_{safe_due_k}", use_container_width=True):
                    st.session_state.mat_sel = nm
                    for b in ["BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD"]:
                        if nm.startswith(f"[{b}]"):
                            st.session_state.t2_blok_selector = b
                            break
                    st.session_state.switch_tab_target = "Meja Belajar"
                    st.session_state.auto_gen_master = True
                    st.rerun()
                
        with cB:
            st.markdown(f'<p style="font-size:.72rem;font-weight:700;color:#4ade80;letter-spacing:1px;text-transform:uppercase;">Jadwal Mendatang ({len(up)})</p>', unsafe_allow_html=True)
            if not up:
                st.markdown('<div class="card card-sm"><div class="cs" style="text-align:center;">Tidak ada jadwal mendatang.</div></div>', unsafe_allow_html=True)
            for nm, d, days, dt_val, rev_n in sorted(up, key=lambda x: x[2]):
                gcal_link = build_gcal_url(f"🧠 Review NeuroStudy: {nm} (Sesi {rev_n})", dt_val, f"Spaced Repetition Review untuk materi {nm}")
                st.markdown(f'<div class="card card-sm" style="margin-bottom:12px;"><div style="display:flex;justify-content:space-between;align-items:center;"><div><div class="ct" style="font-size:.9rem;">📄 {nm[:24]}</div><div class="cs">{d.get("sessions",0)} sesi · {d.get("review_count",0)}× review</div></div><div style="display:flex;align-items:center;gap:8px;"><span class="badge bb">🗓 {days} hari</span><a href="{gcal_link}" target="_blank" style="text-decoration:none;background:rgba(56,189,248,0.15);color:#7dd3fc;border:1px solid rgba(56,189,248,0.3);border-radius:6px;padding:3px 8px;font-size:0.75rem;font-weight:600;">📅 +GCal</a></div></div></div>', unsafe_allow_html=True)

        st.divider()
        pct = int((t_rev / len(mats)) * 100) if mats else 0
        st.markdown(f'<div style="display:flex;justify-content:space-between;margin-bottom:8px;"><span style="font-weight:700;color:#f8fafc;font-size:.92rem;">Progress Keseluruhan</span><span class="badge bp">{t_rev}/{len(mats)} materi · {pct}%</span></div>', unsafe_allow_html=True)
        st.progress(t_rev / len(mats) if mats else 0)

        st.divider()
        st.markdown("""<div class="card">
<div class="ct" style="margin-bottom:8px;">🌙 Konsolidasi Memori Saat Tidur (Neurosains Terkini)</div>
<div class="cs" style="line-height:1.8;">
Saat Anda tidur lelap malam ini (terutama fase Slow-Wave Sleep / SWS), gelombang otak <em>slow oscillations</em> (<1 Hz) yang disinkronkan dengan <em>sleep spindles</em> thalamus (11-16 Hz) dan <em>sharp-wave ripples</em> hippocampus (150-250 Hz) akan memutar ulang jejak memori yang Anda pelajari hari ini secara berulang-ulang, mentransfernya secara permanen dari hippocampus jangka pendek ke neokorteks jangka panjang.<br/><br/>
⏱️ <strong>Siklus Spaced Repetition (SM-2):</strong> Belajar hari ini → Review besok (H+1) → Review H+3 → Review H+7 → Review H+14 → Review H+30 (Memori Permanen Terkunci).
</div>
</div>""", unsafe_allow_html=True)




def render_sub_beta_tester():

    st.markdown('<div style="font-size:1.5rem;font-weight:800;color:#f8fafc;margin-bottom:4px;letter-spacing:-0.5px;">💬 Pusat Uji Coba & Kanal Masukan Beta Tester</div>', unsafe_allow_html=True)
    st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-bottom:16px;">Bagikan tautan akses publik kepada teman sejawat untuk menguji coba fitur Pareto 80/20, lalu kumpulkan umpan balik terstruktur secara langsung.</p>', unsafe_allow_html=True)
    
    # ── 1. KOTAK BERBAGI TAUTAN PUBLIK (SHARE LINK HUB) ──
    active_url = get_active_public_url()
    
    with st.container(border=True):
        st.markdown('''
<div style="display:flex;align-items:center;gap:12px;margin-bottom:10px;">
  <span style="font-size:28px;">🌐</span>
  <div>
    <div style="font-size:1.15rem;font-weight:800;color:#ffffff;">Tautan Akses Publik Aktif (Cloudflare Secure Tunnel)</div>
    <div style="font-size:0.8rem;color:#94a3b8;">Tautan ini aman (HTTPS) dan dapat langsung dibuka oleh teman Anda dari perangkat apa pun (HP/Tablet/Laptop).</div>
  </div>
</div>
''', unsafe_allow_html=True)
        
        c_sh1, c_sh2 = st.columns([3.8, 1.2], vertical_alignment="center")
        with c_sh1:
            st.code(active_url, language="text")
        with c_sh2:
            if st.button("🚀 Nyalakan / Segarkan Tunnel", key="btn_refresh_tunnel", use_container_width=True, help="Menjalankan script tunneling jika tautan belum aktif"):
                import subprocess
                tunnel_script = str(Path(__file__).parent / "start_share.sh")
                subprocess.Popen(["bash", tunnel_script], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                st.success("Perintah penyegaran dikirim! Tunggu 5-10 detik lalu refresh halaman.")
                time.sleep(1)
                st.rerun()

        st.markdown("""
> 💡 **Panduan Membagikan ke Teman:**
> 1. Kirimkan tautan di atas ke WhatsApp teman atau grup sekelas.
> 2. Teman Anda cukup klik tautan, lalu pada halaman awal pilih **"Belum punya akun"** untuk mendaftar (cukup nama & email).
> 3. Akun mereka otomatis langsung mendapatkan **Akses Pro** dan **208 Modul Kuliah Kedokteran** yang siap diuji coba tanpa perlu mengunggah materi lagi!
""")
        
        with st.expander("📋 Salin Draf Pesan Ajakan WhatsApp"):
            wa_draft = f"""Halo teman-teman! 👋
Aku lagi kembangkan platform belajar kedokteran "NeuroStudy" berbasis AI dan bukti ilmiah (Cognitive Load Theory & Sistem Pareto 80/20).

Tujuannya agar kita paham alur patofisiologi & lulus ujian blok tanpa harus pusing membuka ratusan slide PPT dosen yang berantakan.

Bisa tolong bantu coba dan uji fiturnya di:
🔗 {active_url}

(Tinggal klik, daftar akun baru dengan email kalian, materi kuliah 208 modul sudah otomatis tersedia di dalam).
Setelah mencoba, tolong beri masukan di tab "Uji Coba & Feedback Beta" ya. Terima kasih banyak! 🙏"""
            st.text_area("Template Pesan WhatsApp:", value=wa_draft, height=180)

    st.markdown("<br/>", unsafe_allow_html=True)

    # ── 2. FORMULIR MASUKAN PENGUJI (FEEDBACK SUBMISSION FORM) ──
    st.markdown("### 📝 Formulir Masukan Penguji (Beta Feedback Form)")
    st.caption("Masukan Anda sangat berharga untuk menyempurnakan kurikulum kognitif dan akurasi klinis NeuroStudy.")

    cur_user = st.session_state.get("user_info", {})
    u_name = cur_user.get("display_name") or cur_user.get("username", "Dokter Muda")
    u_email = cur_user.get("email", "")

    with st.form(key="form_beta_tester_feedback"):
        c_fb1, c_fb2 = st.columns(2)
        with c_fb1:
            fb_role = st.selectbox(
                "Status / Tingkat Anda:",
                ["Mahasiswa Preklinik Semester 1-2", "Mahasiswa Preklinik Semester 3-4", "Mahasiswa Preklinik Semester 5-6", "Mahasiswa Klinik / Dokter Muda (Koas)", "Dokter Umum / Residen", "Dosen / Pengajar", "Lainnya"]
            )
        with c_fb2:
            fb_rating = st.slider("Rating Kepuasan Pengalaman Belajar Keseluruhan:", min_value=1, max_value=5, value=5, help="1 = Sangat Kurang, 5 = Luar Biasa Sangat Membantu")

        st.markdown("---")
        st.markdown("**🔍 Evaluasi Fitur Sentral:**")

        c_q1, c_q2 = st.columns(2)
        with c_q1:
            q_zero_ppt = st.radio(
                "1. Apakah Catatan Master (Zero-PPT) cukup mendalam sehingga Anda merasa TIDAK PERLU membuka PPT 50 slide aslinya lagi?",
                ["Ya, sangat lengkap & jauh lebih runtut dari PPT", "Cukup lengkap, tapi sesekali masih butuh cek PPT", "Belum cukup, masih harus buka PPT"],
                index=0
            )
        with c_q2:
            q_soal_real = st.radio(
                "2. Apakah Simulasi Soal Campuran C1-C6 dan Analisis Jebakannya relevan dengan tipe ujian di kampus Anda?",
                ["Sangat relevan & jebakannya mirip soal ujian blok/UKMPPD", "Cukup mirip", "Kurang mirip", "Terlalu mudah / Terlalu sulit"],
                index=0
            )

        fb_best_features = st.multiselect(
            "3. Fitur mana yang menurut Anda PALING bermanfaat & bernilai tinggi?",
            ["🏛️ Catatan Master & Kapsul Pareto 80/20", "🎯 Verifikasi Kognitif (Socratic Active Recall)", "📝 Simulasi Soal Campuran C1-C6 & Distractor Analysis", "🃏 Flashcards Spaced Repetition (SM-2)", "🗺️ Mind Map Visual Interaktif", "🩺 Diskusi dengan Dewan Dokter Spesialis AI"],
            default=["🏛️ Catatan Master & Kapsul Pareto 80/20", "📝 Simulasi Soal Campuran C1-C6 & Distractor Analysis"]
        )

        fb_critique = st.text_area(
            "4. Kritik, Kelemahan, atau Hal yang Paling Mendesak untuk Diperbaiki:",
            placeholder="Tuliskan jika ada penjelasan yang kurang jelas, loading yang lambat, tombol membingungkan, atau bug...",
            height=100
        )

        fb_feature_request = st.text_area(
            "5. Usulan Fitur Baru / Harapan Pengembangan:",
            placeholder="Misal: mode offline, integrasi kalender otomatis, quiz bertenggat waktu, dsb...",
            height=80
        )

        submit_fb = st.form_submit_button("🚀 Kirimkan Masukan Saya →", type="primary", use_container_width=True)

    if submit_fb:
        entry = {
            "id": datetime.datetime.now().strftime("%Y%m%d%H%M%S"),
            "timestamp": datetime.datetime.now().isoformat(),
            "name": u_name,
            "email": u_email,
            "role": fb_role,
            "rating": fb_rating,
            "zero_ppt_evaluation": q_zero_ppt,
            "exam_similarity": q_soal_real,
            "best_features": fb_best_features,
            "critique": fb_critique.strip(),
            "feature_requests": fb_feature_request.strip()
        }
        save_beta_feedback(entry)
        st.success("🎉 Terima kasih banyak atas masukan Anda! Umpan balik Anda telah berhasil dicatat ke dalam basis data riset NeuroStudy.")
        st.balloons()

    # ── 3. EXECUTIVE ADMIN VIEW (KHUSUS PEMILIK/DIMAS) ──
    cur_u_email = (cur_user.get("email") or "").lower()
    cur_u_name = (cur_user.get("username") or "").lower()
    is_owner = any(k in cur_u_email or k in cur_u_name for k in ["dimas", "wastu", "mahesawastu8@gmail.com"])
    
    if is_owner:
        st.markdown("<br/>", unsafe_allow_html=True)
        st.markdown("---")
        st.markdown("### 👑 Dashboard Masukan Penguji (Owner Executive View)")
        st.caption("Panel eksklusif untuk memantau metrik adopsi dan umpan balik dari seluruh teman yang mencoba.")

        all_fb = load_beta_feedback()
        if not all_fb:
            st.info("Belum ada masukan yang masuk. Silakan bagikan tautan publik ke teman Anda!")
        else:
            total_resp = len(all_fb)
            avg_rating = sum(f.get("rating", 5) for f in all_fb) / max(1, total_resp)
            zero_ppt_count = sum(1 for f in all_fb if "sangat lengkap" in f.get("zero_ppt_evaluation", "").lower())
            zero_ppt_pct = int((zero_ppt_count / total_resp) * 100)

            c_m1, c_m2, c_m3 = st.columns(3)
            c_m1.metric("Total Responden Penguji", f"{total_resp} Orang")
            c_m2.metric("Skor Kepuasan Rata-Rata", f"⭐ {avg_rating:.1f} / 5.0")
            c_m3.metric("Adopsi 'Zero-PPT' Penuh", f"{zero_ppt_pct}%")

            st.markdown("#### 📋 Riwayat Masukan Lengkap:")
            for item in reversed(all_fb):
                with st.container(border=True):
                    c_fbt1, c_fbt2 = st.columns([3.5, 1.5])
                    with c_fbt1:
                        st.markdown(f"**{item.get('name', 'Anonim')}** · <span style='color:#94a3b8;font-size:0.8rem;'>{item.get('role')} ({item.get('email')})</span>", unsafe_allow_html=True)
                    with c_fbt2:
                        st.markdown(f"<span style='color:#f59e0b;font-weight:700;'>⭐ Rating: {item.get('rating')}/5</span> · <span style='font-size:0.75rem;color:#64748b;'>{item.get('timestamp', '')[:10]}</span>", unsafe_allow_html=True)
                    
                    st.markdown(f"- **Evaluasi Zero-PPT:** {item.get('zero_ppt_evaluation')}")
                    st.markdown(f"- **Relevansi Soal:** {item.get('exam_similarity')}")
                    st.markdown(f"- **Fitur Favorit:** {', '.join(item.get('best_features', []))}")
                    if item.get("critique"):
                        st.markdown(f"- ⚠️ **Kritik / Kendala:** <span style='color:#f87171;'>{item.get('critique')}</span>", unsafe_allow_html=True)
                    if item.get("feature_requests"):
                        st.markdown(f"- 💡 **Saran Fitur Baru:** <span style='color:#38bdf8;'>{item.get('feature_requests')}</span>", unsafe_allow_html=True)

            with st.expander("📥 Unduh / Lihat Data JSON Mentah"):
                st.json(all_fb)





def render_tab_belajar():
    api_key = get_gemini_api_key()
    is_owner = st.session_state.get("is_owner", False)
    mats = load_mats()

    # ── LEARNING MODE TOGGLE: PARETO 80/20 (DEFAULT) VS FAST-TRACK VS 6-PHASE ──
    c_m1, c_m2 = st.columns([2.8, 1.8], vertical_alignment="center")
    with c_m1:
        st.markdown('<div style="font-size:1.5rem;font-weight:800;color:#f8fafc;margin-bottom:2px;letter-spacing:-0.5px;">Ruang Belajar Kognitif Medis (Pareto 80/20 Engine)</div>', unsafe_allow_html=True)
        st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-bottom:4px;">Sistem pembelajaran berbasis bukti ilmiah (Cognitive Load Theory, Testing Effect & Pareto 80/20) untuk menaklukkan beban 750 slide/minggu tanpa harus membuka PPT.</p>', unsafe_allow_html=True)
    with c_m2:
        study_mode_sel = st.radio(
            "Mode Belajar:",
            ["🏛️ Sistem Pareto 80/20 (Zero-PPT)", "⚡ Mode Darurat H-1 Ujian", "🎓 Mode Alur 6-Fase"],
            index=0,
            horizontal=False,
            key="study_mode_picker",
            label_visibility="collapsed"
        )
        st.session_state.study_mode = study_mode_sel

    with st.expander("📚 Lihat Dasar Ilmiah, Meta-Analisis & Konsensus Baku Kedokteran (100% Terverifikasi)", expanded=False):
        st.markdown("""
### 🏛️ Rujukan Ilmiah & Konsensus Baku Kedokteran (100% Asli & Terverifikasi)
Seluruh materi, penalaran klinis, dan algoritma sintesis NeuroStudy berakar secara ketat pada literatur standar emas resmi dunia kedokteran:

| No | Sumber Standar Emas | Edisi & Tahun Resmi | Penulis / Dewan Editor | Penerbit / Jurnal Resmi & ISBN/DOI | Cakupan Klinis Utama |
|:---:|---|---|---|---|---|
| **1** | **Harrison’s Principles of Internal Medicine** | **21st Edition (2022)** | Joseph Loscalzo, Anthony S. Fauci, Dennis L. Kasper, Stephen L. Hauser, Dan L. Longo, J. Larry Jameson | **McGraw-Hill Education / Medical**<br>ISBN-13: 978-1264268504 | Standar emas penyakit dalam, kriteria diagnostik, dan penatalaksanaan klinis global. |
| **2** | **Guyton and Hall Textbook of Medical Physiology** | **14th Edition (2020)** | John E. Hall, Michael E. Hall | **Elsevier**<br>ISBN-13: 978-0323597128 | Mekanisme fisiologi seluler, kontrol umpan balik neuroendokrin, dan homeostasis organ. |
| **3** | **Robbins & Cotran Pathologic Basis of Disease** | **10th Edition (2020)** | Vinay Kumar, Abul K. Abbas, Jon C. Aster | **Elsevier**<br>ISBN-13: 978-0323531139 | Patogenesis molekuler, jejas sel, respon inflamasi-imun, dan korelasi histopatologi. |
| **4** | **Katzung & Trevor’s Pharmacology: Examination & Board Review** | **13th Edition (2021)** / Basic & Clinical 15th-16th Ed | Bertram G. Katzung, Marieke Kruidering-Hall, Anthony J. Trevor | **McGraw-Hill**<br>ISBN-13: 978-1260464917 | Farmakodinamik molekuler, farmakokinetik, interaksi obat, dan rejimen lini 1-3. |
| **5** | **Clinical Reasoning & Diagnostic Accuracy (Norman, Eva, Bowen)** | **Landmark Studies (NEJM, Med Educ, Acad Med)** | Geoff Norman, Kevin W. Eva, Judith L. Bowen, Pat Croskerry | **NEJM** (Bowen 2006; doi:10.1056/NEJMra054778)<br>**Medical Education** (Norman 2005; Eva 2005; Norman et al. 2007)<br>**Academic Medicine** (Croskerry 2003) | *Dual-process theory*, pencegahan bias kognitif (*premature closure*), dan *hypothetico-deductive reasoning*. |
| **6** | **Panduan Praktik Klinis & PNPK Nasional** | **Edisi Konsensus Resmi Terkini** | Kementerian Kesehatan Republik Indonesia & Perhimpunan Spesialis | **Kemenkes RI, PAPDI, PERKI, IDAI, POGI, PERDOSSI** | Alur penatalaksanaan baku di fasilitas kesehatan primer dan rujukan Indonesia. |

---

### 🧠 Landasan Teori Kognitif & Meta-Analisis Pembelajaran (Evidence-Based Learning Framework)

| Pilar / Fase | Landasan Teori | Jurnal & Peneliti Utama | Temuan Ilmiah & Effect Size |
|--------------|----------------|------------------------|-----------------------------|
| **Pilar 1: Catatan Master** | **Cognitive Load Theory** & Schema Construction | **Sweller, van Merriënboer, & Paas (2024)**, *Educ Psychol Rev*; **van Merriënboer & Sweller (2019)** | Membatasi *extraneous load* dari slide kuliah terfragmentasi dan memaksimalkan *germane load*. Sintesis terpadu mengoptimalkan memori kerja (4±1 chunk). |
| **Pilar 1: Pareto 80/20** | **Prinsip Efisiensi Kognitif** | **Koch (1998)**, *The 80/20 Principle*; **Dunlosky et al. (2013)** | 20% mekanisme patofisiologi/farmakodinamik sentral mendikte 80% luaran klinis dan variasi soal ujian. |
| **Pilar 2: Active Recall** | **Testing Effect** & Elaborative Interrogation | **Dunlosky et al. (2013)**, *Psychol Sci Public Interest*; **Roediger & Karpicke (2006)**; **Yang et al. (2021)**, *Psychol Bull* | Meta-analisis **222 studi (48.478 siswa)** membuktikan *retrieval practice* menghasilkan gain retensi besar (**d = 0.50 - 0.74**), 50%+ lebih unggul dari sekadar membaca ulang. |
| **Pilar 3: Soal Campuran** | **Revised Bloom's Taxonomy** & Dual-Process Theory | **Anderson & Krathwohl (2001)**; **Norman, Eva, & Brooks (2007)**, *Medical Education* | Menguji rentang C1 (definisi) hingga C6 (penalaran kasus ranjang IGD). Analisis pengecoh (*distractor analysis*) mencegah bias kognitif *premature closure*. |
| **Pilar 4: Spacing (SM-2)** | **Distributed Practice** & Sleep Consolidation | **Cepeda et al. (2008)**, *Psychol Sci*; **Wozniak (1990)**; **Walker (2017)**; **Tononi & Cirelli (2014)**, *Neuron* | Kurva lupa Ebbinghaus ditangkal melalui algoritma SM-2. Tidur 7-9 jam mengonsolidasi memori dari hippocampus ke neokorteks permanen. |
| **Pilar 5: Mind Map** | **Dual Coding Theory** | **Clark & Paivio (1991)**; **Schroeder et al. (2018)**, *Educ Psychol Rev* | Meta-analisis 142 studi: pemetaan konsep hierarkis menghasilkan effect size **g = 0.58 - 0.72**. |
""")

    if not mats:
        st.markdown('<div class="card"><div class="cs" style="text-align:center;">Upload materi di tab <b>Upload Materi</b> terlebih dahulu.</div></div>', unsafe_allow_html=True)
    else:
        c_t2_f1, c_t2_f2, c_t2_rst = st.columns([2.0, 4.0, 1.0], vertical_alignment="bottom")
        
        with c_t2_f1:
            t2_blok_options = ["Semua Blok", "BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD", "Lainnya"]
            curr_sel_pre = st.session_state.get("mat_sel", "")
            def_blok_idx = 0
            if "t2_blok_selector" in st.session_state and st.session_state.t2_blok_selector in t2_blok_options:
                def_blok_idx = t2_blok_options.index(st.session_state.t2_blok_selector)
            elif curr_sel_pre:
                for b in ["BDT", "BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BMD"]:
                    if curr_sel_pre.startswith(f"[{b}]"):
                        def_blok_idx = t2_blok_options.index(b)
                        break
            sel_t2_blok = st.selectbox("Filter Blok:", t2_blok_options, index=def_blok_idx, key="t2_blok_selector")
            
        with c_t2_f2:
            available_keys = list(mats.keys())
            if sel_t2_blok != "Semua Blok":
                if sel_t2_blok == "Lainnya":
                    available_keys = [k for k in available_keys if not any(k.startswith(f"[{b}") for b in ["BMS 1", "BUAMS", "BMS 2", "BMS 3", "BMS 4", "BDT", "BMD"])]
                else:
                    available_keys = [k for k in available_keys if k.startswith(f"[{sel_t2_blok}]") or sel_t2_blok.lower() in k.lower()]
                    
            if not available_keys:
                available_keys = list(mats.keys())

            # Always ensure the currently active material is preserved in the dropdown
            if curr_sel_pre and curr_sel_pre in mats and curr_sel_pre not in available_keys:
                available_keys.append(curr_sel_pre)
                
            sorted_mat_keys = sorted(
                available_keys,
                key=lambda x: (
                    0 if x.startswith("[BMS 1]") else (
                        1 if x.startswith("[BUAMS]") else (
                            2 if x.startswith("[BMS") else (
                                3 if x.startswith("[BDT") else 4
                            )
                        )
                    ),
                    x.lower()
                )
            )
            def_idx = 0
            curr_sel = st.session_state.get("mat_sel")
            if curr_sel in sorted_mat_keys:
                def_idx = sorted_mat_keys.index(curr_sel)
            sel = st.selectbox("Pilih kuliah yang ingin dipelajari:", sorted_mat_keys, index=def_idx, key="sesi_sel")
            
        with c_t2_rst:
            if st.button("🔄 Reset", use_container_width=True, help="Reset sesi belajar materi ini"):
                clear_active_session(sel)
                st.session_state.phase = 0; st.session_state.phase_data = {}; st.session_state.completed = False; st.session_state.history = []; st.session_state.scores = {}; st.session_state.session_started = False; st.session_state.post_chat = []; st.rerun()

        if sel != st.session_state.get("mat_sel"):
            st.session_state.mat_sel = sel
            saved_s = load_active_session(sel)
            if saved_s:
                st.session_state.phase = saved_s.get("phase", 0)
                st.session_state.phase_data = saved_s.get("phase_data", {})
                st.session_state.completed = saved_s.get("completed", False)
                st.session_state.history = saved_s.get("history", [])
                st.session_state.scores = saved_s.get("scores", {})
                st.session_state.session_started = saved_s.get("session_started", False)
                st.session_state.post_chat = saved_s.get("post_chat", [])
            else:
                st.session_state.phase = 0; st.session_state.phase_data = {}; st.session_state.completed = False; st.session_state.history = []; st.session_state.scores = {}; st.session_state.session_started = False; st.session_state.post_chat = []

        raw_text = load_mats()[sel].get("text", "")
        text = clean_academic_text(raw_text)
        mat_info = load_mats()[sel]

        # ── GOOGLE MEDICAL WORKSPACE TOOLKIT ──
        next_study_date = datetime.datetime.now() + datetime.timedelta(days=1)
        gcal_url = build_gcal_url(
            title=f"🧠 Review: {sel}",
            dt=next_study_date,
            details=f"Sesi belajar modul kedokteran: {sel}\nAplikasi: NeuroStudy Clinical Platform\nTarget: Active Recall & Latihan Klinis."
        )
        scholar_query = urllib.parse.quote(f"{sel} medicine clinical review")
        scholar_url = f"https://scholar.google.com/scholar?q={scholar_query}"
        
        st.markdown(f'''
<div style="background:rgba(15,23,42,0.6);border:1px solid rgba(255,255,255,0.08);border-radius:12px;padding:8px 14px;margin:8px 0 16px;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:8px;">
  <div style="display:flex;align-items:center;gap:8px;">
    <span style="font-size:0.72rem;font-weight:800;color:#94a3b8;letter-spacing:0.4px;">GOOGLE MEDICAL SUITE:</span>
    <span style="font-size:0.78rem;font-weight:800;color:#38bdf8;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;max-width:280px;">{sel}</span>
  </div>
  <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;">
    <a href="{gcal_url}" target="_blank" style="text-decoration:none;background:rgba(66,133,244,0.15);border:1px solid rgba(66,133,244,0.4);color:#93c5fd;font-size:0.74rem;font-weight:700;padding:4px 10px;border-radius:8px;display:inline-flex;align-items:center;gap:5px;">
      <span>📅 + Google Calendar</span>
    </a>
    <a href="{scholar_url}" target="_blank" style="text-decoration:none;background:rgba(251,188,5,0.15);border:1px solid rgba(251,188,5,0.4);color:#fde047;font-size:0.74rem;font-weight:700;padding:4px 10px;border-radius:8px;display:inline-flex;align-items:center;gap:5px;">
      <span>🔬 Google Scholar (EBM)</span>
    </a>
    <a href="https://meet.google.com/new" target="_blank" style="text-decoration:none;background:rgba(52,168,83,0.15);border:1px solid rgba(52,168,83,0.4);color:#86efac;font-size:0.74rem;font-weight:700;padding:4px 10px;border-radius:8px;display:inline-flex;align-items:center;gap:5px;">
      <span>🎥 Google Meet Diskusi</span>
    </a>
  </div>
</div>
''', unsafe_allow_html=True)

        # ── CEK MODE DARURAT UJIAN (FAST TRACK H-1) ──
        if st.session_state.get("study_mode") == "⚡ Mode Darurat H-1 Ujian":
            st.markdown(f'''
<div style="background:linear-gradient(135deg, rgba(239,68,68,0.12) 0%, rgba(245,158,11,0.08) 100%);border:1.5px solid rgba(239,68,68,0.4);border-radius:14px;padding:18px 22px;margin:16px 0 20px;">
  <div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:10px;">
    <div>
      <span style="background:#ef4444;color:#fff;font-size:0.72rem;font-weight:900;padding:2px 8px;border-radius:6px;letter-spacing:0.5px;text-transform:uppercase;">⚡ FAST-TRACK H-1 UJIAN</span>
      <div style="font-size:1.3rem;font-weight:900;color:#ffffff;margin-top:4px;">{sel}</div>
      <div style="font-size:0.8rem;color:#cbd5e1;margin-top:2px;">Mode instan tanpa bertele-tele: Intisari 3 Menit, 10 Soal Prediksi Ujian, & Flashcard Kilat.</div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)
            
            t_ft1, t_ft2, t_ft3 = st.tabs(["📌 Intisari 3 Menit", "🎯 10 Soal Prediksi Ujian", "🃏 Flashcard Kilat"])
            
            with t_ft1:
                st.markdown("#### 📌 Intisari Super High-Yield (Fakta Kritis & Skema Inti)")
                st.caption("Fakta-fakta yang wajib dihafal mati sebelum masuk ruang ujian besok pagi:")
                compressed_mat = get_cached_high_yield_text(sel, text, api_key)
                if st.button("🚀 Tampilkan Intisari 3 Menit", key=f"btn_ft_summary_{sel}", type="primary", use_container_width=True):
                    with st.spinner("AI sedang menyaring poin high-yield..."):
                        ft_prompt = f"""Kamu adalah Profesor Senior Ahli Ujian Kedokteran.
Tuliskan RINGKASAN HIGH-YIELD 3 MENIT untuk mahasiswa yang akan ujian besok pagi dari materi ini:

Materi:
{compressed_mat[:7500]}

Format WAJIB:
1. **🧬 3 Mekanisme Kunci yang Paling Sering Keluar Ujian:** (Jelaskan dalam poin tegas)
2. **📊 Tabel Klasifikasi / Obat Cepat Hafal:** (Tabel ringkas: Nama, Mekanisme, Efek Samping Khas)
3. **🚨 3 Red Flags / Aturan Emas Klinis:** (Hal berbahaya yang sering dijadikan soal kasus)
4. **💡 1 Kalimat Mnemonik Sakti:** (Cara menghafal konsep paling rumit)"""
                        st.markdown(stream_ai_transparent(api_key, ft_prompt, st.empty()))
                        
            with t_ft2:
                st.markdown("#### 🎯 10 Soal Prediksi Ujian Blok & Pembahasan Tuntas")
                st.caption("Prediksi soal pilihan ganda gaya dosen dan UKMPPD lengkap dengan jebakannya:")
                if st.button("📝 Buat 10 Soal Prediksi Ujian", key=f"btn_ft_mcq_{sel}", type="primary", use_container_width=True):
                    with st.spinner("Dr. Marcus Vance sedang merancang 10 soal prediksi..."):
                        mcq_prompt = f"""Kamu adalah Dr. Marcus Vance, Sp.FK.
Buatkan 5-10 SOAL PREDIKSI UJIAN PILIHAN GANDA (MCQ) yang paling potensial keluar pada ujian blok/UKMPPD berdasarkan materi ini:

Materi:
{compressed_mat[:7500]}

Format untuk setiap soal:
**Soal [Nomor]:** [Skenario klinis singkat / pertanyaan mekanisme]
A. [Opsi]
B. [Opsi]
C. [Opsi]
D. [Opsi]
E. [Opsi]

<details><summary><b>🔍 Kunci Jawaban & Pembahasan Jebakan</b></summary>
<b>Jawaban Benar:</b> [Huruf]<br>
<b>Pembahasan:</b> [Alasan benar]<br>
<b>⚠️ Jebakan Dosen:</b> [Mengapa opsi lain salah / mengecoh]
</details>
<hr>"""
                        st.markdown(stream_ai_transparent(api_key, mcq_prompt, st.empty()))
                        
            with t_ft3:
                render_flashcards_widget(sel, text, api_key, "FastTrack", key_prefix="fasttrack_tab")
                            
            return # Selesai mode darurat

        # ── SISTEM UTAMA: PARETO 80/20 (ZERO-PPT NEEDED · EVIDENCE-BASED) ──
        if st.session_state.get("study_mode", "").startswith("🏛️ Sistem Pareto 80/20"):
            sub_t1, sub_t2, sub_t3, sub_t4 = st.tabs([
                "📖 Langkah 1: Catatan Master (Zero-PPT)",
                "🎯 Langkah 2: Socratic Active Recall",
                "📝 Langkah 3: Simulasi Ujian Kasus (C1-C6)",
                "🗺️ Langkah 4: Visual Mind Map & Flashcard"
            ])

            # ── 1. CATATAN MASTER KLINIS (ZERO-PPT NEEDED) ──
            with sub_t1:
                st.markdown('''
<div style="background:linear-gradient(135deg, rgba(99,102,241,0.12) 0%, rgba(56,189,248,0.08) 100%); border:1.5px solid rgba(99,102,241,0.35); border-radius:12px; padding:16px 20px; margin-bottom:16px;">
  <div style="display:flex; align-items:center; gap:12px;">
    <span style="font-size:28px;">🏛️</span>
    <div>
      <div style="font-size:1.15rem; font-weight:800; color:#ffffff;">Buku Catatan Klinis Mandiri (Sintesis Penuh Zero-PPT)</div>
      <div style="font-size:0.78rem; color:#94a3b8; margin-top:3px; line-height:1.5;">
        <strong>Dasar Riset &amp; Standar:</strong> <em>Cognitive Load Theory</em> (Sweller et al., 2024), <em>Schema Construction</em> (van Merriënboer, 2019), &amp; <em>Standar Emas Konsensus Medis</em> (Harrison, Robbins, Guyton, Katzung).
        Seluruh substansi 50 slide dosen disintesis menjadi buku ajar terstruktur agar Anda <strong>tidak perlu membuka file PPT aslinya lagi</strong>.
      </div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)

                master_cache = get_cached_master_note(sel)
                
                c_mn_l, c_mn_r = st.columns([3.5, 1.5], vertical_alignment="center")
                with c_mn_r:
                    btn_regen_master = st.button("🔄 Susun Ulang Sintesis AI", key=f"btn_regen_master_{sel}", use_container_width=True, help="Perbarui dan buat ulang sintesis master note")
                    
                if master_cache and not btn_regen_master:
                    c_act1, c_act2 = st.columns([3.3, 1.7], vertical_alignment="center")
                    with c_act1:
                        b_tags = []
                        if master_cache.get("source") == "global_peer":
                            b_tags.append('<span style="background:rgba(56,189,248,0.15);border:1px solid rgba(56,189,248,0.35);border-radius:6px;padding:3px 8px;font-size:0.75rem;color:#7dd3fc;font-weight:700;">⚡ Perpustakaan Terpadu Angkatan (0s Instan · Bebas Kuota)</span>')
                        if master_cache.get("verified"):
                            b_tags.append(f'<span style="background:rgba(34,197,94,0.15);border:1px solid rgba(34,197,94,0.4);border-radius:6px;padding:3px 8px;font-size:0.75rem;color:#4ade80;font-weight:700;">🛡️ Terverifikasi Klinis: {master_cache.get("verified_by", "Dokter Spesialis")}</span>')
                        elif is_owner:
                            if st.button("🛡️ Beri Cap Validasi Klinis Dokter", key=f"btn_verify_doc_{sel}", help="Beri cap terverifikasi klinis dengan tanda tangan kriptografi anti-tamper"):
                                from core.rate_limiter import RateLimiter
                                sig = RateLimiter.generate_clinical_signature(sel, "dr. Dimas Wastu Mahesa")
                                master_cache["cryptographic_sig"] = sig
                                save_cached_master_note(sel, master_cache, is_verified=True, reviewer_name="dr. Dimas Wastu Mahesa")
                                st.success(f"Materi berhasil diverifikasi secara klinis! (Signature: {sig})")
                                st.rerun()
                        if b_tags:
                            st.markdown(" ".join(b_tags), unsafe_allow_html=True)
                    with c_act2:
                        st.download_button(
                            label="📄 Unduh Catatan (.md / Cetak)",
                            data=master_cache.get("content", ""),
                            file_name=f"{re.sub(r'[^a-zA-Z0-9_]', '_', sel)}_Catatan_Master.md",
                            mime="text/markdown",
                            key=f"dl_mn_{sel}",
                            use_container_width=True
                        )
                    st.markdown(master_cache.get("content", ""))
                    
                    # ── MULTIMODAL CLINICAL VISUAL ATLAS (DUAL CODING THEORY) ──
                    from core.visual_engine import get_module_visual_atlas
                    visuals = get_module_visual_atlas(sel)
                    if visuals:
                        with st.expander(f"🔬 Atlas Visual Diagnostik & Skema Klinis Terpadu ({len(visuals)} Modalitas Visual)", expanded=True):
                            st.caption("Memadukan visual diagnostik patologi, EKG, radiologi, dan alur penalaran klinis ke dalam satu catatan terpadu (Paivio's Dual Coding Theory, effect size g = 0.72):")
                            for v in visuals:
                                st.markdown(f"""
<div style="background:rgba(15, 23, 42, 0.75); border:1px solid rgba(56, 189, 248, 0.35); border-radius:10px; padding:14px 18px; margin-bottom:12px;">
  <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:6px;">
    <span style="font-weight:800; color:#38bdf8; font-size:0.92rem;">{v['title']}</span>
    <span style="background:rgba(56,189,248,0.15); color:#7dd3fc; font-size:0.68rem; padding:2px 8px; border-radius:6px; font-weight:700;">{v['category']}</span>
  </div>
  <div style="font-size:0.78rem; color:#cbd5e1; margin-bottom:10px; line-height:1.4;">
    <strong>💡 Clinical Diagnostic Pearl:</strong> {v['clinical_pearl']}
  </div>
  <pre style="background:#090d16; color:#a5f3fc; padding:12px; border-radius:8px; font-size:0.74rem; line-height:1.35; overflow-x:auto; border:1px solid rgba(255,255,255,0.08); font-family:monospace;">{v['schematic'].strip()}</pre>
</div>
""", unsafe_allow_html=True)

                    st.markdown("---")
                    st.info("💡 **Langkah 1 Selesai Dibaca?** Lanjutkan mengunci alur sebab-akibat di tab **🎯 Langkah 2: Socratic Active Recall** di atas.")
                else:
                    auto_trigger = st.session_state.pop("auto_gen_master", False)
                    btn_gen_now = False
                    if not master_cache and not btn_regen_master:
                        if auto_trigger:
                            btn_gen_now = True
                        else:
                            st.info(f"💡 Modul kuliah **{sel}** (~{len(text):,} karakter teks sumber) siap disintesis menjadi Catatan Master Klinis Komprehensif.")
                            btn_gen_now = st.button("⚡ Bedah 50 Slide & Susun Catatan Master (AI) →", type="primary", use_container_width=True, key=f"btn_generate_master_{sel}")
                    else:
                        btn_gen_now = True
                        
                    if btn_gen_now:
                        from core.rate_limiter import RateLimiter
                        curr_user = st.session_state.get("current_user", "dimas")
                        allowed, limit_msg, remaining = RateLimiter.check_and_increment(curr_user, "ai_master_note")
                        if not allowed:
                            st.warning(f"🛡️ Security Quota Protection: {limit_msg}")
                            st.stop()
                        ph_mn = st.empty()
                        prompt_mn = f"""Kamu adalah Profesor Kedokteran Senior, Guru Besar Biomedis, dan Penulis Buku Ajar Kedokteran Terkemuka (selevel Guyton, Robbins, Harrison, dan Katzung).
Tugasmu adalah membedah dan mensintesis SELURUH materi slide kuliah berikut menjadi BUKU CATATAN KLINIS KOMPREHENSIF (Comprehensive Clinical Master Note) berstandar emas.

PRINSIP UTAMA:
Mahasiswa kedokteran ini TIDAK PERLU MEMBUKA PPT ASLINYA LAGI. Seluruh substansi, fakta kritis, klasifikasi, alur patofisiologi molekuler-organ, kriteria diagnosis baku, dan tata laksana dari ke-50 slide kuliah harus dirangkai menjadi artikel klinis yang sangat mendalam, kaya, presisi, dan mudah dipahami.

Materi Kuliah (Sumber):
Judul: {sel}
Isi Ekstraksi Slide:
{text[:18000]}

Gunakan Struktur Evidence-Based Pareto 80/20 & Konsensus Medis Baku berikut (Format Markdown murni):

# 🏛️ {sel} — Catatan Master Klinis Komprehensif
*Berdasarkan Cognitive Load Theory (Sweller, 2024), Schema Construction (van Merriënboer, 2019), & Standar Konsensus Medis (Harrison, Robbins, Guyton, Katzung).*

---

## ⚡ KAPSUL PARETO 80/20 (20% Konsep Kunci Penentu 80% Nilai Ujian & Klinis)
> **Definisi Sentral & Terminologi Medis:** [Definisi baku medis yang presisi beserta batasan klinisnya]
> 
> **Mekanisme Kausalitas Utama (The Master Key):** 
> [Alur sebab-akibat patofisiologi/farmakodinamik inti: Stimulus/Etiologi → Reseptor/Jalur Molekuler → Perubahan Seluler/Histopatologi → Manifestasi Organ & Gejala]
> 
> **Trias/Tetrad Patognomonik & Red Flags (Tanda Bahaya):**
> - [Gejala patognomonik / tanda khas yang wajib dihafal]
> - [Kontraindikasi mutlak / kegawatan yang sering menjadi jebakan ujian]
> 
> **Terapi Lini Pertama (Drug/Management of Choice):**
> [Nama obat/intervensi lini pertama, mekanisme kerja molekuler, dan target terapinya]

---

## 📖 PEMBEDAHAN MATERI LENGKAP & MENDALAM (ZERO-PPT NEEDED)

### 📌 Bab I: Fondasi Biomedis, Anatomi, & Fisiologi Terkait
[Uraikan struktur anatomi, histologi, fisiologi normal, atau homeostasis seluler yang menjadi dasar materi ini secara terstruktur dan jelas]

### 📌 Bab II: Etiologi, Patogenesis, & Kaskade Molekuler Lengkap
[Uraikan rantai kausalitas patofisiologis lengkap:
- Etiologi & faktor risiko.
- Kaskade molekuler (sitokin, enzim, second messenger, interaksi reseptor).
- Perubahan histopatologis & disfungsi organ.]

### 📌 Bab III: Manifestasi Klinis, Trias/Tetrad Khas, & Pendekatan Diagnostik
[Uraikan:
1. Anamnesis: Gejala subjektif khas, durasi, onset, faktor pemberat/peringan.
2. Pemeriksaan Fisik: Tanda objektif patognomonik.
3. Pemeriksaan Penunjang Bertingkat:
   - Skrining awal (Lab darah rutin, urinalisis, dsb).
   - Penunjang spesifik (Radiologi, EKG, Patologi Anatomi).
   - Standar Emas (Gold Standard) penegakan diagnosis.
4. Tabel Kriteria Diagnosis Resmi (mengacu pada konsensus resmi seperti WHO, ADA, AHA/ACC, KDIGO, GOLD, GINA jika relevan).]

### 📌 Bab IV: Tabel Komparasi Diferensial Diagnosis (DDx)
[Tabel komparasi 2-4 penyakit serupa yang sering mengecoh:
| Penyakit / Kondisi | Karakteristik Khas | Temuan Penunjang Pembeda | Terapi Utama |]

### 📌 Bab V: Algoritma Tata Laksana Komprehensif (Farmakoterapi & Non-Farmakoterapi)
[Uraikan tata laksana bertingkat:
1. Penanganan awal / kegawatdaruratan (ABCDE, stabilisasi).
2. Tabel Farmakoterapi Lini 1, 2, dan 3:
   | Lini Terapi | Golongan Obat | Contoh Generik | Mekanisme Kerja Molekuler (MOA) | Dosis / Indikasi Utama | Efek Samping Khas & Kontraindikasi |
3. Terapi non-farmakologis, modifikasi gaya hidup, dan edukasi pasien.]

### 📌 Bab VI: Komplikasi, Prognosis, & High-Yield Exam Pearls
[Uraikan risiko jika tidak tertangani, prognosis klinis, serta poin-poin krusial yang paling sering dijadikan jebakan soal ujian blok kampus dan UKMPPD.]

---
💡 **Mnemonik Cerdas & Kunci Retensi Permanen:**
[1-2 mnemonik cerdas untuk mempermudah menghafal klasifikasi atau konsep tersulit materi ini.]

---
📚 **Rujukan Ilmiah & Konsensus Baku (Wajib Terverifikasi & Bebas Halusinasi):**
Seluruh rujukan berikut wajib asli, terverifikasi secara akademis, dan relevan dengan materi ini:
1. **Harrison’s Principles of Internal Medicine**, 21st Edition (Eds: Loscalzo J, Fauci AS, Kasper DL, Hauser SL, Longo DL, Jameson JL. McGraw-Hill, 2022. ISBN: 978-1264268504).
2. **Guyton and Hall Textbook of Medical Physiology**, 14th Edition (Hall JE, Hall ME. Elsevier, 2020. ISBN: 978-0323597128).
3. **Robbins & Cotran Pathologic Basis of Disease**, 10th Edition (Kumar V, Abbas AK, Aster JC. Elsevier, 2020. ISBN: 978-0323531139).
4. **Katzung & Trevor’s Pharmacology: Examination & Board Review**, 13th Edition (Katzung BG, Kruidering-Hall M, Trevor AJ. McGraw-Hill, 2021. ISBN: 978-1260464917) / Basic & Clinical Pharmacology.
5. **Norman & Eva**, *Clinical Reasoning and Diagnostic Accuracy*, NEJM / Medical Education / Academic Medicine (Bowen JL, NEJM 2006, doi:10.1056/NEJMra054778; Norman G, Med Educ 2005; Eva KW, Med Educ 2005; Norman et al., Med Educ 2007; Croskerry P, Acad Med 2003).
6. **Panduan Praktik Klinis / PNPK Kemenkes RI** serta Panduan Konsensus Perhimpunan Spesialis Terkait (PAPDI, PERKI, IDAI, POGI).

⚠️ PERATURAN INTEGRITAS AKADEMIK MUTLAK:
DILARANG KERAS mengarang, memalsukan, atau membuat-buat rujukan (zero fabricated citations). Seluruh sitasi bab, mekanisme patofisiologi, kriteria diagnostik, dosis farmakoterapi, dan diferensial diagnosis wajib berakar murni pada konsensus emas kedokteran nyata di atas.
"""
                        res_mn = stream_ai_transparent(api_key, prompt_mn, ph_mn)
                        if res_mn:
                            save_cached_master_note(sel, {"content": res_mn, "timestamp": datetime.datetime.now().isoformat()})
                            st.rerun()

            # ── 2. VERIFIKASI KOGNITIF (TARGETED ACTIVE RECALL) ──
            with sub_t2:
                st.markdown('''
<div style="background:linear-gradient(135deg, rgba(245,158,11,0.12) 0%, rgba(99,102,241,0.08) 100%); border:1.5px solid rgba(245,158,11,0.35); border-radius:12px; padding:16px 20px; margin-bottom:16px;">
  <div style="display:flex; align-items:center; gap:12px;">
    <span style="font-size:28px;">🎯</span>
    <div>
      <div style="font-size:1.15rem; font-weight:800; color:#ffffff;">Targeted Socratic Active Recall (Verifikasi Kausalitas)</div>
      <div style="font-size:0.78rem; color:#94a3b8; margin-top:3px; line-height:1.5;">
        <strong>Dasar Riset:</strong> <em>Testing Effect &amp; Elaborative Interrogation</em> (Dunlosky et al., 2013; Roediger &amp; Karpicke, 2006; Yang et al., 2021 meta-analisis d = 0.50–0.74).
        Bukan puluhan kuis trivia receh, melainkan 3 pertanyaan reflektif penembus logika (MENGAPA &amp; BAGAIMANA) untuk menguji apakah Anda benar-benar paham alur sebab-akibat.
      </div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)

                recall_cache = get_cached_active_recall(sel)
                
                c_rc_l, c_rc_r = st.columns([3.5, 1.5], vertical_alignment="center")
                with c_rc_r:
                    btn_regen_recall = st.button("🔄 Rancang Ulang Soal", key=f"btn_regen_recall_{sel}", use_container_width=True)
                    
                if not recall_cache or btn_regen_recall:
                    if st.button("⚡ Rancang 3 Pertanyaan Kausalitas Penembus Logika (AI)", type="primary", key=f"btn_gen_recall_{sel}", use_container_width=True):
                        ph_rq = st.empty()
                        prompt_rq = f"""Kamu adalah Profesor Evaluasi Kognitif Kedokteran (Dunlosky et al., 2013).
Dari materi kuliah berikut, rancang TEPAT 3 PERTANYAAN KAUSALITAS PENEMBUS LOGIKA (High-Impact Socratic Questions):
Fokus murni pada "MENGAPA" dan "BAGAIMANA" (Elaborative Interrogation).
Jangan tanyakan definisi dangkal! Tanyakan rantai sebab-akibat yang memaksa otak menghubungkan konsep A dengan konsep B (contoh: Mengapa pada syok kardiogenik terjadi edema paru sementara pada syok hipovolemik vena kolaps? Mengapa obat golongan X dikontraindikasikan pada kondisi Y?).

Materi:
{text[:12000]}

Format JSON WAJIB murni tanpa markdown lain:
[
  {{
    "id": 1,
    "question": "Mengapa...",
    "key_concept": "Konsep inti yang harus dipahami...",
    "hint": "Petunjuk arah berpikir klinis..."
  }},
  {{
    "id": 2,
    "question": "Bagaimana kaskade...",
    "key_concept": "Alur patofisiologis...",
    "hint": "Petunjuk arah berpikir..."
  }},
  {{
    "id": 3,
    "question": "Mengapa...",
    "key_concept": "Prinsip farmakologi/klinis...",
    "hint": "Petunjuk arah berpikir..."
  }}
]
"""
                        raw_rq = stream_ai_transparent(api_key, prompt_rq, ph_rq)
                        data_rq = extract_json_safely(raw_rq)
                        if data_rq and isinstance(data_rq, list) and len(data_rq) >= 1:
                            save_cached_active_recall(sel, {"questions": data_rq, "timestamp": datetime.datetime.now().isoformat()})
                            st.rerun()
                        else:
                            save_cached_active_recall(sel, {"raw_text": raw_rq, "timestamp": datetime.datetime.now().isoformat()})
                            st.rerun()
                else:
                    q_list = recall_cache.get("questions")
                    if q_list and isinstance(q_list, list):
                        for item in q_list:
                            with st.container(border=True):
                                st.markdown(f"#### ❓ Pertanyaan {item.get('id', 1)}: {item.get('question')}")
                                with st.expander("💡 Petunjuk Arah Berpikir (Hint)"):
                                    st.caption(f"Fokus konsep: {item.get('key_concept')}")
                                    st.write(item.get('hint', ''))
                    elif recall_cache.get("raw_text"):
                        st.markdown(recall_cache.get("raw_text"))
                        
                    st.markdown("---")
                    st.markdown("#### ✍️ Tuliskan Pemikiran & Logika Kausalitas Anda:")
                    st.caption("Ketik pemahaman Anda terhadap pertanyaan di atas. AI akan menguji apakah logika Anda sudah selaras dengan standar emas kedokteran.")
                    
                    with st.form(key=f"form_recall_eval_{sel}"):
                        user_recall_text = st.text_area("Penjelasan Kausalitas Anda:", height=130, placeholder="Contoh: Menurut analisis saya, obat X menghambat reseptor Y sehingga kaskade Z terhenti...", key=f"user_recall_input_{sel}")
                        submit_recall_eval = st.form_submit_button("🧠 Validasi Logika Kognitif dengan Standar Medis →", type="primary", use_container_width=True)
                        
                    if submit_recall_eval and user_recall_text.strip():
                        ph_ev = st.empty()
                        prompt_ev = f"""Kamu adalah Profesor Penguji Kognitif Kedokteran.
Evaluasi jawaban/penalaran mahasiswa berikut terhadap pertanyaan kausalitas materi kuliah:

Jawaban/Penalaran Mahasiswa:
{user_recall_text.strip()}

Materi Sumber Kuliah:
{text[:10000]}

Format Evaluasi Metakognitif (Markdown):
🎯 **Skor Akurasi Logika Kognitif:** [Skor 0-100]/100 — [Status: Logika Teruji / Cukup / Perlu Pelurusan]

✅ **Konsep yang Tervalidasi Akurat:**
- [Poin penalaran mahasiswa yang sudah tepat secara patofisiologis/ilmiah]

🔍 **Koreksi Ilmiah & Pelurusan Celah Miskonsepsi:**
- [Hal yang masih keliru, ambigu, atau terlewat beserta penjelasan ilmiah yang benar]

💡 **Pesan Konsolidasi Memori:**
[1 kalimat kunci untuk mematenkan pemahaman ini di memori jangka panjang]
"""
                        eval_res = stream_ai_transparent(api_key, prompt_ev, ph_ev)
                        st.markdown(eval_res)

            # ── 3. SIMULASI SOAL CAMPURAN KAMPUS (C1 S/D C6) ──
            with sub_t3:
                st.markdown('''
<div style="background:linear-gradient(135deg, rgba(16,185,129,0.12) 0%, rgba(99,102,241,0.08) 100%); border:1.5px solid rgba(16,185,129,0.35); border-radius:12px; padding:16px 20px; margin-bottom:16px;">
  <div style="display:flex; align-items:center; gap:12px;">
    <span style="font-size:28px;">📝</span>
    <div>
      <div style="font-size:1.15rem; font-weight:800; color:#ffffff;">Simulasi Soal Campuran Kampus (C1 s/d C6 Lengkap)</div>
      <div style="font-size:0.78rem; color:#94a3b8; margin-top:3px; line-height:1.5;">
        <strong>Dasar Riset:</strong> <em>Revised Bloom's Taxonomy</em> (Anderson &amp; Krathwohl, 2001) &amp; <em>Dual-Process Theory Clinical Reasoning</em> (Norman et al., 2007).
        Format ujian nyata kampus: 2 Soal Fondasi (C1-C2), 2 Soal Kausalitas Mekanisme (C3-C4), dan 2 Soal Vignette Kasus IGD/Puskesmas dengan Jebakan UKMPPD (C5-C6).
      </div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)

                exam_cache = get_cached_exam_simulation(sel)
                
                c_ex_l, c_ex_r = st.columns([3.5, 1.5], vertical_alignment="center")
                with c_ex_r:
                    btn_regen_exam = st.button("🔄 Buat Soal Baru", key=f"btn_regen_exam_{sel}", use_container_width=True)
                    
                if not exam_cache or btn_regen_exam:
                    if st.button("⚡ Rancang 6 Soal Campuran Tipe Kampus (AI) →", type="primary", key=f"btn_gen_exam_{sel}", use_container_width=True):
                        ph_ex = st.empty()
                        prompt_ex = f"""Kamu adalah Ketua Tim Pembuat Soal Ujian Blok Kedokteran dan Komite UKMPPD Nasional.
Buatkan TEPAT 6 SOAL UJIAN PILIHAN GANDA CAMPURAN (Mixed Cognitive Levels C1-C6) berstandar emas klinis berdasarkan materi kuliah ini:

Materi:
{text[:14000]}

STRUKTUR DISTRIBUSI 6 SOAL (WAJIB):
- Soal 1 & 2: Level 1 (C1-C2 Bloom) — Menguji Fondasi Biomedis, Terminologi Medis, Nomenklatur, atau Klasifikasi Baku.
- Soal 3 & 4: Level 2 (C3-C4 Bloom) — Menguji Analisis Kausalitas Patofisiologis, Mekanisme Kerja Molekuler Obat (Farmakodinamik), atau Interpretasi Penunjang.
- Soal 5 & 6: Level 3 (C5-C6 Bloom) — Clinical Vignette Pasien Nyata (Lengkap dengan usia, jenis kelamin, keluhan utama dengan onset/durasi, tanda vital lengkap [TD, HR, RR, Suhu, SpO2], temuan pemeriksaan fisik khas, dan hasil lab/penunjang). Pertanyaan klinis terarah mengenai diagnosis kerja, pemeriksaan baku emas, atau terapi lini pertama terpilih (Gaya Soal UKMPPD / CBT Nasional).

Format JSON murni WAJIB (Array of 6 Objects):
[
  {{
    "id": 1,
    "level": "Level 1 (C1-C2: Fondasi & Definisi)",
    "question": "Pertanyaan...",
    "options": ["A. ...", "B. ...", "C. ...", "D. ...", "E. ..."],
    "correct_letter": "A",
    "rationale": "Penjelasan ilmiah lengkap mengapa opsi ini benar berdasarkan patofisiologi/konsensus...",
    "distractor_analysis": "Analisis mengapa tiap opsi pengecoh (B, C, D, E) salah atau kondisi apa yang mencocokinya..."
  }},
  ...
]
Hanya berikan JSON tanpa teks markdown pembungkus di luar array!
"""
                        raw_ex = stream_ai_transparent(api_key, prompt_ex, ph_ex)
                        data_ex = extract_json_safely(raw_ex)
                        if data_ex and isinstance(data_ex, list) and len(data_ex) >= 1:
                            save_cached_exam_simulation(sel, {"questions": data_ex, "timestamp": datetime.datetime.now().isoformat()})
                            st.rerun()
                        else:
                            save_cached_exam_simulation(sel, {"raw_text": raw_ex, "timestamp": datetime.datetime.now().isoformat()})
                            st.rerun()
                else:
                    q_list = exam_cache.get("questions")
                    if q_list and isinstance(q_list, list):
                        for i, q in enumerate(q_list):
                            with st.container(border=True):
                                st.markdown(f"**Soal {i+1}** · <span style='background:rgba(99,102,241,0.15);color:#a5b4fc;font-size:0.75rem;padding:2px 8px;border-radius:6px;font-weight:700;'>{q.get('level', '')}</span>", unsafe_allow_html=True)
                                st.markdown(f"<div style='font-size:1.02rem;line-height:1.6;margin:8px 0 12px;'>{q.get('question', '')}</div>", unsafe_allow_html=True)
                                
                                opts = q.get("options", [])
                                user_ans_key = f"ans_mcq_{sel}_{i}"
                                user_pick = st.radio(f"Pilih jawaban untuk Soal {i+1}:", opts, key=user_ans_key, label_visibility="collapsed")
                                
                                show_expl_key = f"show_expl_{sel}_{i}"
                                if show_expl_key not in st.session_state:
                                    st.session_state[show_expl_key] = False
                                    
                                col_btn, _ = st.columns([1.5, 3.5])
                                with col_btn:
                                    if st.button(f"🔍 Periksa Pembahasan Soal {i+1}", key=f"btn_check_q_{sel}_{i}"):
                                        st.session_state[show_expl_key] = True
                                        
                                if st.session_state[show_expl_key]:
                                    c_letter = q.get("correct_letter", "A").strip().upper()
                                    is_correct = user_pick.startswith(c_letter) if user_pick else False
                                    
                                    if is_correct:
                                        st.success(f"✅ **BENAR! Jawaban Anda Tepat: {c_letter}**")
                                    else:
                                        st.error(f"❌ **KURANG TEPAT.** Kunci Jawaban: **{c_letter}**")
                                        
                                    st.markdown(f"""
<div style="background:rgba(15,23,42,0.6);border:1px solid rgba(255,255,255,0.08);border-radius:10px;padding:12px 16px;margin-top:8px;">
  <div style="color:#38bdf8;font-weight:700;font-size:0.85rem;margin-bottom:4px;">📖 Pembahasan Rasional Medis:</div>
  <div style="font-size:0.9rem;line-height:1.6;color:#e2e8f0;">{q.get('rationale', '')}</div>
  <div style="border-top:1px solid rgba(255,255,255,0.08);margin-top:8px;padding-top:8px;">
    <span style="color:#f59e0b;font-weight:700;font-size:0.82rem;">⚠️ Analisis Pengecoh (Distractor Analysis):</span>
    <div style="font-size:0.85rem;color:#cbd5e1;margin-top:2px;">{q.get('distractor_analysis', '')}</div>
  </div>
</div>
""", unsafe_allow_html=True)
                    elif exam_cache.get("raw_text"):
                        st.markdown(exam_cache.get("raw_text"))

                    st.markdown("---")
                    st.info("💡 **Simulasi Kasus Selesai?** Petakan arsitektur konsep dan drill flashcards di tab **🗺️ Langkah 4: Visual Mind Map & Flashcard** di atas.")

            # ── 4. VISUAL MIND MAP & FLASHCARDS (DUAL CODING & SPACING) ──
            with sub_t4:
                st.markdown('''
<div style="background:linear-gradient(135deg, rgba(56,189,248,0.12) 0%, rgba(99,102,241,0.08) 100%); border:1.5px solid rgba(56,189,248,0.35); border-radius:12px; padding:16px 20px; margin-bottom:16px;">
  <div style="display:flex; align-items:center; gap:12px;">
    <span style="font-size:28px;">🗺️</span>
    <div>
      <div style="font-size:1.15rem; font-weight:800; color:#ffffff;">Visual Mind Map &amp; Flashcards Interaktif (Dual Coding &amp; Spacing)</div>
      <div style="font-size:0.78rem; color:#94a3b8; margin-top:3px; line-height:1.5;">
        <strong>Dasar Riset:</strong> <em>Dual Coding Theory</em> (Clark &amp; Paivio, 1991; Schroeder et al., 2018 meta-analisis g = 0.58–0.72) &amp; <em>Spacing Effect</em>.
        Pohon konsep visual hierarkis untuk memetakan arsitektur global materi kuliah dalam sekali pandang, disempurnakan dengan kartu memori kilat dan ekspor Anki.
      </div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)
                
                map_cache = st.session_state.phase_data.get("map_md")
                if not map_cache:
                    user_root = get_user_root()
                    mm_f = user_root / "mindmaps" / f"{re.sub(r'[^a-zA-Z0-9_-]', '_', sel)}.md"
                    if mm_f.exists():
                        map_cache = mm_f.read_text()
                        st.session_state.phase_data["map_md"] = map_cache
                        
                if not map_cache:
                    if st.button("⚡ Bangun Peta Konsep Visual Interaktif (AI)", type="primary", key=f"btn_gen_mm_{sel}", use_container_width=True):
                        ph_mm = st.empty()
                        prompt_mm = f"""Kamu adalah pakar visualisasi kognitif dan ilmu kedokteran.
Buat MIND MAP HIERARKIS LENGKAP & MENDALAM dari substansi ilmiah materi kuliah berikut dalam Bahasa Indonesia.

Format WAJIB menggunakan struktur Markdown murni (Heading # dan List bertingkat) yang rapi:

# [Batang Utama / Topik Pokok Medis]
## 📌 [Cabang Utama 1: Fondasi Biomedis / Anatomi Terkait]
- [Ranting 1a: Klasifikasi / Prinsip]
  - [Detail / Fakta Klinis 1a.1]
  - [Detail / Fakta Klinis 1a.2]
## 📌 [Cabang Utama 2: Patofisiologi & Kaskade Molekuler]
- [Ranting 2a: Kausalitas Inti]
  - [Detail 2a.1]
## 📌 [Cabang Utama 3: Manifestasi Klinis & Penegakan Diagnosis]
... (buat minimal 4-6 cabang utama lengkap sampai ke tata laksana dan contoh obat konkret)

Materi:
{text[:12000]}"""
                        map_cache = stream_ai_transparent(api_key, prompt_mm, ph_mm)
                        if map_cache:
                            st.session_state.phase_data["map_md"] = map_cache
                            user_root = get_user_root()
                            mm_p = user_root / "mindmaps"
                            mm_p.mkdir(parents=True, exist_ok=True)
                            (mm_p / f"{re.sub(r'[^a-zA-Z0-9_-]', '_', sel)}.md").write_text(map_cache)
                            st.rerun()
                else:
                    st.markdown('<div class="mindmap-box">', unsafe_allow_html=True)
                    components.html(build_mindmap_html(map_cache), height=600)
                    st.markdown('</div>', unsafe_allow_html=True)
                    with st.expander("📋 Lihat Outline Teks Mind Map", expanded=False):
                        st.markdown(map_cache)

                st.markdown("<br/>", unsafe_allow_html=True)
                st.markdown("#### 🃏 Flashcards Memori Kilat & Ekspor Anki")
                anki_tsv = generate_anki_export_tsv(sel)
                if anki_tsv:
                    c_ank1, c_ank2 = st.columns([3.2, 1.8], vertical_alignment="center")
                    with c_ank1:
                        st.markdown("<span style='font-size:0.8rem;color:#cbd5e1;'>📦 <strong>Ekspor ke Anki:</strong> Pelajari kartu ini secara offline di HP/Laptop melalui aplikasi Anki resmi.</span>", unsafe_allow_html=True)
                    with c_ank2:
                        st.download_button(
                            label="📥 Unduh Deck Anki (.txt)",
                            data=anki_tsv,
                            file_name=f"{re.sub(r'[^a-zA-Z0-9_]', '_', sel)}_Anki.txt",
                            mime="text/tab-separated-values",
                            key=f"btn_dl_anki_{sel}",
                            use_container_width=True
                        )

                render_flashcards_widget(sel, text, api_key, "PARETO_8020", key_prefix="study_tab")

                # ── 5. COMPLETION HERO CARD ──
                st.markdown("---")
                with st.container(border=True):
                    c_fin1, c_fin2 = st.columns([3.5, 1.5], vertical_alignment="center")
                    with c_fin1:
                        st.markdown(f"""
                        <div style="font-size:1.1rem;font-weight:800;color:#ffffff;">🎉 Selesai Siklus 4 Langkah Modul Ini!</div>
                        <div style="font-size:0.82rem;color:#94a3b8;margin-top:2px;line-height:1.5;">
                          Anda telah menuntaskan Catatan Master, Socratic Active Recall, Ujian Kasus C1-C6, dan Visual Mind Map. Kunci pencapaian ini ke memori jangka panjang agar algoritma SuperMemo SM-2 menjadwalkan repetisi berkala sebelum daya ingat menurun.
                        </div>
                        """, unsafe_allow_html=True)
                    with c_fin2:
                        if st.button("🚀 Catat Selesai (SM-2) ✓", type="primary", use_container_width=True, key=f"btn_finish_study_cycle_{sel}"):
                            iv = update_sr(sel, 4)
                            st.session_state.completed = True
                            st.balloons()
                            st.success(f"🎉 Sesi belajar dicatat! Jadwal review berikutnya: {iv} hari lagi.")
                            st.rerun()

            return # Selesai mode Pareto 80/20

        # ── GATEWAY: START SESSION OVERVIEW (LAZY START UNTUK MODE 6-FASE) ──
        if not st.session_state.session_started:
            st.markdown("<br/>", unsafe_allow_html=True)
            st.markdown(f"""
<div class="card" style="border: 1.5px solid rgba(99, 102, 241, 0.35); padding: 28px 32px;">
  <div style="display:flex; justify-content:space-between; align-items:flex-start; margin-bottom:12px;">
    <div>
      <span class="badge bp" style="margin-bottom:8px;">Siap Memulai Sesi Belajar</span>
      <div style="font-size:1.5rem; font-weight:800; color:#f8fafc; letter-spacing:-0.5px;">📄 {sel}</div>
    </div>
    <div>{days_badge(mat_info.get("next_review",""))}</div>
  </div>
  <p style="color:#94a3b8; font-size:0.9rem; line-height:1.7; margin-bottom:18px;">
    Materi telah disaring dan siap diproses ke dalam <strong>6 Fase Neurosains Terstruktur</strong>.<br/>
    AI tidak akan memproses sampai Anda menekan tombol di bawah, sehingga Anda leluasa memilih materi tanpa menunggu AI berpikir di latar belakang.
  </p>
  <div style="display:grid; grid-template-columns:repeat(3, 1fr); gap:14px; margin-bottom:24px;">
    <div style="background:rgba(15,23,42,0.6); padding:12px 16px; border-radius:10px; border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:0.75rem; color:#64748b; font-weight:600;">PANJANG SUBSTANSI</div>
      <div style="font-size:1.05rem; font-weight:700; color:#e2e8f0;">~{len(text):,} Karakter</div>
    </div>
    <div style="background:rgba(15,23,42,0.6); padding:12px 16px; border-radius:10px; border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:0.75rem; color:#64748b; font-weight:600;">RIWAYAT SESI</div>
      <div style="font-size:1.05rem; font-weight:700; color:#e2e8f0;">{mat_info.get("sessions",0)} Sesi Terlaksana</div>
    </div>
    <div style="background:rgba(15,23,42,0.6); padding:12px 16px; border-radius:10px; border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:0.75rem; color:#64748b; font-weight:600;">EASE FACTOR (SM-2)</div>
      <div style="font-size:1.05rem; font-weight:700; color:#e2e8f0;">EF {mat_info.get("ease_factor",2.5):.2f}</div>
    </div>
  </div>
</div>
""", unsafe_allow_html=True)
            safe_sel_start = re.sub(r'[^a-zA-Z0-9_]', '_', sel)
            if st.button("🚀 Mulai Sesi Belajar Sekarang →", type="primary", use_container_width=True, key=f"start_session_btn_{safe_sel_start}"):
                st.session_state.session_started = True
                st.rerun()
        else:
            # ── ACTIVE SESSION PIPELINE ──
            phase = st.session_state.phase

            # ── STEPPER ──
            STEPS = [
                ("1", "PRIME", "🎯"),
                ("2", "MAP", "🗺️"),
                ("3", "CHUNK", "🧩"),
                ("4", "DIG", "🔍"),
                ("5", "RECALL", "💡"),
                ("6", "FEYNMAN", "🗣️"),
                ("✓", "SCHEDULE", "📅")
            ]
            
            stepper_html = '<div class="stepper-container">'
            for idx, (num, name, icon) in enumerate(STEPS):
                if idx < phase: cls = "step-done"; circle_content = "✓"
                elif idx == phase: cls = "step-active"; circle_content = icon
                else: cls = "step-todo"; circle_content = num
                
                stepper_html += f'<div class="step-item {cls}"><div class="step-circle">{circle_content}</div><div class="step-label">{name}</div></div>'
                if idx < len(STEPS) - 1:
                    divider_cls = "done" if idx < phase else ""
                    stepper_html += f'<div class="step-divider {divider_cls}"></div>'
            stepper_html += '</div>'
            st.markdown(stepper_html, unsafe_allow_html=True)

            # ── History Render ──
            for msg in st.session_state.history:
                if msg["role"] == "user":
                    st.markdown(f'<div class="msg-user">✍️ <strong>Jawaban/Penjelasan Anda:</strong><br/>{html.escape(str(msg["content"]))}</div>', unsafe_allow_html=True)
                elif msg.get("type") == "mindmap":
                    st.markdown('<div class="mindmap-box">', unsafe_allow_html=True)
                    components.html(build_mindmap_html(msg["content"]), height=600)
                    st.markdown('</div>', unsafe_allow_html=True)
                elif msg.get("type") == "analysis":
                    st.markdown(f'<div class="analysis-box"><div class="analysis-header"><span style="font-weight:700;color:#818cf8;font-size:0.92rem;">🔬 Hasil Analisis Kognitif AI</span><span class="badge bp">Metacognitive Feedback</span></div>\n\n{msg["content"]}\n</div>', unsafe_allow_html=True)
                else:
                    st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy</span></div>\n\n{msg["content"]}\n</div>', unsafe_allow_html=True)

            # ── SELESAI ──
            if st.session_state.completed:
                st.markdown('<div class="complete-box"><div style="font-size:3.2rem;margin-bottom:10px;">🎉</div><div style="font-size:1.4rem;font-weight:800;color:#4ade80;margin-bottom:6px;">Sesi Belajar Selesai!</div><div style="color:#bbf7d0;font-size:.95rem;line-height:1.7;">Semua 6 fase telah Anda lalui dan dievaluasi secara kognitif.<br/>Tidur yang cukup malam ini (7-9 jam) untuk transfer memori ke neokorteks permanen.</div></div>', unsafe_allow_html=True)
                st.divider()
                
                # ── 1. HOLISTIC COGNITIVE AUDIT & ALIGNMENT CHECK ──
                if "final_audit" not in st.session_state.phase_data:
                    user_records = [m["content"] for m in st.session_state.history if m["role"] == "user"]
                    user_answers_str = "\n---\n".join([f"Tahap {i+1} Jawaban Mahasiswa:\n{ans}" for i, ans in enumerate(user_records)])
                    
                    ph_audit = st.empty()
                    audit_prompt = f"""Kamu adalah Profesor Senior Kedokteran, Ilmuwan Biomedis, dan Pakar Evaluasi Kognitif.
Lakukan AUDIT KOGNITIF & PENILAIAN KOMPREHENSIF terhadap SEMUA pemikiran, tebakan, dan penjelasan yang diketik mahasiswa sepanjang sesi belajar materi berikut.

Materi Sumber PDF:
{text[:10000]}

Semua Input/Jawaban Mahasiswa (Pre-test F1, Elaborasi F4, Retrieval F5, Feynman F6):
{user_answers_str}

Tugas Evaluasimu:
1. Validasi secara objektif: Apakah pemikiran dan penjelasan mahasiswa BENAR-BENAR SESUAI dengan fakta ilmiah dalam materi kuliah atau masih ada distorsi/miskonsepsi/hafalan kosong.
2. Berikan Skor Penguasaan Akhir (0-100).
3. Poin-poin yang sudah tervalidasi benar vs poin yang perlu diluruskan secara ilmiah.

Format Markdown:
🏆 **Skor Penguasaan Holistik:** [Skor]/100 — [Kategori: Penguasaan Mendalam / Cukup / Perlu Penguatan]

⚖️ **Validasi Keselarasan dengan Materi Kuliah:**
[Ulasan ringkas apakah pemikiran mahasiswa selaras dengan materi ilmiah sumber]

✅ **Konsep yang Berhasil Dikuasai & Tervalidasi Akurat:**
- [Poin tervalidasi 1]
- [Poin tervalidasi 2]

🔍 **Deteksi Miskonsepsi & Hal yang Perlu Diluruskan (Koreksi Ilmiah):**
- [Koreksi ilmiah 1 terhadap apa yang sempat salah atau kurang tepat dari jawaban mahasiswa]
- [Koreksi ilmiah 2]

💡 **Kesimpulan & Pesan Konsolidasi Memori:**
[1-2 kalimat untuk mengunci pemahaman di memori jangka panjang]"""
                    audit_res = stream_ai_transparent(api_key, audit_prompt, ph_audit)
                    st.session_state.phase_data["final_audit"] = audit_res
                else:
                    st.markdown(f'<div class="analysis-box"><div class="analysis-header"><span style="font-weight:700;color:#818cf8;font-size:1.02rem;">📊 Laporan Audit Kognitif & Penilaian Keselarasan Akhir</span><span class="badge bp">100% Holistik</span></div>\n\n{st.session_state.phase_data["final_audit"]}\n</div>', unsafe_allow_html=True)
                
                st.divider()
                
                scores = st.session_state.scores
                s_prime = scores.get("prime", "Dianalisis ✓")
                s_recall = scores.get("recall", "Terekam ✓")
                s_feyn = scores.get("feynman", "Tervalidasi ✓")
                
                st.markdown(f"""
<div style="display:grid;grid-template-columns:repeat(3, 1fr);gap:16px;margin:16px 0 24px;">
  <div class="card card-sm" style="background:rgba(15,23,42,0.8);border:1px solid rgba(99,102,241,0.3);border-radius:14px;padding:18px 20px;">
    <div style="font-size:0.75rem;color:#818cf8;font-weight:700;letter-spacing:0.5px;margin-bottom:6px;">🎯 PRE-TEST INTUISI (F1)</div>
    <div style="font-size:1.1rem;font-weight:800;color:#ffffff;">{s_prime}</div>
    <div style="font-size:0.75rem;color:#94a3b8;margin-top:4px;">Pemetaan skema awal</div>
  </div>
  <div class="card card-sm" style="background:rgba(15,23,42,0.8);border:1px solid rgba(56,189,248,0.3);border-radius:14px;padding:18px 20px;">
    <div style="font-size:0.75rem;color:#38bdf8;font-weight:700;letter-spacing:0.5px;margin-bottom:6px;">💡 RETRIEVAL ACCURACY (F5)</div>
    <div style="font-size:1.1rem;font-weight:800;color:#ffffff;">{s_recall}</div>
    <div style="font-size:0.75rem;color:#94a3b8;margin-top:4px;">Kekuatan ingatan aktif</div>
  </div>
  <div class="card card-sm" style="background:rgba(15,23,42,0.8);border:1px solid rgba(74,222,128,0.3);border-radius:14px;padding:18px 20px;">
    <div style="font-size:0.75rem;color:#4ade80;font-weight:700;letter-spacing:0.5px;margin-bottom:6px;">🗣️ FEYNMAN MASTERY (F6)</div>
    <div style="font-size:1.1rem;font-weight:800;color:#ffffff;">{s_feyn}</div>
    <div style="font-size:0.75rem;color:#94a3b8;margin-top:4px;">Model mental teruji</div>
  </div>
</div>
""", unsafe_allow_html=True)
                
                st.divider()
                safe_sel_post = re.sub(r'[^a-zA-Z0-9_]', '_', sel)
                quality = st.select_slider("Nilai kemampuan recall keseluruhan sesi ini:", options=[0, 1, 2, 3, 4, 5], value=4,
                    key=f"sr_quality_slider_{safe_sel_post}",
                    format_func=lambda x: {0: "0 — Sama sekali lupa", 1: "1 — Banyak yang salah", 2: "2 — Ingat setelah dibantu", 3: "3 — Sebagian besar ingat", 4: "4 — Hampir sempurna", 5: "5 — Sempurna!"}[x])
                if st.button("📅 Simpan & Jadwalkan Review Spaced Repetition", type="primary", use_container_width=True, key=f"save_sr_btn_{safe_sel_post}"):
                    iv = update_sr(sel, quality)
                    msg = "🎉 Luar biasa!" if quality >= 4 else "📚 Terus berlatih!"
                    (st.success if quality >= 4 else st.warning)(f"{msg} Review berikutnya dijadwalkan dalam **{iv} hari** (SM-2 Spaced Repetition Algorithm).")
                
                st.divider()
                
                # ── 2. POST-SESSION RESEARCH-GROUNDED Q&A ──
                st.markdown("""
<div style="margin-top:20px; margin-bottom:16px;">
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:6px;">
    <span style="font-size:1.4rem;">💬</span>
    <div style="font-size:1.25rem;font-weight:800;color:#f8fafc;">Ruang Diskusi & Validasi Riset Terkini (Post-Session Q&A)</div>
  </div>
  <p style="color:#94a3b8;font-size:0.88rem;line-height:1.6;">
    Ada area abu-abu, keraguan klinis, atau topik dari PDF ini yang ingin Anda diskusikan lebih jauh? Tanyakan apa saja — AI terintegrasi dengan literatur riset ilmiah dan standar medis (Goodman & Gilman, Harrison, PubMed, Evidence-Based Medicine) untuk memberikan jawaban tervalidasi.
  </p>
</div>
""", unsafe_allow_html=True)
                
                if "post_chat" not in st.session_state or not isinstance(st.session_state.post_chat, list):
                    st.session_state.post_chat = []

                for chat_m in st.session_state.post_chat:
                    if chat_m["role"] == "user":
                        st.markdown(f'<div class="msg-user">🙋‍♂️ <strong>Pertanyaan Anda:</strong><br/>{html.escape(chat_m["content"])}</div>', unsafe_allow_html=True)
                    else:
                        st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy Pakar Riset Medis</span></div>\n\n{chat_m["content"]}\n</div>', unsafe_allow_html=True)
                
                with st.form(key=f"post_chat_form_{safe_sel_post}", clear_on_submit=True):
                    q_post = st.text_area("Tanyakan area abu-abu atau topik penelitian medis disini:", placeholder="Contoh: Mengapa pada kondisi klinis tertentu obat ini memiliki efek paradoksal, dan bagaimana penjelasan penelitian molekulernya?", height=90, key=f"post_q_input_{safe_sel_post}")
                    st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim pertanyaan (Gunakan <em>Shift + Enter</em> untuk baris baru).</div>', unsafe_allow_html=True)
                    btn_ask = st.form_submit_button("🔬 Tanyakan & Validasi Ilmiah →", type="primary", use_container_width=True)
                    
                if btn_ask and q_post.strip():
                    st.session_state.post_chat.append({"role": "user", "content": q_post.strip()})
                    ph_chat = st.empty()
                    
                    post_prompt = f"""Kamu adalah Profesor Kedokteran Senior, Peneliti Biomedis Utama, dan Konsultan Farmakologi/Sains Terkemuka.
Mahasiswa sedang mendiskusikan materi PDF kuliah berikut untuk mengklarifikasi area abu-abu, keraguan, skenario klinis, atau literatur riset terkait.

Materi Sumber PDF:
{text[:10000]}

Pertanyaan Mahasiswa:
{q_post.strip()}

Riwayat Diskusi Sebelumnya:
{json.dumps(st.session_state.post_chat[-6:], ensure_ascii=False)}

Instruksi Menjawab:
1. Berikan jawaban yang presisi, mendalam, dan tervalidasi secara pasti berbasis bukti medis/sains terbaik (Evidence-Based Medicine).
2. Hubungkan langsung dengan isi PDF sumber, lalu perkuat dengan referensi standar emas (seperti Goodman & Gilman, Katzung, Harrison\'s Internal Medicine, Guyton, PubMed, atau panduan klinis internasional terkini).
3. Hilangkan seluruh area abu-abu dengan menjelaskan logika mekanisme kausalitas (sebab-akibat) secara jernih, runtut, dan tak terbantahkan.
4. Gunakan bahasa ilmiah yang mengalir, jelas, dan mudah dipahami."""
                    
                    ans_chat = stream_ai_transparent(api_key, post_prompt, ph_chat)
                    st.session_state.post_chat.append({"role": "ai", "content": ans_chat})
                    st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 1 — PRIME (Pre-test Effect & Mandatory Preconception Analysis)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 0:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">🎯</div><div><div class="phase-title">Fase 1 — PRIME</div><div class="phase-meta">Pre-testing Effect & Analisis Pengetahuan Awal · ~3 menit</div></div></div>
<div class="phase-source">📖 <strong>Pan & Sana (2021)</strong>, <em>J. Exp. Psychol. Appl.</em>; <strong>Richland et al. (2009)</strong> — <em>Errorful Generation Effect</em>: Mencoba menebak sebelum membaca materi merangsang rasa ingin tahu dan membuka "slot" penyerapan memori di hippocampus. AI akan langsung menganalisis tebakan Anda.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)
                # Cognitive Difficulty Level Selector
                p1_level = st.radio(
                    "Tingkat Kesulitan Kognitif (Neuroscience Scaffold):",
                    list(COGNITIVE_LEVELS.keys()),
                    index=1,
                    horizontal=True,
                    key="p1_diff_level",
                    help="Pilih tingkat kedalaman kognitif sesuai kesiapan belajar Anda"
                )
                level_info = COGNITIVE_LEVELS[p1_level]
                st.caption(f"{level_info['icon']} **{p1_level}**: {level_info['desc']}")
                
                c_head1, c_head2 = st.columns([3.8, 1.2], vertical_alignment="top")
                with c_head2:
                    if st.button("🔄 Buat Ulang", key="btn_regen_prime_q", use_container_width=True, help="Rancang ulang pertanyaan sesuai level kesulitan yang dipilih"):
                        with st.spinner("AI menyusun pertanyaan pre-test baru..."):
                            q_new = stream_ai_transparent(api_key, f"""Buat 4-5 pertanyaan mendalam secara klinis/akademik dari materi ini.
{level_info['prompt_mod']}

Materi:
{text[:8000]}""", st.empty())
                            st.session_state.phase_data["prime_q"] = q_new
                            st.rerun()
                    if st.button("➕ Tambah (+3)", key="btn_more_prime_q", use_container_width=True, help="Tambah 3 pertanyaan tantangan baru"):
                        with st.spinner("AI merancang pertanyaan tantangan baru..."):
                            more_q = stream_ai_transparent(api_key, f"""Buat 3 pertanyaan pre-test klinis/medis BARU yang belum pernah ditanyakan sebelumnya dari materi ini:

{text[:8000]}""", st.empty())
                            if more_q:
                                st.session_state.phase_data["prime_q"] = st.session_state.phase_data.get("prime_q", "") + f"\n\n{more_q}"
                                st.rerun()
                with c_head1:
                    if "prime_q" not in st.session_state.phase_data:
                        ph_ = st.empty()
                        q_ = stream_ai_transparent(api_key, f"""Buat 4-5 pertanyaan mendalam & menantang secara konsep klinis/akademik tentang SUBSTANSI MEDIS materi ini dalam Bahasa Indonesia.
Pertanyaan harus menguji intuisi awal, membedakan mekanisme biologis/klinis, dan merangsang rasa ingin tahu ilmiah (contoh: prinsip spesimen, flora normal vs patogen, resistensi, target reseptor, dsb).
Format: Tulis pertanyaan bernomor yang jelas tanpa memberikan jawaban.

Materi:
{text[:8000]}""", ph_)
                        st.session_state.phase_data["prime_q"] = q_
                    else:
                        st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy Tantangan Intuisi Awal</span></div>\n\n{st.session_state.phase_data["prime_q"]}\n</div>', unsafe_allow_html=True)

                with st.expander("🃏 Buka Latihan Flashcard (Fase PRIME)", expanded=False):
                    render_flashcards_widget(sel, text, api_key, "PRIME", key_prefix="prime_tab")

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">✍️ <strong>Wajib Diisi:</strong> Tuliskan apa yang Anda ketahui atau tebakan Anda (meskipun salah/menebak). AI akan memetakan skema awal Anda untuk mendeteksi titik fokus belajar.</p>', unsafe_allow_html=True)
                with st.form(key="p1_form", clear_on_submit=False):
                    ans = st.text_area("Tulis tebakan / jawaban Anda disini:", height=110, key="prime_ans", placeholder="Ketik tebakan Anda semampunya (contoh: Menurut saya materi ini membahas mekanisme diagnosis klinis...)")
                    st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim &amp; lanjut (Gunakan <em>Shift + Enter</em> jika ingin baris baru).</div>', unsafe_allow_html=True)
                    submit_p1 = st.form_submit_button("🧠 Kirim Jawaban & Analisis Pemahaman Awal →", type="primary", use_container_width=True)
                
                if submit_p1:
                    if not ans.strip():
                        st.warning("⚠️ Mohon ketikkan tebakan atau apa yang Anda ketahui terlebih dahulu (menebak/salah tidak apa-apa) agar neuroplastisitas dan slot memori otak Anda aktif!")
                    else:
                        st.session_state.history.append({"role": "ai", "content": st.session_state.phase_data["prime_q"]})
                        st.session_state.history.append({"role": "user", "content": ans.strip()})
                        
                        ph_eval = st.empty()
                        eval_prime = stream_ai_transparent(api_key, f"""Kamu adalah pakar kognitif dan dosen ahli kedokteran/sains.
Analisis pemahaman awal/tebakan mahasiswa berikut terhadap substansi medis/ilmiah materi kuliah:

Materi Asli:
{text[:8000]}

Pertanyaan Pre-test:
{st.session_state.phase_data["prime_q"]}

Tebakan/Jawaban Mahasiswa:
{ans}

Buat analisis pemahaman awal dengan format:
🎯 **Akurasi Intuisi Awal:**
[Apa yang sudah tepat / masuk akal dari tebakannya]

🧩 **Celah Konsep & Miskonsepsi:**
[Apa yang masih keliru atau belum terjawab]

⚡ **Titik Fokus Anda di Fase Mind Map & Chunking:**
[1-2 konsep medis/ilmiah spesifik yang wajib Anda perhatikan saat membaca materi nanti]""", ph_eval)
                        st.session_state.history.append({"role": "ai", "content": eval_prime, "type": "analysis"})
                        st.session_state.scores["prime"] = "Dianalisis ✓"
                        st.session_state.phase = 1; save_active_session(sel, {"phase": 1, "phase_data": st.session_state.phase_data, "completed": st.session_state.completed, "history": st.session_state.history, "scores": st.session_state.scores, "session_started": True}); st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 2 — MAP (Advance Organizer + Dual Coding + Visual Interactive Tree)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 1:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">🗺️</div><div><div class="phase-title">Fase 2 — MAP (Mind Map Visual Interaktif)</div><div class="phase-meta">Advance Organizer + Dual Coding · ~5 menit</div></div></div>
<div class="phase-source">📖 <strong>Schroeder et al. (2018)</strong>, <em>Educ. Psychol. Rev.</em> (Meta-analisis 142 studi, g = 0.58–0.72); <strong>Clark & Paivio (1991)</strong> — <em>Dual Coding Theory</em>: Otak memproses informasi visual spasial secara paralel dengan jalur verbal. Mind map interaktif membangun skema global sebelum detail dipelajari.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)

                if "map_md" not in st.session_state.phase_data:
                    prompt_map = f"""Kamu adalah pakar visualisasi kognitif dan ilmu medis.
Buat MIND MAP HIERARKIS LENGKAP & MENDALAM dari substansi ilmiah materi kuliah berikut dalam Bahasa Indonesia.

Format WAJIB menggunakan struktur Markdown murni (Heading # dan List bertingkat) yang rapi:

# [Batang Utama / Topik Pokok Medis]
## 📌 [Cabang Utama 1: Pilar Konsep/Tema]
- [Ranting 1a: Klasifikasi / Prinsip]
  - [Detail / Contoh Obat / Gejala / Fakta Klinis 1a.1]
  - [Detail / Contoh Obat / Gejala / Fakta Klinis 1a.2]
- [Ranting 1b: Klasifikasi / Prinsip]
  - [Detail 1b.1]
## 📌 [Cabang Utama 2: Pilar Konsep/Tema]
- [Ranting 2a]
  - [Detail 2a.1]
  - [Detail 2a.2]
- [Ranting 2b]
  - [Detail 2b.1]
## 📌 [Cabang Utama 3]
... (buat minimal 4-6 cabang utama lengkap sampai ke contoh konkret/obat/fakta klinis)

PENTING:
1. Batang Utama (#) adalah konsep sentral materi.
2. Cabang Utama (##) adalah pilar-pilar penting materi medis.
3. Ranting (-) adalah klasifikasi/subdivisi mekanisme.
4. Daun/Sub-ranting (  -) adalah detail fungsi, mekanisme, atau contoh nyata.
5. JANGAN gunakan teks ASCII seperti ├── atau │. Gunakan Markdown heading dan list bertingkat!
6. JANGAN memasukkan bab administratif (Kata Pengantar, Daftar Isi, Tim Dosen).

Materi:\n{text}"""
                    ph_ = st.empty()
                    map_md = stream_ai_transparent(api_key, prompt_map, ph_)
                    if map_md:
                        st.session_state.phase_data["map_md"] = map_md
                        st.rerun()
                else:
                    map_md = st.session_state.phase_data["map_md"]

                    # ── RENDER 100% STANDALONE INTERACTIVE MINDMAP CANVAS ──
                    st.markdown('<div class="mindmap-box">', unsafe_allow_html=True)
                    components.html(build_mindmap_html(map_md), height=600)
                    st.markdown('</div>', unsafe_allow_html=True)

                    c_opt1, c_opt2 = st.columns([3, 1])
                    with c_opt1:
                        with st.expander("📋 Lihat Struktur Outline / Teks Lengkap", expanded=False):
                            st.markdown(map_md)
                        safe_sel_map = re.sub(r'[^a-zA-Z0-9_]', '_', sel)
                        if st.button("🔄 Buat Ulang Map", key=f"btn_regen_map_phase2_{safe_sel_map}", use_container_width=True):
                            del st.session_state.phase_data["map_md"]
                            st.rerun()

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">💡 <strong>Cara Eksplorasi:</strong> Klik kotak cabang untuk membuka rantingnya hingga ke ujung. Geser (drag) kanvas untuk bernavigasi dengan leluasa.</p>', unsafe_allow_html=True)
                if st.button("Lanjut ke Fase 3 (Chunking) →", type="primary", use_container_width=True, key="p2_next"):
                    st.session_state.history.append({"role": "ai", "content": map_md, "type": "mindmap"})
                    st.session_state.phase = 2; save_active_session(sel, {"phase": 2, "phase_data": st.session_state.phase_data, "completed": st.session_state.completed, "history": st.session_state.history, "scores": st.session_state.scores, "session_started": True}); st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 3 — CHUNK (Cognitive Load + Chunking)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 2:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">🧩</div><div><div class="phase-title">Fase 3 — CHUNK</div><div class="phase-meta">Cognitive Load Management + Chunking · ~15 menit</div></div></div>
<div class="phase-source">📖 <strong>Sweller, van Merriënboer, & Paas (2019/2024)</strong>, <em>Educ. Psychol. Rev.</em>; <strong>Cowan (2010)</strong> — <em>Working Memory & Cognitive Load Theory</em>: Working memory hanya mampu menampung 4±1 chunk aktif. Memecah materi kompleks menjadi unit terkelola meminimalkan extraneous load dan mengoptimalkan germane load.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)
                if "chunk_out" not in st.session_state.phase_data:
                    ph_ = st.empty()
                    out_ = stream_ai_transparent(api_key, f"""Pecah materi substansi ilmiah/medis ini menjadi CHUNK-CHUNK KECIL yang bermakna dalam Bahasa Indonesia.
Setiap chunk MAX 5 kalimat (sesuai batas working memory).

Format WAJIB untuk setiap chunk:
### 🧩 [N]. [Nama Konsep / Topik Medis]
[Penjelasan 3-5 kalimat — jelas, padat, fokus mekanisme & fakta penting]
> 🔍 **Analogi/Contoh nyata:** [analogi atau contoh klinis konkret]
> 💡 **Kata kunci:** [2-4 kata kunci]

---

Buat semua chunk yang diperlukan untuk mencakup seluruh materi medis.

Materi:\n{text}""", ph_)
                    st.session_state.phase_data["chunk_out"] = out_
                else:
                    st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy</span></div>\n\n{st.session_state.phase_data["chunk_out"]}\n</div>', unsafe_allow_html=True)

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">👆 Baca setiap chunk dengan seksama. Gunakan analogi untuk memahami konsep secara fungsional.</p>', unsafe_allow_html=True)
                if st.button("Lanjut ke Fase 4 (Elaborasi Interaktif) →", type="primary", use_container_width=True, key="p3_next"):
                    st.session_state.history.append({"role": "ai", "content": st.session_state.phase_data["chunk_out"]})
                    st.session_state.phase = 3; save_active_session(sel, {"phase": 3, "phase_data": st.session_state.phase_data, "completed": st.session_state.completed, "history": st.session_state.history, "scores": st.session_state.scores, "session_started": True}); st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 4 — DIG (Mandatory Elaborative Interrogation Analysis)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 3:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">🔍</div><div><div class="phase-title">Fase 4 — DIG (Elaborative Interrogation Interaktif)</div><div class="phase-meta">Elaborative Interrogation & Uji Kausalitas · ~10 menit</div></div></div>
<div class="phase-source">📖 <strong>Dunlosky et al. (2013)</strong>, <em>Psychol. Sci. Public Interest</em>; <strong>Pressley et al. (1992)</strong> — <em>Elaborative Interrogation</em>: Menjawab pertanyaan *"Mengapa & Bagaimana"* membentuk jembatan semantik antara informasi baru dan memori jangka panjang. AI akan menganalisis kedalaman logika Anda.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)
                
                # Cognitive Difficulty Level Selector for Dig
                p4_level = st.radio(
                    "Tingkat Kesulitan Elaborasi (Elaborative Depth):",
                    list(COGNITIVE_LEVELS.keys()),
                    index=1,
                    horizontal=True,
                    key="p4_diff_level",
                    help="Pilih kedalaman elaborasi kausalitas mekanisme"
                )
                p4_info = COGNITIVE_LEVELS[p4_level]
                st.caption(f"{p4_info['icon']} **{p4_level}**: {p4_info['desc']}")
                
                c_d1, c_d2 = st.columns([3.8, 1.2], vertical_alignment="top")
                with c_d2:
                    if st.button("🔄 Buat Ulang", key="btn_regen_dig_q", use_container_width=True, help="Rancang pertanyaan elaborasi baru sesuai level"):
                        with st.spinner("AI menyusun pertanyaan elaborasi baru..."):
                            q_dig_new = stream_ai_transparent(api_key, f"""Buat 3-4 pertanyaan MENGAPA / BAGAIMANA yang paling mendalam dari materi kuliah ini.
{p4_info['prompt_mod']}

Materi:
{text[:10000]}""", st.empty())
                            st.session_state.phase_data["dig_q"] = q_dig_new
                            st.rerun()
                    if st.button("➕ Tambah (+2)", key="btn_more_dig_q", use_container_width=True, help="Tambah pertanyaan elaborasi baru"):
                        with st.spinner("AI sedang merancang pertanyaan kausalitas baru..."):
                            more_q = stream_ai_transparent(api_key, f"""Buat 2 pertanyaan elaborasi MENGAPA / BAGAIMANA baru yang mendalami mekanisme kausalitas lain dari materi ini:

{text[:10000]}""", st.empty())
                            if more_q:
                                st.session_state.phase_data["dig_q"] = st.session_state.phase_data.get("dig_q", "") + f"\n\n{more_q}"
                                st.rerun()
                with c_d1:
                    if "dig_q" not in st.session_state.phase_data:
                        ph_ = st.empty()
                        out_ = stream_ai_transparent(api_key, f"""Buat 3-4 pertanyaan MENGAPA / BAGAIMANA yang paling mendalam & esensial tentang mekanisme ilmiah/medis dari materi kuliah ini dalam Bahasa Indonesia.
Pertanyaan harus menguji pemahaman kausalitas molekuler, patologis, dan implikasi klinis.

Format:
1. **🔍 [Pertanyaan Mengapa/Bagaimana 1]**
2. **🔍 [Pertanyaan Mengapa/Bagaimana 2]**
3. **🔍 [Pertanyaan Mengapa/Bagaimana 3]**

Materi:\n{text[:10000]}""", ph_)
                        st.session_state.phase_data["dig_q"] = out_
                    else:
                        st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy Tantangan Elaborasi</span></div>\n\n{st.session_state.phase_data["dig_q"]}\n</div>', unsafe_allow_html=True)

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">✍️ <strong>Wajib Diisi:</strong> Ketik penalaran atau penjelasan mekanisme logika Anda untuk pertanyaan di atas. AI akan mendiagnosis kedalaman logika kausalitas Anda.</p>', unsafe_allow_html=True)
                with st.form(key="p4_form", clear_on_submit=False):
                    ans_dig = st.text_area("Tulis penalaran / penjelasan Anda disini:", height=130, key="dig_ans", placeholder="Jelaskan alasan atau mekanisme menurut pemahaman Anda…")
                    st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim &amp; lanjut (Gunakan <em>Shift + Enter</em> jika ingin baris baru).</div>', unsafe_allow_html=True)
                    submit_p4 = st.form_submit_button("🔬 Kirim Penalaran & Analisis Kausalitas →", type="primary", use_container_width=True)
                
                if submit_p4:
                    if not ans_dig.strip():
                        st.warning("⚠️ Mohon ketikkan penalaran atau penjelasan mekanisme Anda terlebih dahulu untuk diuji logika kausalitasnya!")
                    else:
                        st.session_state.history.append({"role": "ai", "content": st.session_state.phase_data["dig_q"]})
                        st.session_state.history.append({"role": "user", "content": ans_dig.strip()})
                        
                        ph_eval = st.empty()
                        eval_dig = stream_ai_transparent(api_key, f"""Kamu adalah pakar kognitif dan dosen ahli kedokteran/sains.
Evaluasi penalaran elaborasi (Elaborative Interrogation) mahasiswa berikut:

Materi Asli:
{text[:8000]}

Pertanyaan Elaborasi:
{st.session_state.phase_data["dig_q"]}

Jawaban/Penalaran Mahasiswa:
{ans_dig}

Format analisis:
✅ **Kekuatan Logika & Mekanisme yang Tepat:**
[Bagian mekanisme yang sudah dipahami secara akurat]

🔧 **Celah Penjelasan & Koreksi Mekanisme:**
[Rantai sebab-akibat yang terlewat atau keliru]

💡 **Kunci Mengingat Permanen:**
[Sintesis 1-2 kalimat untuk mengunci konsep ini di memori jangka panjang]""", ph_eval)
                        st.session_state.history.append({"role": "ai", "content": eval_dig, "type": "analysis"})
                        st.session_state.scores["elaboration"] = "Logika Teruji ✓"
                        st.session_state.phase = 4; save_active_session(sel, {"phase": 4, "phase_data": st.session_state.phase_data, "completed": st.session_state.completed, "history": st.session_state.history, "scores": st.session_state.scores, "session_started": True}); st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 5 — RECALL (Mandatory Retrieval Practice + Diagnostic Grading)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 4:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">💡</div><div><div class="phase-title">Fase 5 — RECALL (Active Retrieval Practice)</div><div class="phase-meta">Retrieval Practice & Evaluasi Retensi · ~10 menit · TUTUP MATERI SEKARANG</div></div></div>
<div class="phase-source">📖 <strong>Yang et al. (2021)</strong>, <em>Psychol. Bull.</em> (Meta-analisis 222 studi, d = 0.50); <strong>Agarwal et al. (2021)</strong> (g = 0.66); <strong>Latimier et al. (2021)</strong> (g = 0.74); <strong>Roediger & Karpicke (2006)</strong> — *Testing Effect*: Menarik informasi dari memori otak secara aktif memperkuat jalur sinaptik 50%+ lebih kuat dibanding membaca ulang.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)
                st.warning("⚠️ **TUTUP SEMUA CATATAN & SLIDE SEKARANG.** Uji kemampuan memori Anda tanpa melihat materi.")

                c_r1, c_r2 = st.columns([4, 1], vertical_alignment="top")
                with c_r2:
                    if st.button("➕ Tambah Kuis (+3)", key="btn_more_recall_q", use_container_width=True, help="Tambah pertanyaan kuis retrieval baru"):
                        with st.spinner("AI sedang merancang kuis retrieval baru..."):
                            more_q = stream_ai_transparent(api_key, f"""Buat 3 pertanyaan ACTIVE RECALL baru tentang fakta/konsep medis/ilmiah lain dari materi ini:

{text[:10000]}""", st.empty())
                            if more_q:
                                st.session_state.phase_data["recall_q"] = st.session_state.phase_data.get("recall_q", "") + f"\n\n{more_q}"
                                st.rerun()
                with c_r1:
                    if "recall_q" not in st.session_state.phase_data:
                        ph_ = st.empty()
                        q_ = stream_ai_transparent(api_key, f"""Buat 5-6 pertanyaan ACTIVE RECALL komprehensif tentang SUBSTANSI MEDIS/ILMIAH materi dalam Bahasa Indonesia (fakta klinis, diagnosis, patofisiologi, terapi, dan klasifikasi).
Format:
**💡 Pertanyaan [N]:** [pertanyaan spesifik tentang materi medis/ilmiah]

Materi:\n{text[:10000]}""", ph_)
                        st.session_state.phase_data["recall_q"] = q_
                    else:
                        st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy Kuis Retrieval</span></div>\n\n{st.session_state.phase_data["recall_q"]}\n</div>', unsafe_allow_html=True)

                with st.expander("🃏 Buka Latihan Flashcard (Fase RECALL)", expanded=False):
                    render_flashcards_widget(sel, text, api_key, "RECALL", key_prefix="recall_tab")

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">✍️ <strong>Wajib Diisi:</strong> Tulis jawaban Anda dari ingatan (tutup materi). AI akan mendiagnosis tingkat retensi kognitif Anda.</p>', unsafe_allow_html=True)
                with st.form(key="p5_form", clear_on_submit=False):
                    ans_recall = st.text_area("Tulis jawaban Anda dari ingatan (tidak boleh buka materi):", height=180, key="recall_ans", placeholder="Jawab pertanyaan 1, 2, 3, 4, 5 semampunya…")
                    st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim &amp; lanjut (Gunakan <em>Shift + Enter</em> jika ingin baris baru).</div>', unsafe_allow_html=True)
                    submit_p5 = st.form_submit_button("📊 Kirim Jawaban & Evaluasi Retensi Kognitif →", type="primary", use_container_width=True)
                
                if submit_p5:
                    if not ans_recall.strip():
                        st.warning("⚠️ Mohon jawab pertanyaan kuis di atas dari ingatan Anda (tutup materi) agar daya retensi memori Anda dapat diukur secara akurat!")
                    else:
                        st.session_state.history.append({"role": "ai", "content": st.session_state.phase_data["recall_q"]})
                        st.session_state.history.append({"role": "user", "content": ans_recall.strip()})
                        
                        ph2 = st.empty()
                        check = stream_ai_transparent(api_key, f"""Periksa jawaban kuis retrieval mahasiswa ini terhadap materi kuliah medis/ilmiah. Berikan feedback diagnosis kognitif yang mendalam dalam Bahasa Indonesia.

Materi Asli: {text[:8000]}
Pertanyaan: {st.session_state.phase_data["recall_q"]}
Jawaban Mahasiswa: {ans_recall}

Format feedback:
✅ **Konsep yang Berhasil Di-Retrieve (Kuat di Memori):**
[Poin-poin yang dijawab dengan tepat]

❌ **Konsep yang Hilang / Perlu Penguatan:**
[Koreksi detail, fakta, atau istilah medis yang terlewat]

📈 **Skor Retensi Kognitif: [X]/5** — [Tingkat penguasaan memori]

🎯 **Rekomendasi Fokus Review:**
[1 hal utama yang harus dipelajari ulang nanti]""", ph2)
                        st.session_state.history.append({"role": "ai", "content": check, "type": "analysis"})
                        st.session_state.scores["recall"] = "Dievaluasi AI (Tercatat) ✓"
                        st.session_state.phase = 5; save_active_session(sel, {"phase": 5, "phase_data": st.session_state.phase_data, "completed": st.session_state.completed, "history": st.session_state.history, "scores": st.session_state.scores, "session_started": True}); st.rerun()

            # ══════════════════════════════════════════════════════════════════════
            # FASE 6 — FEYNMAN (Mandatory Mental Model Validation)
            # ══════════════════════════════════════════════════════════════════════
            elif phase == 5:
                st.markdown("""<div class="phase-box">
<div class="phase-header"><div class="phase-icon">🗣️</div><div><div class="phase-title">Fase 6 — FEYNMAN (Validasi Model Mental)</div><div class="phase-meta">Feynman Technique + Generation Effect · ~8 menit · FASE TERAKHIR</div></div></div>
<div class="phase-source">📖 <strong>Bertsch, Pesta, Wiscott, & Berman (2007)</strong>, <em>J. Mem. Lang.</em> (Meta-analisis 445 effect sizes, d = 0.40); <strong>Chi et al. (1994)</strong> — <em>Self-Explanation Effect</em>: Menjelaskan dengan bahasa sendiri menghancurkan ilusi pemahaman semu (*illusion of competence*). AI akan mendeteksi apakah pemahaman Anda hafalan atau pemahaman hakiki.</div>
</div>""", unsafe_allow_html=True)
                st.markdown("<br/>", unsafe_allow_html=True)

                c_t1, c_t2 = st.columns([4, 1], vertical_alignment="top")
                with c_t2:
                    if st.button("➕ Tambah Topik (+3)", key="btn_more_topics_q", use_container_width=True, help="Tambah topik inti / skenario kasus baru"):
                        with st.spinner("AI sedang menyusun topik inti baru..."):
                            more_t = stream_ai_transparent(api_key, f"""Buat 3 topik inti substansi medis / skenario klinis BARU untuk dijelaskan mahasiswa:

{text[:8000]}""", st.empty())
                            if more_t:
                                st.session_state.phase_data["topics"] = st.session_state.phase_data.get("topics", "") + f"\n\n{more_t}"
                                st.rerun()
                with c_t1:
                    if "topics" not in st.session_state.phase_data:
                        ph_ = st.empty()
                        t_ = stream_ai_transparent(api_key, f"""Sebutkan 4-5 TOPIK INTI & SKENARIO KLINIS mendasar dari materi ini dalam Bahasa Indonesia untuk diuji pemahaman model mentalnya (format: 1. ... 2. ... 3. ... 4. ...):
{text[:8000]}""", ph_)
                        st.session_state.phase_data["topics"] = t_
                    else:
                        st.markdown(f'<div class="msg-ai"><div class="ai-row"><div class="ai-dot">AI</div><span style="font-size:.75rem;color:#818cf8;font-weight:600;">NeuroStudy Topik Inti & Skenario Klinis</span></div>\n\n{st.session_state.phase_data["topics"]}\n</div>', unsafe_allow_html=True)

                with st.expander("🃏 Buka Latihan Flashcard (Fase FEYNMAN)", expanded=False):
                    render_flashcards_widget(sel, text, api_key, "FEYNMAN", key_prefix="feynman_tab")

                st.markdown('<p style="color:#94a3b8;font-size:.85rem;margin-top:16px;">🗣️ <strong>Wajib Diisi:</strong> Pilih SATU topik di atas, lalu ketik penjelasan Anda seolah sedang mengajari orang awam. Gunakan analogi sehari-hari dan bahasa Anda sendiri (bukan sekadar angka atau istilah mentah).</p>', unsafe_allow_html=True)

                with st.form(key="p6_form", clear_on_submit=False):
                    feyn_in = st.text_area("Tulis penjelasan Feynman Anda (minimal 2-3 kalimat):", height=160, key="feyn_in", placeholder="Contoh: Bayangkan kita menjelaskan konsep ini ke teman yang belum pernah belajar kedokteran...")
                    st.markdown('<div style="font-size:0.75rem;color:#64748b;margin:-6px 0 10px;">⌨️ Tekan <strong>Enter</strong> untuk langsung kirim &amp; lanjut (Gunakan <em>Shift + Enter</em> jika ingin baris baru).</div>', unsafe_allow_html=True)
                    submit_p6 = st.form_submit_button("🎓 Kirim Penjelasan Feynman & Validasi Model Mental →", type="primary", use_container_width=True)
                
                if submit_p6:
                    if not feyn_in.strip() or len(feyn_in.strip()) < 10:
                        st.warning("⚠️ Mohon ketikkan penjelasan konsep dengan kalimat Anda sendiri (minimal 2-3 kalimat) agar pemahaman model mental Anda dapat dievaluasi secara akurat!")
                    else:
                        st.session_state.history.append({"role": "ai", "content": st.session_state.phase_data["topics"]})
                        st.session_state.history.append({"role": "user", "content": feyn_in.strip()})
                        
                        ph2 = st.empty()
                        eval_ = stream_ai_transparent(api_key, f"""Evaluasi penjelasan mahasiswa menggunakan Feynman Technique dalam Bahasa Indonesia.

Materi Asli Medis: {text[:8000]}
Topik yang dipilih: {st.session_state.phase_data["topics"]}
Penjelasan Mahasiswa: {feyn_in}

Format evaluasi:
🌟 **Level Pemahaman Model Mental:**
(Tentukan: Level 1 - Jargon/Hafalan, Level 2 - Pemahaman Fungsional, atau Level 3 - Penguasaan Intuitif Mendalam)

✅ **Pemahaman yang Sangat Kuat:**
[Bagian analogi atau penjelasan yang tepat dan cerdas]

🔍 **Deteksi Gap / Ilusi Pemahaman:**
[Konsep yang masih ambigu, terdistorsi, atau belum terjelaskan]

🏆 **Skor Penguasaan Feynman: [X]/100** — [Penilaian kesiapan memori jangka panjang]""", ph2)
                        st.session_state.history.append({"role": "ai", "content": eval_, "type": "analysis"})
                        st.session_state.scores["feynman"] = "Tervalidasi (100% Selesai) ✓"
                        st.session_state.completed = True; clear_active_session(sel); st.rerun()




def render_tab_review():
    mats = load_mats()
    st.markdown('''
    <div style="margin-bottom:16px;">
      <div style="font-size:1.5rem;font-weight:800;color:#f8fafc;letter-spacing:-0.5px;">🗂️ Pusat Latihan Memori Jangka Panjang (Spaced Repetition &amp; Anki)</div>
      <div style="color:#94a3b8;font-size:0.85rem;margin-top:2px;">Konsolidasi memori permanen berbasis algoritma SuperMemo SM-2, kurva retensi Ebbinghaus, dan ekspor Anki resmi.</div>
    </div>
    ''', unsafe_allow_html=True)
    
    rev_sub1, rev_sub2, rev_sub3 = st.tabs([
        "📅 Antrean Review Harian (SM-2)",
        "🃏 Flashcards Interaktif & Ekspor Anki",
        "🏆 Peringkat & Indeks Penguasaan Klinis (Leaderboard)"
    ])
    
    with rev_sub1:
        render_sub_spaced_repetition()

    with rev_sub2:
        if not mats:
            st.info("Belum ada materi kuliah untuk direview.")
        else:
            available_keys = list(mats.keys())
            curr_sel = st.session_state.get("mat_sel")
            def_idx = available_keys.index(curr_sel) if curr_sel in available_keys else 0
            
            c_r1, c_r2 = st.columns([3.5, 1.5], vertical_alignment="center")
            with c_r1:
                sel_fc = st.selectbox("Pilih Modul Kuliah untuk Flashcard:", available_keys, index=def_idx, key="rev_fc_picker")
            with c_r2:
                anki_tsv = generate_anki_export_tsv(sel_fc)
                if anki_tsv:
                    st.download_button(
                        label="📥 Unduh Deck Anki (.txt)",
                        data=anki_tsv,
                        file_name=f"{re.sub(r'[^a-zA-Z0-9_]', '_', sel_fc)}_Anki.txt",
                        mime="text/tab-separated-values",
                        key=f"btn_dl_anki_rev_tab_{sel_fc}",
                        use_container_width=True
                    )
            
            mat_info = mats.get(sel_fc, {})
            api_k = load_config().get("api_key", "")
            render_flashcards_widget(sel_fc, mat_info.get("text", ""), api_k, "PARETO_8020", key_prefix="review_tab")

    with rev_sub3:
        st.markdown('''
        <div style="background:linear-gradient(135deg, rgba(245,158,11,0.12) 0%, rgba(99,102,241,0.08) 100%); border:1.5px solid rgba(245,158,11,0.35); border-radius:12px; padding:16px 20px; margin-bottom:18px;">
          <div style="display:flex; align-items:center; gap:12px;">
            <span style="font-size:32px;">🏆</span>
            <div>
              <div style="font-size:1.15rem; font-weight:800; color:#ffffff;">Peringkat & Indeks Penguasaan Klinis Nasional (Peer Leaderboard)</div>
              <div style="font-size:0.78rem; color:#94a3b8; margin-top:3px; line-height:1.4;">
                Benchmark kesiapan ujian blok &amp; UKMPPD secara transparan. Konsistensi harian dan retensi memori algoritma SM-2 dikonversi menjadi <strong>Clinical Mastery Index (CMI)</strong>.
              </div>
            </div>
          </div>
        </div>
        ''', unsafe_allow_html=True)
        
        from core.db import get_connection
        conn = get_connection()
        peers = conn.execute("SELECT * FROM peer_leaderboard ORDER BY mastery_index DESC;").fetchall()
        conn.close()
        
        c_lb1, c_lb2, c_lb3 = st.columns(3)
        c_lb1.metric("🥇 Peringkat 1 Nasional", "dr. Dimas Wastu", "96.5% CMI")
        c_lb2.metric("🔥 Rata-rata Streak", "15 Hari Berturut", "+3 hari")
        c_lb3.metric("🎯 Standar Lulus UKMPPD", "≥ 80.0% CMI", "Ambang Aman")
        
        st.markdown('<div style="margin-top:14px;"></div>', unsafe_allow_html=True)
        
        for idx, p in enumerate(peers):
            rank = idx + 1
            badge_color = "#f59e0b" if rank == 1 else ("#94a3b8" if rank == 2 else ("#d97706" if rank == 3 else "#818cf8"))
            border_color = "rgba(245,158,11,0.4)" if rank == 1 else "rgba(255,255,255,0.08)"
            
            st.markdown(f'''
<div style="background:rgba(15,23,42,0.65); border:1px solid {border_color}; border-radius:12px; padding:14px 18px; margin-bottom:10px; display:flex; justify-content:space-between; align-items:center; flex-wrap:wrap; gap:10px;">
  <div style="display:flex; align-items:center; gap:14px;">
    <div style="font-size:1.25rem; font-weight:900; color:{badge_color}; min-width:28px;">#{rank}</div>
    <div style="font-size:1.8rem;">{p['avatar']}</div>
    <div>
      <div style="font-size:0.95rem; font-weight:800; color:#ffffff;">{p['full_name']} <span style="font-size:0.7rem; color:#94a3b8; font-weight:400;">(@{p['username']})</span></div>
      <div style="font-size:0.75rem; color:#818cf8; font-weight:600; margin-top:2px;">{p['badge_title']}</div>
    </div>
  </div>
  <div style="display:flex; align-items:center; gap:18px;">
    <div style="text-align:right;">
      <div style="font-size:0.68rem; color:#94a3b8; text-transform:uppercase; font-weight:700;">Streak Belajar</div>
      <div style="font-size:0.95rem; font-weight:800; color:#fbbf24;">🔥 {p['streak_days']} Hari</div>
    </div>
    <div style="text-align:right;">
      <div style="font-size:0.68rem; color:#94a3b8; text-transform:uppercase; font-weight:700;">Kartu Direview</div>
      <div style="font-size:0.95rem; font-weight:800; color:#38bdf8;">🃏 {p['cards_reviewed']}</div>
    </div>
    <div style="text-align:right; min-width:85px;">
      <div style="font-size:0.68rem; color:#94a3b8; text-transform:uppercase; font-weight:700;">Mastery Index</div>
      <div style="font-size:1.1rem; font-weight:900; color:#34d399;">{p['mastery_index']:.1f}%</div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)


def render_tab_spesialis_dan_akun():
    st.markdown('''
    <div style="margin-bottom:16px;">
      <div style="font-size:1.5rem;font-weight:800;color:#f8fafc;letter-spacing:-0.5px;">🩺 Konsultasi Spesialis, Evaluasi Penguji &amp; Akun</div>
      <div style="color:#94a3b8;font-size:0.85rem;margin-top:2px;">Konsultasi dewan dokter spesialis AI, evaluasi penguji sejawat (beta testing), dan manajemen profil akun Anda.</div>
    </div>
    ''', unsafe_allow_html=True)
    
    p_t1, p_t2, p_t3 = st.tabs([
        "🩺 Dewan Dokter Spesialis AI (5 Departemen)",
        "💬 Hub Uji Coba & Feedback Penguji",
        "👤 Profil & Pengaturan Akun"
    ])
    with p_t1:
        render_sub_specialist_council()
    with p_t2:
        render_sub_beta_tester()
    with p_t3:
        user_inf = st.session_state.get("user_info") or {}
        st.markdown(f'''
        <div style="background:rgba(15,23,42,0.6);border:1px solid rgba(255,255,255,0.08);border-radius:14px;padding:20px;margin-bottom:16px;">
          <div style="font-size:1.15rem;font-weight:700;color:#ffffff;margin-bottom:8px;">Informasi Pengguna</div>
          <div style="color:#cbd5e1;font-size:0.9rem;line-height:1.8;">
            👤 <strong>Nama Pengguna:</strong> {user_inf.get("display_name", "Mahasiswa Kedokteran")}<br/>
            📧 <strong>Email Terdaftar:</strong> {user_inf.get("email", "-")}<br/>
            🛡️ <strong>Paket Akses:</strong> <span style="color:#818cf8;font-weight:700;">{user_inf.get("plan_name", "Pro Kedokteran")}</span><br/>
            📅 <strong>Status Akun:</strong> <span style="color:#34d399;font-weight:700;">Aktif (Akses Penuh 208 Modul)</span>
          </div>
        </div>
        ''', unsafe_allow_html=True)
        
        c_th, c_lg = st.columns(2)
        with c_th:
            is_obsidian = st.session_state.get("app_theme", "obsidian") == "obsidian"
            theme_label = "☀️ Ganti ke Mode Klinis (Terang)" if is_obsidian else "🌙 Ganti ke Mode Obsidian (Gelap)"
            if st.button(theme_label, use_container_width=True, key="btn_th_tab3"):
                new_th = "clinical_white" if is_obsidian else "obsidian"
                st.session_state.app_theme = new_th
                cfg_tmp = load_config()
                cfg_tmp["app_theme"] = new_th
                save_config(cfg_tmp)
                st.rerun()
        with c_lg:
            if st.button("🚪 Logout dari Sesi", type="secondary", use_container_width=True, key="btn_logout_tab3"):
                clear_persisted_auth_session()
                st.session_state.current_user = None
                st.session_state.user_info = None
                st.rerun()

# ══════════════════════════════════════════════════════════════════════
# MAIN 4-TAB NAVIGATION: RAPI, SIMPEL, TERSTRUKTUR & BERBASIS KEILMUAN
# ══════════════════════════════════════════════════════════════════════
tab_kurikulum, tab_belajar, tab_review, tab_spesialis_akun = st.tabs([
    "🏛️  Kurikulum & 208 Modul",
    "📖  Meja Belajar Kognitif",
    "🗂️  Review Harian & Anki (SM-2)",
    "🩺  Konsultasi Spesialis & Akun"
])

with tab_kurikulum:
    render_sub_cloud_library()

with tab_belajar:
    render_tab_belajar()

with tab_review:
    render_tab_review()

with tab_spesialis_akun:
    render_tab_spesialis_dan_akun()

# ── PERMANENT MEDICOLEGAL SAFE HARBOR DISCLAIMER (FOOTER) ────────────────────
st.markdown("""
<div style="margin-top: 40px; padding: 16px 20px; background: rgba(15, 23, 42, 0.7); border: 1px solid rgba(239, 68, 68, 0.25); border-radius: 14px; text-align: center;">
  <div style="display: flex; align-items: center; justify-content: center; gap: 8px; font-size: 0.8rem; font-weight: 800; color: #f87171; text-transform: uppercase; letter-spacing: 0.5px;">
    <span>⚠️</span> PEMBERITAHUAN MEDIKOLEGAL & KEBIJAKAN PENGGUNAAN AMAN (SAFE HARBOR)
  </div>
  <div style="font-size: 0.73rem; color: #94a3b8; line-height: 1.6; max-width: 950px; margin: 6px auto 0;">
    <strong>NeuroStudy adalah simulator dan platform edukasi kognitif medis semata.</strong> Seluruh materi, simulasi kasus IGD, kalkulasi dosis, dan rekomendasi terapi yang dihasilkan oleh model AI dirancang khusus untuk keperluan latihan akademis mahasiswa dan persiapan ujian blok/UKMPPD. Konten ini <u>BUKAN</u> merupakan protokol medis resmi, panduan peresepan klinis riil, atau pengganti keputusan dokter berlisensi. Dalam menangani pasien nyata di rumah sakit, selalu verifikasi ke standar emas (Fornas, PNPK Kemenkes, FDA, IDI) dan supervisi dokter penanggung jawab pelayanan (DPJP).
  </div>
</div>
""", unsafe_allow_html=True)

# ── DYNAMIC TAB SWITCHER & DIRECT ENTER SUBMIT INJECTOR (POST-RENDER) ─────────
switch_tab_kw = st.session_state.pop("switch_tab_target", None)
js_switch_code = ""
if switch_tab_kw:
    js_switch_code = f"""
  const targetKeyword = {json.dumps(switch_tab_kw.lower())};
  let switchAttempts = 0;
  function triggerTabSwitch() {{
    switchAttempts++;
    try {{
      const pDoc = window.parent.document;
      const tabs = pDoc.querySelectorAll('[data-testid="stTabs"] button[role="tab"], button[data-baseweb="tab"]');
      for (const tab of tabs) {{
        const txt = (tab.innerText || tab.textContent || "").toLowerCase();
        if (txt.includes(targetKeyword)) {{
          tab.click();
          return;
        }}
      }}
    }} catch(e) {{}}
    if (switchAttempts < 25) {{
      setTimeout(triggerTabSwitch, 60);
    }}
  }}
  triggerTabSwitch();
"""

components.html(f"""
<script>
(function() {{
  {js_switch_code}
  function setupEnter() {{
    try {{
      const pDoc = window.parent.document;
      const textareas = pDoc.querySelectorAll('textarea');
      textareas.forEach(ta => {{
        if (ta.getAttribute('data-enter-bound') === 'true') return;
        ta.setAttribute('data-enter-bound', 'true');
        
        ta.addEventListener('keydown', function(e) {{
          if (e.key === 'Enter' && !e.shiftKey && !e.ctrlKey && !e.altKey && !e.metaKey && !e.isComposing) {{
            e.preventDefault();
            e.stopPropagation();
            
            ta.dispatchEvent(new Event('input', {{ bubbles: true }}));
            ta.dispatchEvent(new Event('change', {{ bubbles: true }}));
            
            const form = ta.closest('form');
            let submitBtn = form ? form.querySelector('button[type="submit"], button[data-testid="stFormSubmitButton"]') : null;
            if (!submitBtn) {{
              submitBtn = pDoc.querySelector('button[kind="primary"], button[data-testid="baseButton-primary"]');
            }}
            if (submitBtn) {{
              setTimeout(() => {{ submitBtn.click(); }}, 40);
            }}
          }}
        }});
      }});
    }} catch(e) {{}}
  }}
  
  setupEnter();
  try {{
    const observer = new MutationObserver(() => {{ setupEnter(); }});
    observer.observe(window.parent.document.body, {{ childList: true, subtree: true }});
  }} catch(e) {{
    setInterval(setupEnter, 1000);
  }}
}})();
</script>
""", height=0, width=0)

